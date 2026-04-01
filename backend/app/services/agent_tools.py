"""Agent tools — unified file-based tools that give digital employees
access to their own structured workspace.

Design principle:  ONE set of file tools covers EVERYTHING.
The agent's workspace uses well-known paths:
  - tasks.json          → task list (auto-synced from DB)
  - soul.md             → personality definition
  - memory.md           → long-term memory / notes
  - skills/             → skill definitions (markdown files)
  - workspace/          → general working files, reports, etc.

The agent reads/writes these files directly. No per-concept tools needed.
"""

import json
import os
import uuid
from contextvars import ContextVar
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Optional, Any
import re

from loguru import logger
from sqlalchemy import select, or_

from app.database import async_session
from app.models.task import Task
from app.models.agent import Agent as AgentModel
from app.models.org import AgentRelationship, OrgMember, AgentAgentRelationship
from app.models.audit import ChatMessage, AuditLog
from app.models.chat_session import ChatSession
from app.models.channel_config import ChannelConfig
from app.models.user import User as UserModel
from app.services.auth_registry import auth_provider_registry
from app.services.channel_session import find_or_create_channel_session
from app.services.channel_user_service import get_platform_user_by_org_member
from app.config import get_settings


_settings = get_settings()
WORKSPACE_ROOT = Path(_settings.AGENT_DATA_DIR)

# ─── Tool Config Cache ──────────────────────────────────────────
# Cache tool configurations to avoid frequent DB queries
# Key: (agent_id, tool_name), Value: (config, expiry_time)
_tool_config_cache: dict[tuple, tuple[dict, datetime]] = {}
_TOOL_CONFIG_CACHE_TTL_SECONDS = 60

# Sensitive field keys that should be encrypted/decrypted
SENSITIVE_FIELD_KEYS = {"api_key", "private_key", "auth_code", "password", "secret", "atlassian_api_key"}


def _decrypt_sensitive_fields(config: dict, config_schema: dict | None = None) -> dict:
    """Decrypt sensitive fields in config dict.

    When config_schema is provided, also decrypts fields with type='password'
    (e.g. smithery_api_key) that are not in the hardcoded SENSITIVE_FIELD_KEYS.
    """
    if not config:
        return config

    from app.core.security import decrypt_data
    from app.config import get_settings

    settings = get_settings()
    result = dict(config)

    # Build the set of sensitive keys: hardcoded + schema-derived
    sensitive_keys = set(SENSITIVE_FIELD_KEYS)
    if config_schema:
        for field in config_schema.get("fields", []):
            if field.get("type") == "password":
                key = field.get("key", "")
                if key:
                    sensitive_keys.add(key)

    for key in sensitive_keys:
        if key in result and result[key]:
            value = result[key]
            if isinstance(value, str) and value:
                try:
                    result[key] = decrypt_data(value, settings.SECRET_KEY)
                except Exception:
                    # If decryption fails, assume it's plaintext
                    pass

    return result


def _get_cached_tool_config(agent_id: Optional[uuid.UUID], tool_name: str) -> Optional[dict]:
    """获取缓存的工具配置，过期返回 None。"""
    cache_key = (str(agent_id) if agent_id else None, tool_name)
    if cache_key in _tool_config_cache:
        config, expiry = _tool_config_cache[cache_key]
        if datetime.now() < expiry:
            return config
        # 过期，删除
        del _tool_config_cache[cache_key]
    return None


def _set_cached_tool_config(agent_id: Optional[uuid.UUID], tool_name: str, config: dict):
    """设置工具配置缓存。"""
    cache_key = (str(agent_id) if agent_id else None, tool_name)
    expiry = datetime.now() + timedelta(seconds=_TOOL_CONFIG_CACHE_TTL_SECONDS)
    _tool_config_cache[cache_key] = (config, expiry)


async def _get_tool_config(agent_id: Optional[uuid.UUID], tool_name: str) -> Optional[dict]:
    """Get merged tool config (with caching).

    Priority:
    1. agent_tools.config (per-agent override)
    2. tools.config (company/global config)

    Both configs are decrypted using the tool's config_schema for
    schema-aware field detection (e.g. smithery_api_key with type=password).
    """
    # Check cache first
    cached = _get_cached_tool_config(agent_id, tool_name)
    if cached is not None:
        logger.debug(f"[ToolConfig] Cache hit for {tool_name}, agent_id={agent_id}: {cached}")
        return cached

    from app.models.tool import Tool, AgentTool

    async with async_session() as db:
        # 1. Try per-agent + global config together
        if agent_id:
            result = await db.execute(
                select(AgentTool.config, Tool.config, Tool.config_schema)
                .join(Tool, AgentTool.tool_id == Tool.id)
                .where(AgentTool.agent_id == agent_id, Tool.name == tool_name)
            )
            row = result.first()
            if row:
                agent_config, global_config, config_schema = row
                # Merge: agent overrides global
                merged = {**(global_config or {}), **(agent_config or {})}
                if merged:
                    # Decrypt with schema awareness
                    merged = _decrypt_sensitive_fields(merged, config_schema)
                    logger.info(f"[ToolConfig] DB merged config for {tool_name}, agent_id={agent_id}")
                    _set_cached_tool_config(agent_id, tool_name, merged)
                    return merged

        # 2. Fallback to global config only
        result = await db.execute(select(Tool).where(Tool.name == tool_name))
        tool = result.scalar_one_or_none()
        if tool and tool.config:
            # Decrypt with schema awareness
            decrypted = _decrypt_sensitive_fields(tool.config, tool.config_schema)
            logger.info(f"[ToolConfig] DB global config for {tool_name}")
            _set_cached_tool_config(agent_id, tool_name, decrypted)
            return decrypted

    logger.error(f"[ToolConfig] No DB config found for {tool_name}, agent_id={agent_id}")
    return None

# ContextVar set by each channel handler so send_channel_file knows where to send
# Value: async callable(file_path: Path) -> None  |  None for web chat (returns URL)
channel_file_sender: ContextVar = ContextVar('channel_file_sender', default=None)
# For web chat: agent_id needed to build download URL
channel_web_agent_id: ContextVar = ContextVar('channel_web_agent_id', default=None)
# Set by Feishu channel handler — open_id of the message sender so calendar tool
# can auto-invite them as attendee when no explicit attendee list is given
channel_feishu_sender_open_id: ContextVar = ContextVar('channel_feishu_sender_open_id', default=None)

# ─── Tool Definitions (OpenAI function-calling format) ──────────

AGENT_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "List files and folders in a directory within my workspace. Can also list enterprise_info/ for shared company information.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Directory path to list, defaults to root (empty string). e.g.: '', 'skills', 'workspace', 'enterprise_info', 'enterprise_info/knowledge_base'",
                    }
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Read file contents from the workspace. Can read tasks.json for tasks, soul.md for personality, memory/memory.md for memory, skills/ for skill files, and enterprise_info/ for shared company info.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path, e.g.: tasks.json, soul.md, memory/memory.md, skills/xxx.md, enterprise_info/company_profile.md",
                    }
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "Write or update a file in the workspace. Can update memory/memory.md, focus.md, task_history.md, create documents in workspace/, create skills in skills/.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path, e.g.: memory/memory.md, workspace/report.md, skills/data_analysis.md",
                    },
                    "content": {
                        "type": "string",
                        "description": "File content to write",
                    },
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "delete_file",
            "description": "Delete a file from the workspace. Cannot delete soul.md or tasks.json.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path to delete",
                    }
                },
                "required": ["path"],
            },
        },
    },
    # --- Trigger management tools (Aware engine) ---
    {
        "type": "function",
        "function": {
            "name": "set_trigger",
            "description": "Set a new trigger to wake yourself up at a specific time or condition. Use this to schedule future actions, monitor changes, or wait for messages. The trigger will fire and invoke you with the reason text as context. Trigger types: 'cron' (recurring schedule), 'once' (fire once at a time), 'interval' (every N minutes), 'poll' (HTTP monitoring), 'on_message' (when another agent or a human user replies — use from_agent_name for agents, or from_user_name for human users on Feishu/Slack/Discord), 'webhook' (receive external HTTP POST — system generates a unique URL, give it to the user so they can configure it in external services like GitHub, Grafana, etc.).",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "Unique name for this trigger, e.g. 'daily_briefing' or 'wait_morty_reply'",
                    },
                    "type": {
                        "type": "string",
                        "enum": ["cron", "once", "interval", "poll", "on_message", "webhook"],
                        "description": "Trigger type",
                    },
                    "config": {
                        "type": "object",
                        "description": "Type-specific config. cron: {\"expr\": \"0 9 * * *\"}. once: {\"at\": \"2026-03-10T09:00:00+08:00\"}. interval: {\"minutes\": 30}. poll: {\"url\": \"...\", \"json_path\": \"$.status\", \"fire_on\": \"change\", \"interval_min\": 5}. on_message: {\"from_agent_name\": \"Morty\"} or {\"from_user_name\": \"张三\"} (for human users on Feishu/Slack/Discord). webhook: {\"secret\": \"optional_hmac_secret\"} (system auto-generates the URL)",
                    },
                    "reason": {
                        "type": "string",
                        "description": "What you should do when this trigger fires. This will be shown to you as context when you wake up.",
                    },
                    "focus_ref": {
                        "type": "string",
                        "description": "Optional: identifier of the focus item in focus.md that this trigger relates to (use the checklist identifier, e.g. 'daily_news_check')",
                    },
                },
                "required": ["name", "type", "config", "reason"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "update_trigger",
            "description": "Update an existing trigger's configuration or reason. Use this to adjust timing, change parameters, etc. For example, change interval from 5 minutes to 30 minutes.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "Name of the trigger to update",
                    },
                    "config": {
                        "type": "object",
                        "description": "New config (replaces existing config)",
                    },
                    "reason": {
                        "type": "string",
                        "description": "New reason text",
                    },
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "cancel_trigger",
            "description": "Cancel (disable) a trigger by name. Use this when a task is completed and the trigger is no longer needed.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "Name of the trigger to cancel",
                    },
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_triggers",
            "description": "List all your active triggers. Shows name, type, config, reason, fire count, and status.",
            "parameters": {
                "type": "object",
                "properties": {},
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_channel_file",
            "description": "Send a file to a specific person or back to the current conversation. If member_name is provided, the system resolves the recipient across all connected channels (Feishu, Slack, etc.) and delivers the file via the appropriate channel. If member_name is omitted, the file is sent back through the current conversation channel.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Workspace-relative path to the file, e.g. workspace/report.md",
                    },
                    "member_name": {
                        "type": "string",
                        "description": "Name of the person to send the file to. If provided, the system looks up this person across all configured channels and delivers via the appropriate one.",
                    },
                    "message": {
                        "type": "string",
                        "description": "Optional message to accompany the file",
                    },
                },
                "required": ["file_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_feishu_message",
            "description": (
                "Send a Feishu IM message to a colleague. "
                "You can provide either the colleague's name (will auto-search their open_id) "
                "or their open_id directly. "
                "To contact digital employees use send_message_to_agent instead."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "member_name": {
                        "type": "string",
                        "description": "Recipient's name, e.g. '覃睿'. Will be looked up automatically.",
                    },
                    "user_id": {
                        "type": "string",
                        "description": "Recipient's Feishu user_id (preferred, tenant-stable). Get from feishu_user_search.",
                    },
                    "open_id": {
                        "type": "string",
                        "description": "Recipient's Feishu open_id (fallback, per-app). Use user_id instead when available.",
                    },
                    "message": {
                        "type": "string",
                        "description": "Message content to send",
                    },
                },
                "required": ["message"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_channel_message",
            "description": (
                "Send a message to a colleague via their configured channel (Feishu, DingTalk, WeCom). "
                "Automatically detects the recipient's channel based on their org relationship. "
                "Use this as the primary method to send messages to colleagues in your relationship network."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "member_name": {
                        "type": "string",
                        "description": "Recipient's name as shown in relationships, e.g. '张三'. Must be a person in your relationship network.",
                    },
                    "message": {
                        "type": "string",
                        "description": "Message content to send",
                    },
                    "channel": {
                        "type": "string",
                        "description": "Optional: Specific channel to use (feishu, dingtalk, wecom). Use this if multiple people have the same name in different channels.",
                        "enum": ["feishu", "dingtalk", "wecom"]
                    },
                },
                "required": ["member_name", "message"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_web_message",
            "description": "Send a message to a user on the Clawith web platform. The message will appear in their web chat history and be pushed in real-time if they are online. Use this to proactively notify web users.",
            "parameters": {
                "type": "object",
                "properties": {
                    "username": {
                        "type": "string",
                        "description": "Username or display name of the recipient (must be a registered platform user)",
                    },
                    "message": {
                        "type": "string",
                        "description": "Message content to send",
                    },
                },
                "required": ["username", "message"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_message_to_agent",
            "description": "Send a message to a digital employee colleague and receive a reply. The recipient is another AI agent, not a human. This triggers the recipient's LLM reasoning and returns their response. Suitable for asking questions, delegating tasks, or collaboration. Your relationships.md lists available digital employees under 'Digital Employee Colleagues'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "agent_name": {
                        "type": "string",
                        "description": "Target digital employee's name",
                    },
                    "message": {
                        "type": "string",
                        "description": "Message content to send",
                    },
                    "msg_type": {
                        "type": "string",
                        "enum": ["notify", "consult", "task_delegate"],
                        "description": "Message type: notify (notification), consult (ask a question), task_delegate (delegate a task). Defaults to notify.",
                    },
                },
                "required": ["agent_name", "message"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_file_to_agent",
            "description": "Send a workspace file to another digital employee. The file is copied into the target agent's workspace/inbox/files/ directory and a delivery note is created in their inbox.",
            "parameters": {
                "type": "object",
                "properties": {
                    "agent_name": {
                        "type": "string",
                        "description": "Target digital employee's name",
                    },
                    "file_path": {
                        "type": "string",
                        "description": "Workspace-relative path of the source file, e.g. workspace/report.md",
                    },
                    "message": {
                        "type": "string",
                        "description": "Optional delivery note for the target digital employee",
                    },
                },
                "required": ["agent_name", "file_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "jina_search",
            "description": "Search the internet using Jina AI Search (s.jina.ai). Returns high-quality search results with full page content, not just snippets. Ideal for research, news, technical docs, and any real-time information lookup.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Search query, e.g. 'Python asyncio best practices' or '苏州通道人工智能科技有限公司'",
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Number of results to return, default 5, max 10",
                    },
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "jina_read",
            "description": "Read and extract the full content from a web page URL using Jina AI Reader (r.jina.ai). Returns clean, well-structured markdown including article text, tables, and key information. Better than jina_search when you already have a specific URL to read.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {
                        "type": "string",
                        "description": "The full URL of the web page to read, e.g. 'https://example.com/article'",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Max characters to return (default 8000, max 20000)",
                    },
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_document",
            "description": "Read office document contents (PDF, Word, Excel, PPT, etc.) and extract text. Suitable for reading knowledge base documents.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Document file path, e.g.: workspace/knowledge_base/report.pdf, enterprise_info/knowledge_base/policy.docx",
                    }
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "execute_code",
            "description": "Execute code (Python, Bash, or Node.js) in a sandboxed environment within the agent's root directory. Useful for data processing, calculations, file transformations, and automation scripts. Code runs with the agent root as the working directory, so you can access skills/, workspace/, memory/ etc. directly. Security restrictions apply: no network access commands, no system-level operations, 30-second timeout.",
            "parameters": {
                "type": "object",
                "properties": {
                    "language": {
                        "type": "string",
                        "enum": ["python", "bash", "node"],
                        "description": "Programming language to execute",
                    },
                    "code": {
                        "type": "string",
                        "description": "Code to execute. For Python, you can import standard libraries (json, csv, math, re, collections, etc.). Working directory is the agent root (skills/, workspace/, memory/ are accessible).",
                    },
                    "timeout": {
                        "type": "integer",
                        "description": "Max execution time in seconds (default 30, max 60)",
                    },
                },
                "required": ["language", "code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "upload_image",
            "description": "Upload an image file from your workspace (or from a public URL) to a cloud CDN and get a permanent public URL. Use this when you need to share images externally, embed them in messages/reports, or make workspace images accessible via URL. Supports common formats: PNG, JPG, GIF, WebP, SVG.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Workspace-relative path to the image file, e.g. workspace/chart.png or workspace/knowledge_base/diagram.jpg",
                    },
                    "url": {
                        "type": "string",
                        "description": "Alternative: a public URL of an image to upload (e.g. https://example.com/photo.jpg). Use this instead of file_path when the image is not in your workspace.",
                    },
                    "file_name": {
                        "type": "string",
                        "description": "Optional custom filename for the uploaded image. If omitted, the original filename is used.",
                    },
                    "folder": {
                        "type": "string",
                        "description": "Optional CDN folder path, e.g. /agents/reports. Defaults to /clawith.",
                    },
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "generate_image_siliconflow",
            "description": "Generate an image via SiliconFlow (FLUX). Save to workspace. Fast and China-friendly.",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "Detailed image description in English.",
                    },
                    "size": {
                        "type": "string",
                        "description": "Image size. Default: 1024x1024. Options: 1024x1024, 1024x768, 768x1024",
                    },
                    "save_path": {
                        "type": "string",
                        "description": "Workspace path to save the image (e.g. workspace/images/sunset.png).",
                    },
                },
                "required": ["prompt"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "generate_image_openai",
            "description": "Generate an image via OpenAI. Save to workspace.",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "Detailed image description in English.",
                    },
                    "size": {
                        "type": "string",
                        "description": "Image size. Default: 1024x1024.",
                    },
                    "save_path": {
                        "type": "string",
                        "description": "Workspace path to save the image.",
                    },
                },
                "required": ["prompt"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "generate_image_google",
            "description": "Generate an image via Google Gemini Image (Nano Banana) or Vertex AI. Save to workspace.",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "Detailed image description in English.",
                    },
                    "size": {
                        "type": "string",
                        "description": "Image size. Default: 1024x1024.",
                    },
                    "save_path": {
                        "type": "string",
                        "description": "Workspace path to save the image.",
                    },
                },
                "required": ["prompt"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "discover_resources",
            "description": "Search public MCP registries (Smithery) for tools and capabilities that can extend your abilities. Use this when you encounter a task you cannot handle with your current tools.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Semantic description of the capability needed, e.g. 'send email', 'query SQL database', 'generate images'",
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Max results to return (default 5, max 10)",
                    },
                },
                "required": ["query"],
            },
        },
    },
    # ── Feishu Bitable (多维表格) Tools ──────────────────────
    {
        "type": "function",
        "function": {
            "name": "bitable_list_tables",
            "description": "列出飞书多维表格内的所有数据表 (Tables)。url 支持表格链接或 Wiki 链接。使用此工具了解请求的多维表格中有哪些表。",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "多维表格的 URL 链接。"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bitable_list_fields",
            "description": "列出飞书多维表格指定数据表中的所有字段 (Fields)。url 支持表格链接或 Wiki 链接。在查询或修改数据前，必须先调用此工具了解字段名称和类型。",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "多维表格的 URL 链接。"},
                    "table_id": {"type": "string", "description": "具体的数据表 ID，如果 url 中包含 tbl 则可以不填。"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bitable_query_records",
            "description": "查询飞书多维表格中的数据行。可以提供过滤条件 (filter)。",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "多维表格的 URL 链接。"},
                    "table_id": {"type": "string", "description": "具体的数据表 ID，如果 url 中包含 tbl 则可以不填。"},
                    "filter_info": {"type": "string", "description": "可选，FQL 语法的过滤条件，例如 'CurrentValue.[Status]=\"Done\"'。如不确定过滤语法，可以不填，由你臺己在本地过滤返回的所有数据。"},
                    "max_results": {"type": "integer", "description": "最大返回条数 (默认 100)"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bitable_create_record",
            "description": "在飞书多维表格中新增一行数据。fields 参数是一个字典，key 是字段名 (需要先通过 bitable_list_fields 获取)，value 是对应的值。",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "多维表格的 URL 链接。"},
                    "table_id": {"type": "string", "description": "具体的数据表 ID，如果 url 中包含 tbl 则可以不填。"},
                    "fields": {"type": "string", "description": "一个 JSON 字符串，代表要插入的 fields。例如：'{\"Name\": \"张三\", \"Age\": 30}'"},
                },
                "required": ["url", "fields"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bitable_update_record",
            "description": "更新飞书多维表格中的指定行数据。",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "多维表格的 URL 链接。"},
                    "table_id": {"type": "string", "description": "具体的数据表 ID，如果 url 中包含 tbl 则可以不填。"},
                    "record_id": {"type": "string", "description": "要更新的 record_id，通过 bitable_query_records 获取。"},
                    "fields": {"type": "string", "description": "一个 JSON 字符串，代表要更新的 fields。例如：'{\"Status\": \"Done\"}'"},
                },
                "required": ["url", "record_id", "fields"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bitable_delete_record",
            "description": "删除飞书多维表格中的指定行数据。",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "多维表格的 URL 链接。"},
                    "table_id": {"type": "string", "description": "具体的数据表 ID，如果 url 中包含 tbl 则可以不填。"},
                    "record_id": {"type": "string", "description": "要删除的 record_id，通过 bitable_query_records 获取。"},
                },
                "required": ["url", "record_id"],
            },
        },
    },
    # ── Feishu Document Tools ──────────────────────────────────────
    {
        "type": "function",
        "function": {
            "name": "feishu_wiki_list",
            "description": (
                "List all sub-pages (child nodes) of a Feishu Wiki (知识库) page. "
                "Works with wiki URLs like 'https://xxx.feishu.cn/wiki/NodeToken'. "
                "Use this when a wiki page has child pages you need to explore. "
                "Returns titles, node_tokens, and obj_tokens for each sub-page. "
                "Each sub-page can then be read with feishu_doc_read using its node_token."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "node_token": {
                        "type": "string",
                        "description": "Wiki node token from the URL, e.g. 'HrGawgXxLiqoS5kT6pUczya3nEc' from 'https://xxx.feishu.cn/wiki/HrGawgXxLiqoS5kT6pUczya3nEc'",
                    },
                    "recursive": {
                        "type": "boolean",
                        "description": "If true, also list sub-pages of sub-pages (up to 2 levels deep). Default false.",
                    },
                },
                "required": ["node_token"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_doc_read",
            "description": (
                "Read the text content of a Feishu document or Wiki page. "
                "Works with both regular docx URLs (https://xxx.feishu.cn/docx/Token) "
                "and Wiki page URLs (https://xxx.feishu.cn/wiki/Token). "
                "Automatically handles wiki node tokens. "
                "If the page has sub-pages, use feishu_wiki_list to list them."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "document_token": {
                        "type": "string",
                        "description": "Feishu document token (from document URL)",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Max characters to return (default 6000, max 20000)",
                    },
                },
                "required": ["document_token"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_doc_create",
            "description": "Create a new Feishu document with a given title. Returns the new document token and URL.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Document title",
                    },
                    "folder_token": {
                        "type": "string",
                        "description": "Optional: parent folder token. Leave empty to create in root My Drive.",
                    },
                },
                "required": ["title"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_doc_append",
            "description": "Append text content to an existing Feishu document. Content is appended as one or more new paragraphs at the end.",
            "parameters": {
                "type": "object",
                "properties": {
                    "document_token": {
                        "type": "string",
                        "description": "Feishu document token",
                    },
                    "content": {
                        "type": "string",
                        "description": "Text content to append. Supports multiple lines separated by \\n.",
                    },
                },
                "required": ["document_token", "content"],
            },
        },
    },
    # ── Feishu Calendar Tools ──────────────────────────────────────
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_list",
            "description": "查询飞书日历。**自动读取当前对话用户的真实忙碌时段（freebusy）**，同时列出 bot 创建的日程。用于查询某人是否有空、安排日程时避开冲突。",
            "parameters": {
                "type": "object",
                "properties": {
                    "start_time": {
                        "type": "string",
                        "description": "查询起始时间，ISO 8601 格式，例如 '2026-03-13T00:00:00+08:00'。默认：当前时间。",
                    },
                    "end_time": {
                        "type": "string",
                        "description": "查询截止时间，ISO 8601 格式。默认：7天后。",
                    },
                    "user_open_id": {
                        "type": "string",
                        "description": "要查询 freebusy 的用户 open_id。不填则自动使用当前对话发送者。",
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Max events to return (default 20)",
                    },
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_create",
            "description": "Create a Feishu calendar event immediately. The current user is automatically invited as attendee — no email or authorization required. Just provide the title and time.",
            "parameters": {
                "type": "object",
                "properties": {
                    "summary": {
                        "type": "string",
                        "description": "Event title",
                    },
                    "start_time": {
                        "type": "string",
                        "description": "Event start in ISO 8601 with timezone, e.g. '2026-03-15T14:00:00+08:00'",
                    },
                    "end_time": {
                        "type": "string",
                        "description": "Event end in ISO 8601 with timezone, e.g. '2026-03-15T15:00:00+08:00'",
                    },
                    "description": {
                        "type": "string",
                        "description": "Event description or agenda",
                    },
                    "attendee_names": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Names of colleagues to invite, e.g. ['覃睿', '张三']. Will be looked up automatically via feishu_user_search.",
                    },
                    "attendee_open_ids": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Feishu open_ids to invite directly (if you already have them from feishu_user_search).",
                    },
                    "attendee_emails": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Additional attendee emails to invite (use attendee_names if you only have the name).",
                    },
                    "location": {
                        "type": "string",
                        "description": "Event location or meeting room",
                    },
                    "timezone": {
                        "type": "string",
                        "description": "Timezone, e.g. 'Asia/Shanghai'. Defaults to Asia/Shanghai.",
                    },
                },
                "required": ["summary", "start_time", "end_time"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_update",
            "description": "Update an existing Feishu calendar event. Provide only the fields you want to change.",
            "parameters": {
                "type": "object",
                "properties": {
                    "user_email": {"type": "string", "description": "Calendar owner's email"},
                    "event_id": {"type": "string", "description": "Event ID from feishu_calendar_list"},
                    "summary": {"type": "string", "description": "New title"},
                    "description": {"type": "string", "description": "New description"},
                    "start_time": {"type": "string", "description": "New start time (ISO 8601)"},
                    "end_time": {"type": "string", "description": "New end time (ISO 8601)"},
                    "location": {"type": "string", "description": "New location"},
                },
                "required": ["user_email", "event_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_delete",
            "description": "Delete (cancel) a Feishu calendar event.",
            "parameters": {
                "type": "object",
                "properties": {
                    "user_email": {"type": "string", "description": "Calendar owner's email"},
                    "event_id": {"type": "string", "description": "Event ID to delete"},
                },
                "required": ["user_email", "event_id"],
            },
        },
    },
    # ── Feishu Drive Share (collaborator management for all file types) ──
    {
        "type": "function",
        "function": {
            "name": "feishu_drive_share",
            "description": (
                "Manage Feishu Drive file collaborators and permissions. "
                "Supports ALL file types: docx, bitable, sheet, doc, folder, mindnote, slides. "
                "Can add or remove collaborators with viewer/editor/full_access roles, "
                "or get the current collaborator list. "
                "Accepts colleague names (auto-searched) or open_ids directly."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "document_token": {
                        "type": "string",
                        "description": "File token (from feishu_doc_create, bitable_create_app, or URL)",
                    },
                    "doc_type": {
                        "type": "string",
                        "enum": ["docx", "bitable", "sheet", "doc", "folder", "mindnote", "slides"],
                        "description": "File type. Default: 'docx'. Use 'bitable' for Bitable, 'sheet' for Spreadsheet, etc.",
                    },
                    "action": {
                        "type": "string",
                        "enum": ["add", "remove", "list"],
                        "description": "'add' to grant access, 'remove' to revoke, 'list' to view current collaborators",
                    },
                    "member_names": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Colleague names to add/remove, e.g. ['覃睿', '张三']. Auto-searched.",
                    },
                    "member_open_ids": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Feishu open_ids to add/remove directly (if already known).",
                    },
                    "permission": {
                        "type": "string",
                        "enum": ["view", "edit", "full_access"],
                        "description": "Permission level: 'view' (read-only), 'edit' (can edit), 'full_access' (can manage). Default: 'edit'",
                    },
                },
                "required": ["document_token", "action"],
            },
        },
    },
    # ── Feishu Drive Delete (delete files from cloud space) ──
    {
        "type": "function",
        "function": {
            "name": "feishu_drive_delete",
            "description": (
                "Delete a file or folder from Feishu Drive (cloud space). "
                "The file will be moved to the recycle bin, not permanently deleted. "
                "For folders, the deletion is asynchronous. "
                "Requires ownership + parent folder edit permission, or parent folder full_access."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_token": {
                        "type": "string",
                        "description": "Token of the file or folder to delete (from URL or previous tool output)",
                    },
                    "file_type": {
                        "type": "string",
                        "enum": ["file", "docx", "bitable", "folder", "doc", "sheet", "mindnote", "shortcut", "slides"],
                        "description": "Type of the file to delete. Use 'docx' for documents, 'bitable' for multitable, 'sheet' for spreadsheets, 'file' for uploaded files, 'folder' for folders.",
                    },
                },
                "required": ["file_token", "file_type"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_user_search",
            "description": (
                "Search for a colleague in the Feishu (Lark) directory by name. "
                "Returns their open_id, email, and department so you can send messages, "
                "invite them to calendar events, or share documents. "
                "Use this whenever you need to find a colleague's Feishu identity."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "The colleague's name to search for, e.g. '覃睿' or '张三'",
                    },
                },
                "required": ["name"],
            },
        },
    },
    # ── Feishu Approval Tools ──────────────────────────────────────
    {
        "type": "function",
        "function": {
            "name": "feishu_approval_create",
            "description": "发起一个飞书审批流实例。你需要知道审批定义的 approval_code 和表单对应字段的内容。",
            "parameters": {
                "type": "object",
                "properties": {
                    "approval_code": {
                        "type": "string",
                        "description": "审批定义的唯一代码 (approval_code)",
                    },
                    "user_id": {
                        "type": "string",
                        "description": "发起人的 open_id。可以通过 feishu_user_search 获取。",
                    },
                    "form_data": {
                        "type": "string",
                        "description": "表单内容的 JSON 字符串，例如 '[{\"id\":\"widget1\",\"type\":\"input\",\"value\":\"这是内容\"}]'",
                    },
                },
                "required": ["approval_code", "user_id", "form_data"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_approval_query",
            "description": "查询指定的飞书审批实例列表。可以支持按状态查询（PENDING, APPROVED, REJECTED, CANCELED, DELETED）。",
            "parameters": {
                "type": "object",
                "properties": {
                    "approval_code": {
                        "type": "string",
                        "description": "审批定义的唯一代码 (approval_code)",
                    },
                    "status": {
                        "type": "string",
                        "description": "可选过滤状态：PENDING, APPROVED, REJECTED, CANCELED, DELETED",
                    },
                },
                "required": ["approval_code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_approval_get",
            "description": "获取指定飞书审批实例的详细信息与当前审批状态。",
            "parameters": {
                "type": "object",
                "properties": {
                    "instance_id": {
                        "type": "string",
                        "description": "审批实例的 instance_id",
                    },
                },
                "required": ["instance_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "import_mcp_server",
            "description": "Import an MCP server from Smithery registry into the platform. The server's tools become available for use. Use discover_resources first to find the server ID. If previously imported tools stopped working (e.g. OAuth expired), set reauthorize=true to re-run the authorization flow.",
            "parameters": {
                "type": "object",
                "properties": {
                    "server_id": {
                        "type": "string",
                        "description": "Smithery server ID, e.g. '@anthropic/brave-search' or '@anthropic/fetch'",
                    },
                    "config": {
                        "type": "object",
                        "description": "Optional server configuration (e.g. API keys required by the server)",
                    },
                    "reauthorize": {
                        "type": "boolean",
                        "description": "Set to true to force re-authorization of existing tools (e.g. when OAuth token has expired)",
                    },
                },
                "required": ["server_id"],
            },
        },
    },
    # ─── Email Tools ────────────────────────
    {
        "type": "function",
        "function": {
            "name": "send_email",
            "description": "Send an email to one or more recipients. Supports subject, body text, CC, and file attachments from workspace. Requires email configuration in tool settings.",
            "parameters": {
                "type": "object",
                "properties": {
                    "to": {
                        "type": "string",
                        "description": "Recipient email address(es), comma-separated for multiple",
                    },
                    "subject": {
                        "type": "string",
                        "description": "Email subject line",
                    },
                    "body": {
                        "type": "string",
                        "description": "Email body text",
                    },
                    "cc": {
                        "type": "string",
                        "description": "CC recipients, comma-separated (optional)",
                    },
                    "attachments": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of workspace-relative file paths to attach (optional)",
                    },
                },
                "required": ["to", "subject", "body"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_emails",
            "description": "Read emails from your inbox. Can limit the number returned and search by criteria (e.g. FROM, SUBJECT, SINCE date). Requires email configuration in tool settings.",
            "parameters": {
                "type": "object",
                "properties": {
                    "limit": {
                        "type": "integer",
                        "description": "Max number of emails to return (default 10, max 30)",
                    },
                    "search": {
                        "type": "string",
                        "description": "IMAP search criteria, e.g. 'FROM \"john@example.com\"', 'SUBJECT \"meeting\"', 'SINCE 01-Mar-2026'. Default: all emails.",
                    },
                    "folder": {
                        "type": "string",
                        "description": "Mailbox folder, default INBOX",
                    },
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "reply_email",
            "description": "Reply to an email by its Message-ID. Maintains the email thread with proper In-Reply-To headers. Requires email configuration in tool settings.",
            "parameters": {
                "type": "object",
                "properties": {
                    "message_id": {
                        "type": "string",
                        "description": "Message-ID of the email to reply to (from read_emails output)",
                    },
                    "body": {
                        "type": "string",
                        "description": "Reply body text",
                    },
                },
                "required": ["message_id", "body"],
            },
        },
    },
    # --- Pages: public HTML hosting ---
    {
        "type": "function",
        "function": {
            "name": "publish_page",
            "description": "Publish an HTML file from workspace as a public page. Returns a public URL that anyone can access without login. Only .html/.htm files can be published.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path in workspace, e.g. 'workspace/output.html'",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_published_pages",
            "description": "List all pages published by this agent, showing their public URLs and view counts.",
            "parameters": {
                "type": "object",
                "properties": {},
            },
        },
    },
    # --- Skill Management ---
    {
        "type": "function",
        "function": {
            "name": "search_clawhub",
            "description": "Search the ClawHub skill registry for skills matching a query. Returns a list of available skills with name, description, and last updated date. Use this to help users find skills to install.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Search query, e.g. 'research', 'code review', 'market analysis'",
                    },
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "install_skill",
            "description": "Install a skill into this agent's workspace. Accepts either a ClawHub skill slug (e.g. 'market-research') or a GitHub URL (e.g. 'https://github.com/user/repo'). The skill files will be downloaded and saved to skills/<name>/ in your workspace.",
            "parameters": {
                "type": "object",
                "properties": {
                    "source": {
                        "type": "string",
                        "description": "ClawHub skill slug (e.g. 'market-research') or GitHub URL (e.g. 'https://github.com/user/repo')",
                    },
                },
                "required": ["source"],
            }
        }
    },
    # ── AgentBay Tools ────────────────────────────────────────────
    {
        "type": "function",
        "function": {
            "name": "agentbay_browser_navigate",
            "description": "使用 AgentBay 浏览器环境访问指定 URL。访问后会自动截图以便你观察当前页面状态。Tip: after navigating, use browser_observe to identify elements, then browser_type/browser_click to interact. IMPORTANT: Do NOT call navigate again after clicking or typing — that will refresh the page and lose all your progress. Use agentbay_browser_screenshot instead.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "要访问的网址，如 https://example.com"},
                    "wait_for": {"type": "string", "description": "等待特定元素出现的选择器（可选）"},
                    "save_to_workspace": {
                        "type": "boolean",
                        "description": "CRITICAL: Set to True IF AND ONLY IF the user explicitly asked you to SHOW them a screenshot or save it (e.g. \"截图给我看\", \"截图看看\", \"把截图发出来\"). If True, the image is saved to their workspace and you get a Markdown link. Default is False (internal in-memory analysis only, completely invisible to the user).",
                        "default": False,
                    },
                },
                "required": ["url"],
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "agentbay_browser_screenshot",
            "description": "Take a screenshot of the CURRENT browser page without navigating anywhere. Use this after clicking, typing, or submitting a form to verify the result — it preserves the current page state. Never call browser_navigate just to take a screenshot.",
            "parameters": {
                "type": "object",
                "properties": {
                    "save_to_workspace": {
                        "type": "boolean",
                        "description": "CRITICAL: Set to True IF AND ONLY IF the user explicitly asked you to SHOW them a screenshot or save it (e.g. \"截图给我看\", \"截图看看\", \"把截图发出来\"). If True, the image is saved to their workspace and you get a Markdown link. Default is False (internal in-memory analysis only, completely invisible to the user).",
                        "default": False,
                    },
                },
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "agentbay_browser_click",
            "description": "在 AgentBay 浏览器中点击指定元素。selector 可以是 CSS 选择器（如 #btn）或自然语言描述（如 'the Send button' 或 '发送验证码按钮'）。",
            "parameters": {
                "type": "object",
                "properties": {
                    "selector": {"type": "string", "description": "CSS selector (e.g. #button) or natural language description of the element (e.g. 'the blue Submit button')"},
                },
                "required": ["selector"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "agentbay_browser_type",
            "description": "在 AgentBay 浏览器的输入框中输入文本。selector 可以是 CSS 选择器或自然语言描述（如 'phone number input' 或 '手机号输入框'）。",
            "parameters": {
                "type": "object",
                "properties": {
                    "selector": {"type": "string", "description": "CSS selector or natural language description of the input field (e.g. 'the phone number input' or 'input[type=tel]')"},
                    "text": {"type": "string", "description": "要输入的文本"},
                },
                "required": ["selector", "text"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "agentbay_browser_login",
            "description": "Use AgentBay's AI-driven login skill to automate complex login flows (CAPTCHAs, OTP, multi-step auth). Requires a login_config JSON with AgentBay skill credentials. Navigate to the login page and execute the login skill.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "The login page URL to navigate to"},
                    "login_config": {"type": "string", "description": "JSON string with login config, e.g. '{\"api_key\": \"xxx\", \"skill_id\": \"yyy\"}'"},
                },
                "required": ["url", "login_config"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "agentbay_code_execute",
            "description": "在 AgentBay 代码空间中执行代码。支持 Python、Bash、Node.js。需要先配置 AgentBay 通道。",
            "parameters": {
                "type": "object",
                "properties": {
                    "language": {"type": "string", "enum": ["python", "bash", "node"], "description": "编程语言"},
                    "code": {"type": "string", "description": "要执行的代码"},
                    "timeout": {"type": "integer", "description": "超时时间（秒，默认 30）", "default": 30},
                },
                "required": ["language", "code"],
            },
        },
    },
]


# Core tools that should always be available to agents regardless of
# DB configuration.
_ALWAYS_INCLUDE_CORE = {
    "send_channel_file",
    "send_file_to_agent",
    "write_file",
    "send_channel_message",
}
# Channel message tool - available when any channel (Feishu/DingTalk/WeCom) is configured
_CHANNEL_MESSAGE_TOOL_NAMES = {
    "send_channel_message",
}
# Feishu tools are ONLY included when the agent has a configured Feishu channel,
# to avoid exposing unnecessary tools to non-Feishu agents (reduces hallucination risk).
_FEISHU_TOOL_NAMES = {
    "send_feishu_message",
    "feishu_user_search",
    "bitable_create_app",
    "bitable_list_tables",
    "bitable_list_fields",
    "bitable_query_records",
    "bitable_create_record",
    "bitable_update_record",
    "bitable_delete_record",
    "feishu_wiki_list",
    "feishu_doc_read",
    "feishu_doc_create",
    "feishu_doc_append",
    "feishu_drive_share",
    "feishu_drive_delete",
    "feishu_calendar_list",
    "feishu_calendar_create",
    "feishu_calendar_update",
    "feishu_calendar_delete",
    "feishu_approval_create",
    "feishu_approval_query",
    "feishu_approval_get",
}
_always_core_tools = [t for t in AGENT_TOOLS if t["function"]["name"] in _ALWAYS_INCLUDE_CORE]
_feishu_tools = [t for t in AGENT_TOOLS if t["function"]["name"] in _FEISHU_TOOL_NAMES]
_channel_tools = [t for t in AGENT_TOOLS if t["function"]["name"] in _CHANNEL_MESSAGE_TOOL_NAMES]


async def _agent_has_feishu(agent_id: uuid.UUID) -> bool:
    """Check if agent has a configured Feishu channel."""
    try:
        from app.models.channel_config import ChannelConfig
        async with async_session() as db:
            r = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "feishu",
                    ChannelConfig.is_configured == True,
                )
            )
            return r.scalar_one_or_none() is not None
    except Exception:
        return False


async def _agent_has_any_channel(agent_id: uuid.UUID) -> bool:
    """Check if agent has any configured channel (Feishu/DingTalk/WeCom)."""
    try:
        from app.models.channel_config import ChannelConfig
        async with async_session() as db:
            r = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.is_configured == True,
                )
            )
            return r.scalar_one_or_none() is not None
    except Exception:
        return False


# ─── Dynamic Tool Loading from DB ──────────────────────────────

async def get_agent_tools_for_llm(agent_id: uuid.UUID) -> list[dict]:
    """Load enabled tools for an agent from DB (OpenAI function-calling format).

    Falls back to hardcoded AGENT_TOOLS if DB not ready.
    Always includes core system tools (send_channel_file, write_file).
    Feishu tools are only included when the agent has a configured Feishu channel.
    send_channel_message is included when any channel (Feishu/DingTalk/WeCom) is configured.
    """
    has_feishu = await _agent_has_feishu(agent_id)
    has_any_channel = await _agent_has_any_channel(agent_id)
    _always_tools = _always_core_tools + (_feishu_tools if has_feishu else []) + (_channel_tools if has_any_channel else [])

    try:
        from app.models.tool import Tool, AgentTool

        async with async_session() as db:
            # Get all globally enabled tools
            all_tools_r = await db.execute(select(Tool).where(Tool.enabled == True))
            all_tools = all_tools_r.scalars().all()

            # Get agent-specific assignments
            agent_tools_r = await db.execute(select(AgentTool).where(AgentTool.agent_id == agent_id))
            assignments = {str(at.tool_id): at for at in agent_tools_r.scalars().all()}

            result = []
            db_tool_names = set()
            for t in all_tools:
                tid = str(t.id)
                at = assignments.get(tid)
                enabled = at.enabled if at else t.is_default
                if not enabled:
                    continue

                # Skip feishu tools if the agent has no Feishu channel configured
                if t.category == "feishu" and not has_feishu:
                    continue

                # Build OpenAI function-calling format
                tool_def = {
                    "type": "function",
                    "function": {
                        "name": t.name,
                        "description": t.description,
                        "parameters": t.parameters_schema or {"type": "object", "properties": {}},
                    },
                }
                result.append(tool_def)
                db_tool_names.add(t.name)

            if result:
                # Append always-available system tools that aren't already in the DB list
                for t in _always_tools:
                    if t["function"]["name"] not in db_tool_names:
                        result.append(t)
                return result
    except Exception as e:
        logger.error(f"[Tools] DB load failed, using fallback: {e}")

    # Fallback to hardcoded tools
    return AGENT_TOOLS


# ─── Workspace initialization ──────────────────────────────────

async def ensure_workspace(agent_id: uuid.UUID, tenant_id: str | None = None) -> Path:
    """Initialize agent workspace with standard structure."""
    ws = WORKSPACE_ROOT / str(agent_id)
    ws.mkdir(parents=True, exist_ok=True)

    # Create standard directories
    (ws / "skills").mkdir(exist_ok=True)
    (ws / "workspace").mkdir(exist_ok=True)
    (ws / "workspace" / "knowledge_base").mkdir(exist_ok=True)
    (ws / "memory").mkdir(exist_ok=True)

    # Ensure tenant-scoped enterprise_info directory exists
    if tenant_id:
        enterprise_dir = WORKSPACE_ROOT / f"enterprise_info_{tenant_id}"
    else:
        enterprise_dir = WORKSPACE_ROOT / "enterprise_info"
    enterprise_dir.mkdir(parents=True, exist_ok=True)
    (enterprise_dir / "knowledge_base").mkdir(exist_ok=True)
    # Create default company profile if missing
    profile_path = enterprise_dir / "company_profile.md"
    if not profile_path.exists():
        profile_path.write_text("# Company Profile\n\n_Edit company information here. All digital employees can access this._\n\n## Basic Info\n- Company Name:\n- Industry:\n- Founded:\n\n## Business Overview\n\n## Organization Structure\n\n## Company Culture\n", encoding="utf-8")

    # Migrate: move root-level memory.md into memory/ directory
    if (ws / "memory.md").exists() and not (ws / "memory" / "memory.md").exists():
        import shutil
        shutil.move(str(ws / "memory.md"), str(ws / "memory" / "memory.md"))

    # Create default memory file if missing
    if not (ws / "memory" / "memory.md").exists():
        (ws / "memory" / "memory.md").write_text("# Memory\n\n_Record important information and knowledge here._\n", encoding="utf-8")

    if not (ws / "soul.md").exists():
        # Try to load from DB
        try:
            async with async_session() as db:

                r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
                agent = r.scalar_one_or_none()
                if agent and agent.role_description:
                    (ws / "soul.md").write_text(
                        f"# Personality\n\n{agent.role_description}\n",
                        encoding="utf-8",
                    )
                else:
                    (ws / "soul.md").write_text("# Personality\n\n_Describe your role and responsibilities._\n", encoding="utf-8")
        except Exception:
            (ws / "soul.md").write_text("# Personality\n\n_Describe your role and responsibilities._\n", encoding="utf-8")

    # Always sync tasks from DB
    await _sync_tasks_to_file(agent_id, ws)

    return ws


async def _sync_tasks_to_file(agent_id: uuid.UUID, ws: Path):
    """Sync tasks from DB to tasks.json in workspace."""
    try:
        async with async_session() as db:
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id).order_by(Task.created_at.desc())
            )
            tasks = result.scalars().all()

        task_list = []
        for t in tasks:
            task_list.append({
                "title": t.title,
                "status": t.status,
                "priority": t.priority,
                "description": t.description or "",
                "created_at": t.created_at.isoformat() if t.created_at else "",
                "completed_at": t.completed_at.isoformat() if t.completed_at else "",
            })

        (ws / "tasks.json").write_text(
            json.dumps(task_list, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    except Exception as e:
        logger.error(f"[AgentTools] Failed to sync tasks: {e}")


# ─── Tool Executors ─────────────────────────────────────────────

# Mapping from tool_name to autonomy action_type
_TOOL_AUTONOMY_MAP = {
    "write_file": "write_workspace_files",
    "delete_file": "delete_files",
    "send_feishu_message": "send_feishu_message",
    "send_message_to_agent": "send_feishu_message",
    "send_file_to_agent": "send_feishu_message",
    "web_search": "web_search",
    "execute_code": "execute_code",
}


async def _get_agent_tenant_id(agent_id: uuid.UUID) -> str | None:
    """Get the agent tenant ID for tenant-scoped shared paths."""
    try:
        async with async_session() as db:

            r = await db.execute(select(AgentModel.tenant_id).where(AgentModel.id == agent_id))

            tenant_id = r.scalar_one_or_none()
            if tenant_id:
                return str(tenant_id)
    except Exception:
        pass
    return None


async def _execute_tool_direct(
    tool_name: str,
    arguments: dict,
    agent_id: uuid.UUID,
) -> str:
    """Execute a tool directly, bypassing autonomy checks.

    Used by the approval post-processing hook after an action
    has been approved and needs to actually run.
    """
    _agent_tenant_id = await _get_agent_tenant_id(agent_id)
    ws = await ensure_workspace(agent_id, tenant_id=_agent_tenant_id)
    try:
        if tool_name == "delete_file":
            return _delete_file(ws, arguments.get("path", ""))
        elif tool_name == "write_file":
            path = arguments.get("path")
            content = arguments.get("content", "")
            if not path:
                return "Missing path"
            return _write_file(ws, path, content, tenant_id=_agent_tenant_id)
        elif tool_name == "execute_code":
            logger.info(f"[DirectTool] Executing code with arguments: {arguments}")
            return await _execute_code(agent_id, ws, arguments)
        elif tool_name == "web_search":
            return await _web_search(arguments, agent_id)
        elif tool_name == "jina_search":
            return await _jina_search(arguments)
        elif tool_name == "send_feishu_message":
            return await _send_feishu_message(agent_id, arguments)
        elif tool_name == "send_message_to_agent":
            return await _send_message_to_agent(agent_id, arguments)
        elif tool_name == "send_file_to_agent":
            return await _send_file_to_agent(agent_id, ws, arguments)
        else:
            return f"Tool {tool_name} does not support post-approval execution"
    except Exception as e:
        logger.exception(f"[DirectTool] Error executing {tool_name}: {e}")
        return f"Error executing {tool_name}: {e}"


async def execute_tool(
    tool_name: str,
    arguments: dict,
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    session_id: str = "",
) -> str:
    """Execute a tool call and return the result as a string.

    Args:
        session_id: The ChatSession ID, used to isolate AgentBay instances
                    per conversation. Passed through to agentbay_* tools.
    """
    _agent_tenant_id = await _get_agent_tenant_id(agent_id)

    ws = await ensure_workspace(agent_id, tenant_id=_agent_tenant_id)

    # ── Autonomy boundary check ──
    action_type = _TOOL_AUTONOMY_MAP.get(tool_name)
    if action_type:
        try:
            from app.services.autonomy_service import autonomy_service
            from app.models.agent import Agent as AgentModel
            async with async_session() as _adb:
                _ar = await _adb.execute(select(AgentModel).where(AgentModel.id == agent_id))
                _agent = _ar.scalar_one_or_none()
                if _agent:
                    result_check = await autonomy_service.check_and_enforce(
                        _adb, _agent, action_type, {"tool": tool_name, "args": str(arguments)[:200], "requested_by": str(user_id)}
                    )
                    await _adb.commit()
                    if not result_check.get("allowed"):
                        level = result_check.get("level", "L3")
                        logger.info(f"[Autonomy] Tool {tool_name} denied, level: {level}")
                        if level == "L3":
                            return f"⏳ This action requires approval. An approval request has been sent. Please wait for approval before retrying. (Approval ID: {result_check.get('approval_id', 'N/A')})"
                        return f"❌ Action denied: {result_check.get('message', 'unknown reason')}"
        except Exception as e:
            logger.exception(f"[Autonomy] Check failed: {e}")
            return f"⚠️ Autonomy check failed ({e}). Operation blocked for safety. Please retry or contact admin."

    # Pre-inject session_id into arguments for AgentBay tools so each
    # _agentbay_* handler can pass it to get_agentbay_client_for_agent()
    # for per-ChatSession isolation of cloud instances.
    if tool_name.startswith("agentbay_"):
        arguments["_session_id"] = session_id

        # Take Control lock: block automatic tool execution while a human
        # is manually controlling the browser/desktop session. This prevents
        # input collisions between human clicks and agent-initiated actions.
        from app.api.agentbay_control import is_session_locked
        if is_session_locked(str(agent_id), session_id):
            return (
                "⏸️ A human operator is currently controlling this browser session "
                "(Take Control mode). Please wait for them to finish before retrying "
                "browser/computer operations."
            )

    try:
        if tool_name == "list_files":
            result = _list_files(ws, arguments.get("path", ""), tenant_id=_agent_tenant_id)
        elif tool_name == "read_file":
            path = arguments.get("path")
            if not path:
                return "❌ Missing required argument 'path' for read_file"
            result = _read_file(ws, path, tenant_id=_agent_tenant_id)
        elif tool_name == "read_document":
            path = arguments.get("path")
            if not path:
                return "❌ Missing required argument 'path' for read_document"
            max_chars = min(int(arguments.get("max_chars", 8000)), 20000)
            result = await _read_document(ws, path, max_chars=max_chars, tenant_id=_agent_tenant_id)
        elif tool_name == "write_file":
            path = arguments.get("path")
            content = arguments.get("content")
            if not path:
                return "❌ Missing required argument 'path' for write_file. Please provide a file path like 'skills/my-skill/SKILL.md'"
            if content is None:
                return "❌ Missing required argument 'content' for write_file"
            result = _write_file(ws, path, content, tenant_id=_agent_tenant_id)
        elif tool_name == "delete_file":
            result = _delete_file(ws, arguments.get("path", ""))
        elif tool_name == "manage_tasks":
            result = await _manage_tasks(agent_id, user_id, ws, arguments)
        elif tool_name == "set_trigger":
            result = await _handle_set_trigger(agent_id, arguments)
        elif tool_name == "update_trigger":
            result = await _handle_update_trigger(agent_id, arguments)
        elif tool_name == "cancel_trigger":
            result = await _handle_cancel_trigger(agent_id, arguments)
        elif tool_name == "list_triggers":
            result = await _handle_list_triggers(agent_id)
        elif tool_name == "send_feishu_message":
            result = await _send_feishu_message(agent_id, arguments)
        elif tool_name == "send_web_message":
            result = await _send_web_message(agent_id, arguments)
        elif tool_name == "send_channel_message":
            result = await _send_channel_message(agent_id, arguments)
        elif tool_name == "send_message_to_agent":
            result = await _send_message_to_agent(agent_id, arguments)
        elif tool_name == "send_file_to_agent":
            result = await _send_file_to_agent(agent_id, ws, arguments)
        elif tool_name == "send_channel_file":
            result = await _send_channel_file(agent_id, ws, arguments)
        elif tool_name == "web_search":
            result = await _web_search(arguments, agent_id)
        elif tool_name == "jina_search":
            result = await _jina_search(arguments)
        elif tool_name == "bing_search":
            result = await _jina_search(arguments)  # redirect legacy to jina
        elif tool_name == "jina_read":
            result = await _jina_read(arguments)
        elif tool_name == "read_webpage":
            result = await _jina_read(arguments)  # redirect legacy to jina
        elif tool_name == "plaza_get_new_posts":
            result = await _plaza_get_new_posts(agent_id, arguments)
        elif tool_name == "plaza_create_post":
            result = await _plaza_create_post(agent_id, arguments)
        elif tool_name == "plaza_add_comment":
            result = await _plaza_add_comment(agent_id, arguments)
        elif tool_name == "execute_code":
            logger.info(f"[DirectTool] Executing code with arguments: {arguments}")
            result = await _execute_code(agent_id, ws, arguments)
        elif tool_name == "upload_image":
            result = await _upload_image(agent_id, ws, arguments)
        elif tool_name == "generate_image_siliconflow":
            result = await _generate_image(agent_id, ws, arguments, "siliconflow")
        elif tool_name == "generate_image_openai":
            result = await _generate_image(agent_id, ws, arguments, "openai")
        elif tool_name == "generate_image_google":
            result = await _generate_image(agent_id, ws, arguments, "google")
        elif tool_name == "discover_resources":
            result = await _discover_resources(arguments)
        elif tool_name == "import_mcp_server":
            result = await _import_mcp_server(agent_id, arguments)
        # ── Feishu Bitable Tools ──
        elif tool_name == "bitable_create_app":
            result = await _bitable_create_app(agent_id, arguments)
        elif tool_name == "bitable_list_tables":
            result = await _bitable_list_tables(agent_id, arguments)
        elif tool_name == "bitable_list_fields":
            result = await _bitable_list_fields(agent_id, arguments)
        elif tool_name == "bitable_query_records":
            result = await _bitable_query_records(agent_id, arguments)
        elif tool_name == "bitable_create_record":
            result = await _bitable_create_record(agent_id, arguments)
        elif tool_name == "bitable_update_record":
            result = await _bitable_update_record(agent_id, arguments)
        elif tool_name == "bitable_delete_record":
            result = await _bitable_delete_record(agent_id, arguments)
        # ── Feishu Document Tools ──
        elif tool_name == "feishu_wiki_list":
            result = await _feishu_wiki_list(agent_id, arguments)
        elif tool_name == "feishu_doc_read":
            result = await _feishu_doc_read(agent_id, arguments)
        elif tool_name == "feishu_doc_create":
            result = await _feishu_doc_create(agent_id, arguments)
        elif tool_name == "feishu_doc_append":
            result = await _feishu_doc_append(agent_id, arguments)
        # ── Feishu Calendar Tools ──
        elif tool_name == "feishu_drive_share":
            result = await _feishu_drive_share(agent_id, arguments)
        elif tool_name == "feishu_drive_delete":
            result = await _feishu_drive_delete(agent_id, arguments)
        elif tool_name == "feishu_user_search":
            result = await _feishu_user_search(agent_id, arguments)
        elif tool_name == "feishu_calendar_list":
            result = await _feishu_calendar_list(agent_id, arguments)
        elif tool_name == "feishu_calendar_create":
            result = await _feishu_calendar_create(agent_id, arguments)
        elif tool_name == "feishu_calendar_update":
            result = await _feishu_calendar_update(agent_id, arguments)
        elif tool_name == "feishu_calendar_delete":
            result = await _feishu_calendar_delete(agent_id, arguments)
        elif tool_name == "feishu_approval_create":
            result = await _feishu_approval_create(agent_id, arguments)
        elif tool_name == "feishu_approval_query":
            result = await _feishu_approval_query(agent_id, arguments)
        elif tool_name == "feishu_approval_get":
            result = await _feishu_approval_get(agent_id, arguments)
        # ── Email Tools ──
        elif tool_name in ("send_email", "read_emails", "reply_email"):
            result = await _handle_email_tool(tool_name, agent_id, ws, arguments)
        # ── Pages: public HTML hosting ──
        elif tool_name == "publish_page":
            result = await _publish_page(agent_id, user_id, ws, arguments)
        elif tool_name == "list_published_pages":
            result = await _list_published_pages(agent_id)
        # ── AgentBay Tools ──
        elif tool_name == "agentbay_browser_navigate":
            result = await _agentbay_browser_navigate(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_screenshot":
            result = await _agentbay_browser_screenshot(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_click":
            result = await _agentbay_browser_click(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_type":
            result = await _agentbay_browser_type(agent_id, ws, arguments)
        elif tool_name == "agentbay_code_execute":
            result = await _agentbay_code_execute(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_extract":
            result = await _agentbay_browser_extract(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_observe":
            result = await _agentbay_browser_observe(agent_id, ws, arguments)
        elif tool_name == "agentbay_browser_login":
            result = await _agentbay_browser_login(agent_id, ws, arguments)
        elif tool_name == "agentbay_command_exec":
            result = await _agentbay_command_exec(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_screenshot":
            result = await _agentbay_computer_screenshot(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_click":
            result = await _agentbay_computer_click(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_input_text":
            result = await _agentbay_computer_input_text(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_press_keys":
            result = await _agentbay_computer_press_keys(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_scroll":
            result = await _agentbay_computer_scroll(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_move_mouse":
            result = await _agentbay_computer_move_mouse(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_drag_mouse":
            result = await _agentbay_computer_drag_mouse(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_get_screen_size":
            result = await _agentbay_computer_get_screen_size(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_start_app":
            result = await _agentbay_computer_start_app(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_get_cursor_position":
            result = await _agentbay_computer_get_cursor_position(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_get_active_window":
            result = await _agentbay_computer_get_active_window(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_activate_window":
            result = await _agentbay_computer_activate_window(agent_id, ws, arguments)
        elif tool_name == "agentbay_computer_list_visible_apps":
            result = await _agentbay_computer_list_visible_apps(agent_id, ws, arguments)
        # ── Skill Management ──
        elif tool_name == "search_clawhub":
            result = await _search_clawhub(agent_id, arguments)
        elif tool_name == "install_skill":
            result = await _install_skill(agent_id, ws, arguments)
        else:
            # Try MCP tool execution
            result = await _execute_mcp_tool(tool_name, arguments, agent_id=agent_id)

        # Log tool call activity (skip noisy read operations)
        if tool_name not in ("list_files", "read_file", "read_document"):
            from app.services.activity_logger import log_activity
            await log_activity(
                agent_id, "tool_call",
                f"Called tool {tool_name}: {result[:80]}",
                detail={"tool": tool_name, "args": {k: str(v)[:100] for k, v in arguments.items()}, "result": result[:300]},
            )
        return result
    except Exception as e:
        logger.exception(f"[Tool] Execution failed: {tool_name}")
        return f"Tool execution error ({tool_name}): {type(e).__name__}: {str(e)[:200]}"


async def _web_search(arguments: dict, agent_id: uuid.UUID | None = None) -> str:
    """Search the web using a configurable search engine.

    Config resolution priority: Agent config > Company config > Defaults.
    """
    import httpx
    import re

    query = arguments.get("query", "")
    if not query:
        return "❌ Please provide search keywords"

    # Use the standard _get_tool_config helper (Agent > Company, cached, decrypted)
    config = await _get_tool_config(agent_id, "web_search") or {}

    engine = config.get("search_engine", "duckduckgo")
    api_key = config.get("api_key", "")
    max_results = min(arguments.get("max_results", config.get("max_results", 5)), 10)
    language = config.get("language", "zh-CN")

    try:
        if engine == "tavily" and api_key:
            return await _search_tavily(query, api_key, max_results)
        elif engine == "google" and api_key:
            return await _search_google(query, api_key, max_results, language)
        elif engine == "bing" and api_key:
            return await _search_bing(query, api_key, max_results, language)
        else:
            return await _search_duckduckgo(query, max_results)
    except Exception as e:
        return f"❌ Search error ({engine}): {str(e)[:200]}"


async def _search_duckduckgo(query: str, max_results: int) -> str:
    """Search via DuckDuckGo HTML (free, no API key)."""
    import httpx, re

    async with httpx.AsyncClient(follow_redirects=True) as client:
        resp = await client.get(
            "https://html.duckduckgo.com/html/",
            params={"q": query},
            headers={"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"},
            timeout=10,
        )

    results = []
    blocks = re.findall(
        r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>.*?'
        r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>',
        resp.text, re.DOTALL,
    )
    for url, title, snippet in blocks[:max_results]:
        title = re.sub(r'<[^>]+>', '', title).strip()
        snippet = re.sub(r'<[^>]+>', '', snippet).strip()
        if "uddg=" in url:
            from urllib.parse import unquote, parse_qs, urlparse
            parsed = parse_qs(urlparse(url).query)
            url = unquote(parsed.get("uddg", [url])[0])
        results.append(f"**{title}**\n{url}\n{snippet}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 DuckDuckGo results for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)

async def _get_jina_api_key() -> str:
    """Read Jina API key from DB system_settings first, then fall back to env."""
    try:
        from app.database import async_session
        from app.models.system_settings import SystemSetting
        from sqlalchemy import select
        async with async_session() as db:
            result = await db.execute(select(SystemSetting).where(SystemSetting.key == "jina_api_key"))
            setting = result.scalar_one_or_none()
            if setting and setting.value.get("api_key"):
                return setting.value["api_key"]
    except Exception:
        pass
    from app.config import get_settings
    return get_settings().JINA_API_KEY


async def _jina_search(arguments: dict) -> str:
    """Search via Jina AI Search API (s.jina.ai). Returns full content per result, not just snippets."""
    import httpx

    query = arguments.get("query", "").strip()
    if not query:
        return "❌ Please provide search keywords"

    max_results = min(arguments.get("max_results", 5), 10)
    api_key = await _get_jina_api_key()

    headers: dict = {
        "Accept": "application/json",
        "X-Respond-With": "no-content",  # return snippets/descriptions, not full pages (faster)
        "X-Return-Format": "markdown",
    }
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
            resp = await client.get(
                f"https://s.jina.ai/{__import__('urllib.parse', fromlist=['quote']).quote(query)}",
                headers=headers,
            )

        if resp.status_code != 200:
            return f"❌ Jina Search error HTTP {resp.status_code}: {resp.text[:200]}"

        data = resp.json()
        items = data.get("data", [])[:max_results]

        if not items:
            return f'🔍 No results found for "{query}"'

        parts = []
        for i, item in enumerate(items, 1):
            title = item.get("title", "Untitled")
            url = item.get("url", "")
            description = item.get("description", "") or item.get("content", "")[:500]
            parts.append(f"**{i}. {title}**\n{url}\n{description}")

        return f'🔍 Jina Search results for "{query}" ({len(items)} items):\n\n' + "\n\n---\n\n".join(parts)

    except Exception as e:
        return f"❌ Jina Search error: {str(e)[:300]}"


async def _jina_read(arguments: dict) -> str:
    """Read web page via Jina AI Reader API (r.jina.ai). Returns clean structured markdown."""
    import httpx
    from app.config import get_settings

    url = arguments.get("url", "").strip()
    if not url:
        return "❌ Please provide a URL"
    if not url.startswith("http"):
        url = "https://" + url

    max_chars = min(arguments.get("max_chars", 8000), 20000)
    api_key = await _get_jina_api_key()

    headers: dict = {
        "Accept": "text/plain, text/markdown, */*",
        "X-Return-Format": "markdown",
        "X-Remove-Selector": "header, footer, nav, aside, .ads, .advertisement",
    }
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
            resp = await client.get(
                f"https://r.jina.ai/{url}",
                headers=headers,
            )

        if resp.status_code != 200:
            return f"❌ Jina Reader error HTTP {resp.status_code}: {resp.text[:200]}"

        text = resp.text.strip()
        if not text or len(text) < 100:
            return f"❌ Jina Reader returned empty content for {url}"

        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n[... truncated at {max_chars} chars]"

        return f"📄 **Content from: {url}**\n\n{text}"

    except Exception as e:
        return f"❌ Jina Reader error: {str(e)[:300]}"



async def _search_tavily(query: str, api_key: str, max_results: int) -> str:
    """Search via Tavily API (AI-optimized search)."""
    import httpx

    async with httpx.AsyncClient() as client:
        resp = await client.post(
            "https://api.tavily.com/search",
            json={"query": query, "max_results": max_results, "search_depth": "basic"},
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            timeout=15,
        )
        data = resp.json()

    if "results" not in data:
        return f"❌ Tavily search failed: {data.get('error', str(data)[:200])}"

    results = []
    for r in data["results"][:max_results]:
        results.append(f"**{r.get('title', '')}**\n{r.get('url', '')}\n{r.get('content', '')[:200]}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 Tavily search for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)


async def _search_google(query: str, api_key: str, max_results: int, language: str) -> str:
    """Search via Google Custom Search JSON API."""
    import httpx

    # api_key format: "API_KEY:CX_ID"
    parts = api_key.split(":", 1)
    if len(parts) != 2:
        return "❌ Google search requires API key in format 'API_KEY:SEARCH_ENGINE_ID'"

    gapi_key, cx = parts
    async with httpx.AsyncClient() as client:
        resp = await client.get(
            "https://www.googleapis.com/customsearch/v1",
            params={"key": gapi_key, "cx": cx, "q": query, "num": max_results, "lr": f"lang_{language[:2]}"},
            timeout=10,
        )
        data = resp.json()

    results = []
    for item in data.get("items", [])[:max_results]:
        results.append(f"**{item.get('title', '')}**\n{item.get('link', '')}\n{item.get('snippet', '')}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 Google search for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)


async def _search_bing(query: str, api_key: str, max_results: int, language: str) -> str:
    """Search via Bing Web Search API."""
    import httpx

    async with httpx.AsyncClient() as client:
        resp = await client.get(
            "https://api.bing.microsoft.com/v7.0/search",
            params={"q": query, "count": max_results, "mkt": language},
            headers={"Ocp-Apim-Subscription-Key": api_key},
            timeout=10,
        )
        data = resp.json()

    results = []
    for item in data.get("webPages", {}).get("value", [])[:max_results]:
        results.append(f"**{item.get('name', '')}**\n{item.get('url', '')}\n{item.get('snippet', '')}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 Bing search for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)


async def _send_channel_file(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Send a file to a person or back to the current channel.
    
    Priority:
    1. If member_name is provided, resolve the recipient across all configured channels
       and deliver via the appropriate one (Feishu, Slack, etc.).
    2. If channel_file_sender ContextVar is set (channel-initiated), use it directly.
    3. Fall back to web chat download URL when no explicit recipient is requested.
    """
    rel_path = arguments.get("file_path", "").strip()
    accompany_msg = arguments.get("message", "")
    member_name = (arguments.get("member_name") or "").strip()
    if not rel_path:
        return "Error: file_path is required"

    # Resolve file path within agent workspace
    file_path = (ws / rel_path).resolve()
    ws_resolved = ws.resolve()
    if not str(file_path).startswith(str(ws_resolved)):
        file_path = (WORKSPACE_ROOT / str(agent_id) / rel_path).resolve()
        if not file_path.exists():
            return f"Error: File not found: {rel_path}"
    if not file_path.exists():
        return f"Error: File not found: {rel_path}"

    # Priority 1: explicit recipient - resolve member across channels
    if member_name:
        result = await _send_file_to_recipient(agent_id, file_path, member_name, accompany_msg)
        if result:
            return result
        return (
            f"Failed to send file to '{member_name}': recipient not reachable via configured channels. "
            "Use send_message_to_agent for digital employees, or omit member_name to return a download link."
        )

    # Priority 2: channel-initiated (ContextVar set by channel webhook handler)
    sender = channel_file_sender.get()
    if sender is not None:
        try:
            await sender(file_path, accompany_msg)
            return f"File '{file_path.name}' sent to user via channel."
        except Exception as e:
            return f"Failed to send file: {e}"

    # Priority 3: Web chat fallback — return download URL
    aid = channel_web_agent_id.get() or str(agent_id)
    base_abs = (WORKSPACE_ROOT / str(agent_id)).resolve()
    try:
        file_rel = str(file_path.resolve().relative_to(base_abs))
    except ValueError:
        file_rel = rel_path
    from app.config import get_settings as _gs
    _s = _gs()
    base_url = getattr(_s, 'BASE_URL', '').rstrip('/') or ''
    download_url = f"{base_url}/api/agents/{aid}/files/download?path={file_rel}"
    msg = f"File ready: [{file_path.name}]({download_url})"
    if accompany_msg:
        msg = accompany_msg + "\n\n" + msg
    return msg


async def _send_file_to_recipient(
    agent_id: uuid.UUID, file_path: Path, member_name: str, message: str = ""
) -> str | None:
    """Resolve a recipient by name and send file via their reachable channel.
    
    Checks Feishu and Slack channels configured for this agent.
    Returns a result string, or None if no channel found.
    """
    from app.models.channel_config import ChannelConfig

    async with async_session() as db:
        # Load all channel configs for this agent
        result = await db.execute(
            select(ChannelConfig).where(ChannelConfig.agent_id == agent_id)
        )
        configs = {c.channel_type: c for c in result.scalars().all()}

    # --- Try Feishu ---
    feishu_config = configs.get("feishu")
    if feishu_config:
        feishu_result = await _send_file_via_feishu(agent_id, feishu_config, file_path, member_name, message)
        if feishu_result:
            return feishu_result

    # --- Try Slack ---
    slack_config = configs.get("slack")
    if slack_config:
        slack_result = await _send_file_via_slack(agent_id, slack_config, file_path, member_name, message)
        if slack_result:
            return slack_result

    return None  # No channel could reach this recipient


async def _resolve_feishu_recipient(agent_id: uuid.UUID, config, member_name: str) -> tuple[str, str] | None:
    """Resolve a Feishu recipient by name. Returns (receive_id, id_type) or None."""
    # 1. Try feishu_user_search (checks cache, OrgMember, User table)
    import re as _re
    search_result = await _feishu_user_search(agent_id, {"name": member_name})
    
    uid_match = _re.search(r'user_id: `([A-Za-z0-9]+)`', search_result)
    oid_match = _re.search(r'open_id: `(ou_[A-Za-z0-9]+)`', search_result)
    
    if uid_match:
        return (uid_match.group(1), "user_id")
    if oid_match:
        return (oid_match.group(1), "open_id")
    
    # 2. Try AgentRelationship
    from app.models.org import AgentRelationship
    from sqlalchemy.orm import selectinload
    async with async_session() as db:
        result = await db.execute(
            select(AgentRelationship)
            .where(AgentRelationship.agent_id == agent_id)
            .options(selectinload(AgentRelationship.member))
        )
        for r in result.scalars().all():
            if r.member and r.member.name == member_name:
                if r.member.external_id:
                    return (r.member.external_id, "user_id")
                if r.member.open_id:
                    return (r.member.open_id, "open_id")
                break
    return None


async def _send_file_via_feishu(agent_id, config, file_path: Path, member_name: str, message: str) -> str | None:
    """Send file to a person via Feishu. Returns result string or None."""
    recipient = await _resolve_feishu_recipient(agent_id, config, member_name)
    if not recipient:
        return None
    
    receive_id, id_type = recipient
    from app.services.feishu_service import feishu_service
    try:
        await feishu_service.upload_and_send_file(
            config.app_id, config.app_secret,
            receive_id, file_path,
            receive_id_type=id_type,
            accompany_msg=message,
        )
        return f"File '{file_path.name}' sent to {member_name} via Feishu."
    except Exception as e:
        # If upload fails, try sending a download link as fallback
        import json as _j
        from app.config import get_settings as _gs
        _s = _gs()
        base_url = getattr(_s, 'BASE_URL', '').rstrip('/') or ''
        base_abs = (WORKSPACE_ROOT / str(agent_id)).resolve()
        try:
            _rel = str(file_path.resolve().relative_to(base_abs))
        except ValueError:
            _rel = file_path.name
        parts = []
        if message:
            parts.append(message)
        if base_url:
            dl_url = f"{base_url}/api/agents/{agent_id}/files/download?path={_rel}"
            parts.append(f"{file_path.name}\n{dl_url}")
        parts.append(f"File upload failed ({e}). If you need direct file sending, enable im:resource permission in Feishu.")
        try:
            await feishu_service.send_message(
                config.app_id, config.app_secret,
                receive_id, "text",
                _j.dumps({"text": "\n\n".join(parts)}, ensure_ascii=False),
                receive_id_type=id_type,
            )
            return f"File upload to Feishu failed, sent download link to {member_name} instead."
        except Exception:
            return f"Failed to send file to {member_name} via Feishu: {e}"


async def _send_file_via_slack(agent_id, config, file_path: Path, member_name: str, message: str) -> str | None:
    """Send file to a person via Slack DM. Returns result string or None."""
    import httpx
    bot_token = config.app_secret or ""
    if not bot_token:
        return None
    
    # Resolve Slack user by name
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(
                "https://slack.com/api/users.list",
                headers={"Authorization": f"Bearer {bot_token}"},
                params={"limit": 200},
            )
            data = resp.json()
            if not data.get("ok"):
                return None
            slack_user_id = None
            for u in data.get("members", []):
                profile = u.get("profile", {})
                display = profile.get("display_name", "") or profile.get("real_name", "") or u.get("real_name", "")
                if display == member_name or u.get("name") == member_name:
                    slack_user_id = u.get("id")
                    break
            if not slack_user_id:
                return None
            
            # Open a DM channel
            dm_resp = await client.post(
                "https://slack.com/api/conversations.open",
                headers={"Authorization": f"Bearer {bot_token}", "Content-Type": "application/json"},
                json={"users": slack_user_id},
            )
            dm_data = dm_resp.json()
            if not dm_data.get("ok"):
                return None
            channel_id = dm_data["channel"]["id"]
            
            # Upload file
            upload_url_resp = await client.post(
                "https://slack.com/api/files.getUploadURLExternal",
                headers={"Authorization": f"Bearer {bot_token}"},
                data={"filename": file_path.name, "length": str(file_path.stat().st_size)},
            )
            ud = upload_url_resp.json()
            if not ud.get("ok"):
                return f"Slack file upload failed: {ud.get('error')}"
            await client.post(ud["upload_url"], content=file_path.read_bytes(),
                            headers={"Content-Type": "application/octet-stream"})
            complete = await client.post(
                "https://slack.com/api/files.completeUploadExternal",
                headers={"Authorization": f"Bearer {bot_token}"},
                json={"files": [{"id": ud["file_id"]}], "channel_id": channel_id,
                      "initial_comment": message or ""},
            )
            if not complete.json().get("ok"):
                return f"Slack file upload complete failed: {complete.json().get('error')}"
            return f"File '{file_path.name}' sent to {member_name} via Slack."
    except Exception as e:
        return f"Failed to send file via Slack: {e}"


async def _execute_mcp_tool(tool_name: str, arguments: dict, agent_id=None) -> str:
    """Execute a tool via MCP if it exists in the DB as an MCP tool."""
    try:
        from app.models.tool import Tool, AgentTool
        from app.services.mcp_client import MCPClient

        async with async_session() as db:
            result = await db.execute(select(Tool).where(Tool.name == tool_name, Tool.type == "mcp"))
            tool = result.scalar_one_or_none()

            if not tool:
                logger.warning(f"[MCP] Unknown tool: {tool_name}")
                return f"Unknown tool: {tool_name}"

            # Load per-agent config override
            agent_config = {}
            if tool and agent_id:
                at_r = await db.execute(
                    select(AgentTool).where(
                        AgentTool.agent_id == agent_id,
                        AgentTool.tool_id == tool.id,
                    )
                )
                at = at_r.scalar_one_or_none()
                agent_config = (at.config or {}) if at else {}

        if not tool.mcp_server_url:
            logger.error(f"[MCP] Tool {tool_name} has no server URL configured")
            return f"❌ MCP tool {tool_name} has no server URL configured"

        # Merge global config + agent override
        merged_config = {**(tool.config or {}), **agent_config}
        merged_config = _decrypt_sensitive_fields(merged_config)

        mcp_url = tool.mcp_server_url
        mcp_name = tool.mcp_tool_name or tool_name

        # Detect Smithery-hosted MCP servers (*.run.tools URLs)
        # These need Smithery Connect to route tool calls
        if ".run.tools" in mcp_url and merged_config:
            return await _execute_via_smithery_connect(mcp_url, mcp_name, arguments, merged_config, agent_id=agent_id)

        # Direct MCP call for non-Smithery servers
        # Priority for API key:
        # 1. Per-agent tool config (api_key / atlassian_api_key)
        # 2. Agent's Atlassian channel config (for atlassian_* tools)
        direct_api_key = merged_config.get("api_key") or merged_config.get("atlassian_api_key")
        if not direct_api_key and tool.mcp_server_name == "Atlassian Rovo":
            try:
                from app.api.atlassian import get_atlassian_api_key_for_agent
                direct_api_key = await get_atlassian_api_key_for_agent(agent_id)
            except Exception:
                pass
        client = MCPClient(mcp_url, api_key=direct_api_key)
        return await client.call_tool(mcp_name, arguments)

    except Exception as e:
        logger.exception(f"[MCP] Tool execution error: {tool_name}")
        return f"❌ MCP tool execution error: {str(e)[:200]}"


async def _execute_via_smithery_connect(mcp_url: str, tool_name: str, arguments: dict, config: dict, agent_id=None) -> str:
    """Execute an MCP tool via Smithery Connect API.

    Uses stored namespace/connection or falls back to creating one.
    Smithery Connect returns SSE-format responses that need special parsing.
    """
    import httpx
    import json as json_mod

    # Get Smithery API key centrally (from discover_resources/import_mcp_server AgentTool config)
    from app.services.resource_discovery import _get_smithery_api_key
    api_key = await _get_smithery_api_key(agent_id)
    if not api_key:
        return (
            "❌ Smithery API key not configured.\n\n"
            "请提供你的 Smithery API Key，你可以通过以下步骤获取：\n"
            "1. 注册/登录 https://smithery.ai\n"
            "2. 前往 https://smithery.ai/account/api-keys 创建 API Key\n"
            "3. 将 Key 提供给我，我会帮你配置"
        )

    # Get namespace + connection from tool config, or use defaults
    namespace = config.pop("smithery_namespace", None)
    connection_id = config.pop("smithery_connection_id", None)

    if not namespace or not connection_id:
        # Fallback: try to get from Smithery settings
        try:
            from app.models.tool import Tool
            async with async_session() as db:
                r = await db.execute(select(Tool).where(Tool.name == "discover_resources"))
                disc_tool = r.scalar_one_or_none()
                if disc_tool and disc_tool.config:
                    namespace = namespace or disc_tool.config.get("smithery_namespace")
                    connection_id = connection_id or disc_tool.config.get("smithery_connection_id")
        except Exception:
            pass

    if not namespace or not connection_id:
        return (
            "❌ Smithery Connect namespace/connection not configured. "
            "Please set smithery_namespace and smithery_connection_id in the tool configuration."
        )

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    try:
        async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
            # Call the tool via the existing connection
            tool_resp = await client.post(
                f"https://api.smithery.ai/connect/{namespace}/{connection_id}/mcp",
                json={
                    "jsonrpc": "2.0",
                    "id": 1,
                    "method": "tools/call",
                    "params": {
                        "name": tool_name,
                        "arguments": arguments,
                    },
                },
                headers=headers,
            )

            # Detect auth/connection failures and attempt auto-recovery
            if tool_resp.status_code in (401, 403, 404):
                recovery_result = await _smithery_auto_recover(
                    api_key, mcp_url, namespace, connection_id, agent_id
                )
                if recovery_result:
                    return recovery_result
                # If recovery returned None, fall through to normal parsing

            # Smithery Connect returns SSE format: "event: message\ndata: {...}\n"
            raw = tool_resp.text
            data = None

            # Parse SSE response
            for line in raw.split("\n"):
                line = line.strip()
                if line.startswith("data: "):
                    try:
                        data = json_mod.loads(line[6:])
                        break
                    except json_mod.JSONDecodeError:
                        pass

            # Fallback: try parsing as plain JSON
            if data is None:
                try:
                    data = json_mod.loads(raw)
                except json_mod.JSONDecodeError:
                    return f"❌ Unexpected response from Smithery: {raw[:300]}"

            if "error" in data:
                err = data["error"]
                msg = err.get("message", str(err)) if isinstance(err, dict) else str(err)
                # Check if error indicates auth/connection issue
                auth_keywords = ["auth", "unauthorized", "forbidden", "expired", "not found", "connection"]
                if any(kw in msg.lower() for kw in auth_keywords):
                    recovery_result = await _smithery_auto_recover(
                        api_key, mcp_url, namespace, connection_id, agent_id
                    )
                    if recovery_result:
                        return recovery_result
                return f"❌ MCP tool error: {msg[:300]}"

            result = data.get("result", {})
            if isinstance(result, str):
                return result

            content_blocks = result.get("content", []) if isinstance(result, dict) else []
            texts = []
            for block in content_blocks:
                if isinstance(block, str):
                    texts.append(block)
                elif isinstance(block, dict):
                    if block.get("type") == "text":
                        texts.append(block.get("text", ""))
                    elif block.get("type") == "image":
                        texts.append(f"[Image: {block.get('mimeType', 'image')}]")
                    else:
                        texts.append(str(block))
                else:
                    texts.append(str(block))

            return "\n".join(texts) if texts else str(result)

    except Exception as e:
        return f"❌ Smithery Connect error: {str(e)[:200]}"


async def _smithery_auto_recover(api_key: str, mcp_url: str, namespace: str, connection_id: str, agent_id=None) -> str | None:
    """Attempt to auto-recover a failed Smithery connection.

    Re-creates the Smithery Connect connection. If OAuth is needed,
    returns the auth URL for the user. Returns None if recovery fails silently.
    """
    try:
        from app.services.resource_discovery import _ensure_smithery_connection
        display_name = connection_id.replace("-", " ").title() if connection_id else "MCP Server"

        conn_result = await _ensure_smithery_connection(api_key, mcp_url, display_name)
        if "error" in conn_result:
            return (
                f"❌ MCP tool connection expired and auto-recovery failed: {conn_result['error']}\n\n"
                f"💡 Please re-authorize by telling me: `import_mcp_server(server_id=\"...\", reauthorize=true)`"
            )

        # Update stored config with new connection info
        new_config = {
            "smithery_namespace": conn_result["namespace"],
            "smithery_connection_id": conn_result["connection_id"],
        }
        if agent_id:
            try:
                from app.models.tool import Tool, AgentTool
                async with async_session() as db:
                    # Update all MCP tools for this server URL
                    r = await db.execute(
                        select(Tool).where(Tool.mcp_server_url == mcp_url, Tool.type == "mcp")
                    )
                    for tool in r.scalars().all():
                        at_r = await db.execute(
                            select(AgentTool).where(
                                AgentTool.agent_id == agent_id,
                                AgentTool.tool_id == tool.id,
                            )
                        )
                        at = at_r.scalar_one_or_none()
                        if at:
                            at.config = {**(at.config or {}), **new_config}
                    await db.commit()
            except Exception:
                pass  # Non-critical — connection may still work

        if conn_result.get("auth_url"):
            return (
                f"🔐 MCP tool connection expired. Re-authorization needed.\n\n"
                f"Please visit the following URL to re-authorize:\n"
                f"{conn_result['auth_url']}\n\n"
                f"After completing authorization, the tools will work again automatically."
            )

        # Connection re-created without OAuth — should work now
        return None  # Signal caller to retry (but we don't retry here to avoid loops)

    except Exception as e:
        return f"❌ Auto-recovery failed: {str(e)[:200]}"


def _list_files(ws: Path, rel_path: str, tenant_id: str | None = None) -> str:
    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        # Remap: enterprise_info/... → enterprise_info_{tenant_id}/...
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        target = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(target).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        target = (ws / rel_path) if rel_path else ws
        target = target.resolve()
        if not str(target).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not target.exists():
        return f"Directory not found: {rel_path or '/'}"

    items = []
    # If listing root, also show enterprise_info entry
    if not rel_path:
        if tenant_id:
            enterprise_dir = WORKSPACE_ROOT / f"enterprise_info_{tenant_id}"
        else:
            enterprise_dir = WORKSPACE_ROOT / "enterprise_info"
        if enterprise_dir.exists():
            items.append("  📁 enterprise_info/ (shared company info)")

    dir_count = 0
    file_count = 0
    for p in sorted(target.iterdir()):
        if p.name.startswith("."):
            continue
        if p.is_dir():
            dir_count += 1
            child_count = len([c for c in p.iterdir() if not c.name.startswith(".")])
            items.append(f"  📁 {p.name}/ ({child_count} items)")
        elif p.is_file():
            file_count += 1
            size_bytes = p.stat().st_size
            if size_bytes < 1024:
                size_str = f"{size_bytes}B"
            else:
                size_str = f"{size_bytes/1024:.1f}KB"
            items.append(f"  📄 {p.name} ({size_str})")

    if not items:
        return f"📂 {rel_path or 'root'}: Empty directory (0 files, 0 folders)"

    header = f"📂 {rel_path or 'root'}: {dir_count} folder(s), {file_count} file(s)\n"
    return header + "\n".join(items)


def _read_file(ws: Path, rel_path: str, tenant_id: str | None = None) -> str:
    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        file_path = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        content = file_path.read_text(encoding="utf-8", errors="replace")
        if len(content) > 6000:
            content = content[:6000] + f"\n\n...[truncated, {len(content)} chars total]"
        return content
    except Exception as e:
        return f"Read failed: {e}"


async def _read_document(ws: Path, rel_path: str, max_chars: int = 8000, tenant_id: str | None = None) -> str:
    """Read content from office documents (PDF, DOCX, XLSX, PPTX)."""
    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        file_path = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not file_path.exists():
        return f"File not found: {rel_path}"

    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            import pdfplumber
            text_parts = []
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages[:50]):  # Limit to 50 pages
                    page_text = page.extract_text() or ""
                    if page_text:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
            content = "\n\n".join(text_parts) if text_parts else "(PDF is empty or text extraction failed)"

        elif ext == ".docx":
            from docx import Document
            from docx.oxml.ns import qn
            doc = Document(str(file_path))
            lines: list[str] = []

            def _extract_para_text(para) -> str:
                return para.text.strip()

            def _extract_table(table) -> str:
                """Flatten a table into readable text."""
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    # Remove duplicate adjacent cells (merged cells repeat)
                    deduped = [cells[0]] + [c for i, c in enumerate(cells[1:]) if c != cells[i]]
                    row_str = " | ".join(c for c in deduped if c)
                    if row_str:
                        rows.append(row_str)
                return "\n".join(rows)

            # 1. Main paragraphs
            for para in doc.paragraphs:
                t = _extract_para_text(para)
                if t:
                    lines.append(t)

            # 2. Tables in main body
            for table in doc.tables:
                t = _extract_table(table)
                if t:
                    lines.append(t)

            # 3. Text boxes / drawing shapes (wmf/shapes in body XML)
            for shape in doc.element.body.iter(qn("w:txbxContent")):
                for child in shape.iter(qn("w:t")):
                    if child.text and child.text.strip():
                        lines.append(child.text.strip())

            # 4. Headers and footers
            for section in doc.sections:
                for hf in [section.header, section.footer]:
                    if hf and hf.is_linked_to_previous is False:
                        for para in hf.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)

            content = "\n".join(lines) if lines else "(Document is empty or uses unsupported formatting)"

        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            sheets = []
            for ws_name in wb.sheetnames[:10]:  # Limit to 10 sheets
                sheet = wb[ws_name]
                rows = []
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheets.append(f"=== Sheet: {ws_name} ===\n" + "\n".join(rows))
            wb.close()
            content = "\n\n".join(sheets) if sheets else "(Excel is empty)"

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides[:50]):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                if texts:
                    slides.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
            content = "\n\n".join(slides) if slides else "(PPT is empty)"

        elif ext in (".txt", ".md", ".json", ".csv", ".log"):
            content = file_path.read_text(encoding="utf-8", errors="replace")

        else:
            return f"Unsupported file format: {ext}. Supported: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV"

        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n...[truncated, {len(content)} chars total]"
        return content

    except ImportError as e:
        return f"Missing dependency: {e}. Install: pip install pdfplumber python-docx openpyxl python-pptx"
    except Exception as e:
        return f"Document read failed: {str(e)[:200]}"


def _write_file(ws: Path, rel_path: str, content: str, tenant_id: str | None = None) -> str:
    # Protect tasks.json from direct writes
    if rel_path.strip("/") == "tasks.json":
        return "tasks.json is read-only. Use manage_tasks tool to manage tasks."

    # Handle enterprise_info/ as shared directory (tenant-scoped)
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        if not sub:
            return "Write failed: please provide a file path under enterprise_info/, e.g. enterprise_info/knowledge_base/report.md"
        file_path = (enterprise_root / sub).resolve()
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    try:
        file_path.parent.mkdir(parents=True, exist_ok=True)
        file_path.write_text(content, encoding="utf-8")
        return f"✅ Written to {rel_path} ({len(content)} chars)"
    except Exception as e:
        return f"Write failed: {e}"


def _delete_file(ws: Path, rel_path: str) -> str:
    protected = {"tasks.json", "soul.md"}
    if rel_path.strip("/") in protected:
        return f"{rel_path} cannot be deleted (protected)"

    file_path = (ws / rel_path).resolve()
    if not str(file_path).startswith(str(ws.resolve())):
        return "Access denied for this path"
    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        if file_path.is_dir():
            import shutil
            shutil.rmtree(file_path)
            return f"✅ Deleted directory {rel_path}"
        else:
            file_path.unlink()
            return f"✅ Deleted {rel_path}"
    except Exception as e:
        return f"Delete failed: {e}"


async def _manage_tasks(
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    ws: Path,
    args: dict,
) -> str:
    """Create / update / delete tasks in DB and sync to workspace."""
    from app.models.task import TaskLog
    from datetime import datetime, timezone

    action = args["action"]
    title = args["title"]

    async with async_session() as db:
        if action == "create":
            task_type = args.get("task_type", "todo")
            task = Task(
                agent_id=agent_id,
                title=title,
                description=args.get("description"),
                type=task_type,
                priority=args.get("priority", "medium"),
                created_by=user_id,
                status="pending",
                supervision_target_name=args.get("supervision_target_name"),
                supervision_channel=args.get("supervision_channel", "feishu"),
                remind_schedule=args.get("remind_schedule"),
            )
            db.add(task)
            await db.commit()
            await db.refresh(task)

            if task_type == "todo":
                # Trigger auto-execution for todo tasks
                import asyncio
                from app.services.task_executor import execute_task
                asyncio.create_task(execute_task(task.id, agent_id))
                await _sync_tasks_to_file(agent_id, ws)
                return f"✅ Task created: {title} — auto-execution started"
            else:
                # Supervision task — reminder engine will pick it up
                target = args.get('supervision_target_name', 'someone')
                schedule = args.get('remind_schedule', 'not set')
                await _sync_tasks_to_file(agent_id, ws)
                return f"✅ Supervision task created: '{title}' — will remind {target} on schedule ({schedule})"

        elif action == "update_status":
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id, Task.title.ilike(f"%{title}%"))
            )
            task = result.scalars().first()
            if not task:
                return f"No task found matching '{title}'"
            old = task.status
            task.status = args["status"]
            if args["status"] == "done":
                task.completed_at = datetime.now(timezone.utc)
            await db.commit()
            await _sync_tasks_to_file(agent_id, ws)
            return f"✅ Updated '{task.title}' from {old} to {args['status']}"

        elif action == "delete":
            from sqlalchemy import delete as sa_delete
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id, Task.title.ilike(f"%{title}%"))
            )
            task = result.scalars().first()
            if not task:
                return f"No task found matching '{title}'"
            task_title = task.title
            await db.execute(sa_delete(TaskLog).where(TaskLog.task_id == task.id))
            await db.delete(task)
            await db.commit()
            await _sync_tasks_to_file(agent_id, ws)
            return f"✅ Task deleted: {task_title}"

        return f"Unknown action: {action}"


async def _send_feishu_message(agent_id: uuid.UUID, args: dict) -> str:
    """Send a Feishu message to a person in the agent's relationship list."""
    member_name = (args.get("member_name") or "").strip()
    direct_open_id = (args.get("open_id") or "").strip()
    direct_user_id = (args.get("user_id") or "").strip()
    message_text = (args.get("message") or "").strip()

    if not message_text:
        return "❌ Please provide message content"
    if not member_name and not direct_open_id and not direct_user_id:
        return "❌ Please provide member_name, user_id, or open_id"

    try:
        from app.services.feishu_service import feishu_service
        from sqlalchemy.orm import selectinload

        async with async_session() as db:

            # ── Shortcut: if caller provided user_id or open_id directly ──
            config_result = await db.execute(
                select(ChannelConfig).where(ChannelConfig.agent_id == agent_id, ChannelConfig.channel_type == "feishu")
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return "❌ This agent has no Feishu channel configured"
            if (direct_user_id or direct_open_id) and not member_name:
                import json as _j
                # Prefer user_id over open_id
                if direct_user_id:
                    resp = await feishu_service.send_message(
                        config.app_id, config.app_secret,
                        receive_id=direct_user_id, msg_type="text",
                        content=_j.dumps({"text": message_text}, ensure_ascii=False),
                        receive_id_type="user_id",
                    )
                    if resp.get("code") == 0:
                        # Save to history session
                        await _save_outgoing_to_feishu_session(direct_user_id or direct_open_id)
                        return f"✅ 消息已发送（user_id: {direct_user_id}）"
                    # Fallback to open_id if user_id fails
                    logger.info(f"❌ 发送失败：{resp.get('msg')} (code {resp.get('code')})")
                    if direct_open_id:
                        resp = await feishu_service.send_message(
                            config.app_id, config.app_secret,
                            receive_id=direct_open_id, msg_type="text",
                            content=_j.dumps({"text": message_text}, ensure_ascii=False),
                            receive_id_type="open_id",
                        )
                        if resp.get("code") == 0:
                            await _save_outgoing_to_feishu_session(direct_open_id)
                            return f"✅ 消息已发送（open_id: {direct_open_id}）"
                    return f"❌ 发送失败：{resp.get('msg')} (code {resp.get('code')})"
                else:
                    resp = await feishu_service.send_message(
                        config.app_id, config.app_secret,
                        receive_id=direct_open_id, msg_type="text",
                        content=_j.dumps({"text": message_text}, ensure_ascii=False),
                        receive_id_type="open_id",
                    )
                    if resp.get("code") == 0:
                        await _save_outgoing_to_feishu_session(direct_open_id)
                        return f"✅ 消息已发送（open_id: {direct_open_id}）"
                    logger.info(f"❌ 发送失败：{resp.get('msg')} (code {resp.get('code')})")
                    return f"❌ 发送失败：{resp.get('msg')} (code {resp.get('code')})"

            # Find the relationship member by name
            result = await db.execute(
                select(AgentRelationship)
                .where(AgentRelationship.agent_id == agent_id)
                .options(selectinload(AgentRelationship.member))
            )
            rels = result.scalars().all()

            target_member = None
            for r in rels:
                if r.member and r.member.name == member_name:
                    target_member = r.member
                    break

            if not target_member:
                logger.info(f"❌ {member_name} has no Feishu user_id in relationship")   
                return f"❌ {member_name} 不是我的关系"
                
            logger.info(f"target_member={target_member.external_id}, {target_member.open_id}, {target_member.email}, {target_member.phone}")
            if not target_member.external_id and not target_member.open_id and not target_member.email and not target_member.phone:
                logger.error(f"❌ {member_name} has no linked Feishu account (no user_id, open_id, email, or phone)")
                return f"❌ {member_name} has no linked Feishu account (no user_id, open_id, email, or phone)"

            content = json.dumps({"text": message_text}, ensure_ascii=False)

            async def _try_send(app_id: str, app_secret: str, receive_id: str, id_type: str = "open_id") -> dict:
                return await feishu_service.send_message(
                    app_id, app_secret,
                    receive_id=receive_id, msg_type="text",
                    content=content, receive_id_type=id_type,
                )

            async def _save_outgoing_to_feishu_session(open_id: str):
                """Save the outgoing message to the Feishu P2P chat session."""
                try:
                    from datetime import datetime as _dt, timezone as _tz


                    agent_r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
                    agent_obj = agent_r.scalar_one_or_none()
                    creator_id = agent_obj.creator_id if agent_obj else agent_id

                    # Get or create platform user from OrgMember (unified logic)
                    platform_user = await get_platform_user_by_org_member(
                        db=db,
                        org_member=target_member,
                        agent_tenant_id=agent_obj.tenant_id if agent_obj else None,
                    )
                    user_id = platform_user.id

                    ext_conv_id = f"feishu_p2p_{open_id}"
                    sess = await find_or_create_channel_session(
                        db=db,
                        agent_id=agent_id,
                        user_id=user_id,
                        external_conv_id=ext_conv_id,
                        source_channel="feishu",
                        first_message_title=f"[Agent → {member_name or open_id}]",
                    )
                    db.add(ChatMessage(
                        agent_id=agent_id,
                        user_id=user_id,
                        role="assistant",
                        content=message_text,
                        conversation_id=str(sess.id),
                    ))
                    sess.last_message_at = _dt.now(_tz.utc)
                    await db.commit()
                    logger.info(f"[Feishu] Saved outgoing message to session {sess.id} (ID: {open_id})")
                except Exception as e:
                    logger.error(f"[Feishu] Failed to save outgoing message to history: {e}")

            # Step 1: Try using feishu_user_id (tenant-stable, works across apps)
            if target_member.external_id:
                resp = await _try_send(config.app_id, config.app_secret, target_member.external_id, "user_id")
                if resp.get("code") == 0:
                    await _save_outgoing_to_feishu_session(target_member.external_id or target_member.open_id)
                    return f"✅ Successfully sent message to {member_name}"
                logger.info(f"❌ Failed to send message to {target_member.external_id} via Feishu (user_id): {resp}")
                
                # Fallback to open_id if user_id fails (e.g., due to missing employee_id:readonly permission)
                if target_member.open_id:
                    resp_open = await _try_send(config.app_id, config.app_secret, target_member.open_id, "open_id")
                    if resp_open.get("code") == 0:
                        await _save_outgoing_to_feishu_session(target_member.open_id)
                        return f"✅ Successfully sent message to {member_name}"
                    logger.info(f"❌ Failed to send message to {target_member.open_id} via Feishu (open_id): {resp_open}")
                    return f"发送失败 (user_id: {resp.get('code')}, open_id: {resp_open.get('code')}): {resp_open.get('msg')}"
                return f"发送失败 {resp}"
            
            # Step 2: If no external_id, try open_id directly
            elif target_member.open_id:
                resp = await _try_send(config.app_id, config.app_secret, target_member.open_id, "open_id")
                if resp.get("code") == 0:
                    await _save_outgoing_to_feishu_session(target_member.open_id)
                    return f"✅ Successfully sent message to {member_name}"
                logger.info(f"❌ Failed to send message to {target_member.open_id} via Feishu (open_id): {resp}")
                return f"发送失败 {resp}"
    except Exception as e:
        return f"❌ Message send error: {str(e)[:200]}"


async def _send_channel_message(agent_id: uuid.UUID, args: dict) -> str:
    """Send message via the recipient's configured channel (Feishu/DingTalk/WeCom).

    1. Find target user from relationships (AgentRelationship -> OrgMember)
    2. Determine user's provider type (via OrgMember.provider_id -> IdentityProvider)
    3. Find corresponding channel config (ChannelConfig)
    4. Send via the appropriate channel
    """
    from sqlalchemy import select
    from sqlalchemy.orm import selectinload
    from app.models.org import AgentRelationship, OrgMember
    from app.models.identity import IdentityProvider

    member_name = (args.get("member_name") or "").strip()
    message_text = (args.get("message") or "").strip()
    target_channel = (args.get("channel") or "").strip().lower()

    if not member_name:
        return "❌ Please provide member_name"
    if not message_text:
        return "❌ Please provide message content"

    try:
        async with async_session() as db:
            # 1. Find target member from relationships with provider info (only active members)
            result = await db.execute(
                select(AgentRelationship, OrgMember, IdentityProvider)
                .join(OrgMember, AgentRelationship.member_id == OrgMember.id)
                .outerjoin(IdentityProvider, OrgMember.provider_id == IdentityProvider.id)
                .where(AgentRelationship.agent_id == agent_id, OrgMember.name == member_name, OrgMember.status == "active")
                .options(selectinload(AgentRelationship.member))
            )
            rows = result.all()

            if not rows:
                return f"❌ {member_name} is not in your relationship network"

            target_member = None
            provider_type = None

            # Handle multiple matches across different providers
            if target_channel:
                for rel, member, provider in rows:
                    if provider and provider.provider_type == target_channel:
                        target_member = member
                        provider_type = target_channel
                        break
                if not target_member:
                    available = [p.provider_type for _, _, p in rows if p]
                    return f"❌ {member_name} not found in {target_channel} channel. Available channels: {', '.join(available)}"
            else:
                if len(rows) > 1:
                    available = [p.provider_type for _, _, p in rows if p]
                    logger.warning(f"[ChannelMessage] Ambiguous member '{member_name}' found in multiple channels: {available}")
                    # Pick the first one as before, but mention others if possible
                
                rel, member, provider = rows[0]
                target_member = member
                provider_type = provider.provider_type if provider else None

            # 2. Determine channel based on provider type
            if not provider_type:
                # Fallback: check which channel configs exist and has user info
                if target_member.external_id or target_member.open_id:
                    # Try Feishu as default
                    provider_type = "feishu"
                else:
                    return f"❌ {member_name} has no linked channel (no provider info)"

            logger.info(f"[ChannelMessage] Sending to {member_name} via {provider_type}")

            # 3. Route to appropriate channel
            if provider_type == "feishu":
                return await _send_feishu_message(agent_id, {"member_name": member_name, "message": message_text})
            elif provider_type == "dingtalk":
                return await _send_dingtalk_message(agent_id, member_name, message_text, target_member)
            elif provider_type == "wecom":
                return await _send_wecom_message(agent_id, member_name, message_text, target_member)
            else:
                return f"❌ Unsupported channel type: {provider_type}"

    except Exception as e:
        logger.exception("[ChannelMessage] Error")
        return f"❌ Channel message error: {str(e)[:200]}"


async def _send_dingtalk_message(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> str:
    """Send message via DingTalk channel using Open API."""
    from app.services.dingtalk_service import send_dingtalk_message


    try:
        async with async_session() as db:
            # 1. Get DingTalk channel config
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "dingtalk",
                    ChannelConfig.is_configured == True,
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return "❌ This agent has no DingTalk channel configured"

            # 2. Get recipient's user_id (external_id)
            user_id = target_member.external_id
            if not user_id:
                # Try to use unionid or openid as fallback
                user_id = target_member.unionid or target_member.open_id
                if not user_id:
                    return f"❌ {member_name} has no DingTalk user_id"

            logger.info(f"[DingTalk] Sending to user_id: {user_id}")

            # Get agent_id from extra_config (required for DingTalk API)
            agent_id_dingtalk = config.extra_config.get("agent_id") if config.extra_config else None

            # 3. Send message via DingTalk service
            result = await send_dingtalk_message(
                app_id=config.app_id,
                app_secret=config.app_secret,
                user_id=user_id,
                message=message_text,
                agent_id=agent_id_dingtalk,
            )

            if result.get("errcode") == 0:
                try:
                    # Get agent tenant context
                    agent_r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
                    agent_obj = agent_r.scalar_one_or_none()


                    # Get or create platform user from OrgMember (unified logic)
                    platform_user = await get_platform_user_by_org_member(
                        db=db,
                        org_member=target_member,
                        agent_tenant_id=agent_obj.tenant_id if agent_obj else None,
                    )


                    conv_id = f"dingtalk_p2p_{user_id}"
                    # 2. Get/Create session
                    sess = await find_or_create_channel_session(
                        db=db,
                        agent_id=agent_id,
                        user_id=platform_user.id,
                        external_conv_id=conv_id,
                        source_channel="dingtalk",
                        first_message_title=message_text[:30],
                    )
                    # 3. Save assistant message
                    db.add(ChatMessage(
                        agent_id=agent_id,
                        user_id=platform_user.id,
                        role="assistant",
                        content=message_text,
                        conversation_id=str(sess.id),
                    ))
                    sess.last_message_at = datetime.now(timezone.utc)
                    await db.commit()
                    logger.info(f"[DingTalk] Proactive message saved to session {sess.id}")
                except Exception as ex:
                    logger.error(f"[DingTalk] Failed to save proactive message to session: {ex}")

                return f"✅ Message sent to {member_name} via DingTalk"
            else:
                errmsg = result.get("errmsg", "Unknown error")
                logger.error(f"[DingTalk] Send failed: {result}")
                return f"❌ DingTalk send failed: {errmsg}"

    except Exception as e:
        logger.exception("[DingTalk] Error")
        return f"❌ DingTalk message error: {str(e)[:200]}"


async def _send_wecom_message(
    agent_id: uuid.UUID,
    member_name: str,
    message_text: str,
    target_member: "OrgMember",
) -> str:
    """Send message via WeCom channel using Open API."""
    from app.services.wecom_service import send_wecom_message


    try:
        async with async_session() as db:
            # 1. Get WeCom channel config
            config_result = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "wecom",
                    ChannelConfig.is_configured == True,
                )
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return "❌ This agent has no WeCom channel configured"

            # 2. Get recipient's user_id
            user_id = target_member.external_id
            if not user_id:
                user_id = target_member.open_id
                if not user_id:
                    return f"❌ {member_name} has no WeCom user_id"

            logger.info(f"[WeCom] Sending to user_id: {user_id}")

            # 3. Send message via WeCom service
            result = await send_wecom_message(
                config.app_id,
                config.app_secret,
                user_id,
                message_text,
            )

            if result.get("errcode") == 0:
                # Save proactive message to session so it appears in UI
                try:

                    # Get agent tenant context
                    agent_r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
                    agent = agent_r.scalar_one_or_none()


                    # Get or create platform user from OrgMember (unified logic)
                    platform_user = await get_platform_user_by_org_member(
                        db=db,
                        org_member=target_member,
                        agent_tenant_id=agent.tenant_id if agent else None,
                    )

                    conv_id = f"wecom_p2p_{user_id}"
                    sess = await find_or_create_channel_session(
                        db=db,
                        agent_id=agent_id,
                        user_id=platform_user.id,
                        external_conv_id=conv_id,
                        source_channel="wecom",
                        first_message_title=message_text[:30],
                    )
                    db.add(ChatMessage(
                        agent_id=agent_id,
                        user_id=platform_user.id,
                        role="assistant",
                        content=message_text,
                        conversation_id=str(sess.id),
                    ))
                    sess.last_message_at = datetime.now(timezone.utc)
                    await db.commit()
                    logger.info(f"[WeCom] Proactive message saved to session {sess.id}")
                except Exception as ex:
                    logger.error(f"[WeCom] Failed to save proactive message to session: {ex}")

                return f"✅ Message sent to {member_name} via WeCom"
            else:
                errmsg = result.get("errmsg", "Unknown error")
                logger.error(f"[WeCom] Send failed: {result}")
                return f"❌ WeCom send failed: {errmsg}"

    except Exception as e:
        logger.exception("[WeCom] Error")
        return f"❌ WeCom message error: {str(e)[:200]}"


async def _send_web_message(agent_id: uuid.UUID, args: dict) -> str:
    """Send a proactive message to a web platform user."""
    username = args.get("username", "").strip()
    message_text = args.get("message", "").strip()

    if not username or not message_text:
        return "❌ Please provide recipient username and message content"

    try:
        from datetime import datetime as _dt, timezone as _tz


        async with async_session() as db:
            # 0. Get agent's tenant_id for scoping
            agent_res = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = agent_res.scalar_one_or_none()
            if not agent:
                return "❌ Agent not found"

            # 1. Look up target user by username or display_name within tenant

            query = select(UserModel).where(
                or_(
                    UserModel.username == username,
                    UserModel.display_name == username,
                )
            )
            if agent.tenant_id:
                query = query.where(UserModel.tenant_id == agent.tenant_id)

            u_result = await db.execute(query)
            target_user = u_result.scalar_one_or_none()
            if not target_user:
                # List available users for the agent to pick from (within the same tenant)
                list_query = select(UserModel.username, UserModel.display_name).limit(20)
                if agent.tenant_id:
                    list_query = list_query.where(UserModel.tenant_id == agent.tenant_id)
                
                all_r = await db.execute(list_query)
                names = [f"{r.display_name or r.username}" for r in all_r.all()]
                return f"❌ No user named '{username}' found in your organization. Available users: {', '.join(names) if names else 'none'}"

            # Find or create a web session between the agent and this user
            sess_r = await db.execute(
                select(ChatSession).where(
                    ChatSession.agent_id == agent_id,
                    ChatSession.user_id == target_user.id,
                    ChatSession.source_channel == "web",
                ).order_by(ChatSession.created_at.desc()).limit(1)
            )
            session = sess_r.scalar_one_or_none()

            if not session:
                # Create a new session for this user
                session = ChatSession(
                    agent_id=agent_id,
                    user_id=target_user.id,
                    title=f"[Agent Message] {_dt.now(_tz.utc).strftime('%m-%d %H:%M')}",
                    source_channel="web",
                    created_at=_dt.now(_tz.utc),
                )
                db.add(session)
                await db.flush()

            # Save the message
            db.add(ChatMessage(
                agent_id=agent_id,
                user_id=target_user.id,
                role="assistant",
                content=message_text,
                conversation_id=str(session.id),
            ))
            session.last_message_at = _dt.now(_tz.utc)
            await db.commit()

            # Push via WebSocket if user has an active connection
            try:
                from app.api.websocket import manager as ws_manager
                agent_id_str = str(agent_id)
                if agent_id_str in ws_manager.active_connections:
                    for ws, sid in list(ws_manager.active_connections[agent_id_str]):
                        try:
                            await ws.send_json({
                                "type": "trigger_notification",
                                "content": message_text,
                                "triggers": ["web_message"],
                            })
                        except Exception:
                            pass
            except Exception:
                pass

            display = target_user.display_name or target_user.username
            return f"✅ Message sent to {display} on web platform. It has been saved to their chat history."

    except Exception as e:
        return f"❌ Web message send error: {str(e)[:200]}"


async def _send_file_to_agent(from_agent_id: uuid.UUID, ws: Path, args: dict) -> str:
    """Send a workspace file to another digital employee (agent)."""
    agent_name = (args.get("agent_name") or "").strip()
    rel_path = (args.get("file_path") or "").strip()
    delivery_note = (args.get("message") or "").strip()

    if not agent_name or not rel_path:
        return "❌ Please provide both agent_name and file_path"

    # Resolve source file path inside sender workspace
    source_file_path = (ws / rel_path).resolve()
    ws_resolved = ws.resolve()
    sender_root = (WORKSPACE_ROOT / str(from_agent_id)).resolve()
    if not str(source_file_path).startswith(str(ws_resolved)):
        source_file_path = (sender_root / rel_path).resolve()
    if not str(source_file_path).startswith(str(sender_root)):
        return "❌ Access denied: source path is outside your workspace"

    if not source_file_path.exists():
        return f"❌ Source file not found: {rel_path}"
    if not source_file_path.is_file():
        return f"❌ Source path is not a file: {rel_path}"

    # File size limit (50 MB)
    MAX_FILE_SIZE = 50 * 1024 * 1024
    file_size = source_file_path.stat().st_size
    if file_size > MAX_FILE_SIZE:
        size_mb = file_size / (1024 * 1024)
        return f"❌ File too large ({size_mb:.1f} MB). Maximum allowed is 50 MB."

    try:
        from app.services.activity_logger import log_activity
        import shutil

        async with async_session() as db:
            src_result = await db.execute(select(AgentModel).where(AgentModel.id == from_agent_id))
            source_agent = src_result.scalar_one_or_none()
            source_name = source_agent.name if source_agent else "Unknown agent"
            source_tenant_id = source_agent.tenant_id if source_agent else None

            # Build base filter: same tenant + not self
            base_filter = [AgentModel.id != from_agent_id]
            if source_tenant_id:
                base_filter.append(AgentModel.tenant_id == source_tenant_id)

            # Try exact name match first, then fuzzy
            target_agent = None
            exact_result = await db.execute(
                select(AgentModel).where(AgentModel.name == agent_name, *base_filter)
            )
            target_agent = exact_result.scalars().first()
            if not target_agent:
                # Sanitize SQL wildcards in user input
                safe_name = agent_name.replace("%", "").replace("_", r"\_")
                fuzzy_result = await db.execute(
                    select(AgentModel).where(AgentModel.name.ilike(f"%{safe_name}%"), *base_filter)
                )
                target_agent = fuzzy_result.scalars().first()

            if not target_agent:
                # Only show agents from relationships, not all agents
                from app.models.org import AgentAgentRelationship
                rel_r = await db.execute(
                    select(AgentModel.name).join(
                        AgentAgentRelationship,
                        (AgentAgentRelationship.target_agent_id == AgentModel.id) & (AgentAgentRelationship.agent_id == from_agent_id)
                    )
                )
                rel_names = [n for (n,) in rel_r.all()]
                return f"❌ No agent found matching '{agent_name}'. Your connected colleagues: {', '.join(rel_names) if rel_names else 'none — ask your administrator to set up relationships'}"

            if target_agent.is_expired or (target_agent.expires_at and datetime.now(timezone.utc) >= target_agent.expires_at):
                return f"⚠️ {target_agent.name} is currently unavailable — their service period has ended. Please contact the platform administrator."

            # Enforce relationship: only allow file transfer with agents in relationships
            rel_check = await db.execute(
                select(AgentAgentRelationship.id).where(
                    ((AgentAgentRelationship.agent_id == from_agent_id) & (AgentAgentRelationship.target_agent_id == target_agent.id))
                    | ((AgentAgentRelationship.agent_id == target_agent.id) & (AgentAgentRelationship.target_agent_id == from_agent_id))
                ).limit(1)
            )
            if not rel_check.scalar_one_or_none():
                return f"❌ You do not have a relationship with {target_agent.name}. Only agents in your relationship list can receive files. Ask your administrator to add a relationship if needed."

            target_tenant_id = str(target_agent.tenant_id) if target_agent.tenant_id else None
            target_name = target_agent.name
            target_id = target_agent.id

        target_ws = await ensure_workspace(target_id, tenant_id=target_tenant_id)
        inbox_dir = (target_ws / "workspace" / "inbox").resolve()
        files_dir = (inbox_dir / "files").resolve()
        target_ws_resolved = target_ws.resolve()
        if not str(inbox_dir).startswith(str(target_ws_resolved)) or not str(files_dir).startswith(str(target_ws_resolved)):
            return "❌ Access denied for target agent inbox path"

        inbox_dir.mkdir(parents=True, exist_ok=True)
        files_dir.mkdir(parents=True, exist_ok=True)

        ts = datetime.now(timezone.utc)
        stamp = ts.strftime("%Y%m%d_%H%M%S_%f")
        delivered_name = source_file_path.name
        delivered_path = files_dir / delivered_name
        while delivered_path.exists():
            delivered_name = f"{stamp}_{source_file_path.name}"
            delivered_path = files_dir / delivered_name

        shutil.copy2(source_file_path, delivered_path)

        sender_short = str(from_agent_id)[:8]
        note_path = inbox_dir / f"{stamp}_{sender_short}_file_delivery.md"
        target_rel_path = f"workspace/inbox/files/{delivered_name}"
        note_lines = [
            f"# File delivery from {source_name}",
            "",
            f"- Time (UTC): {ts.isoformat()}",
            f"- Sender: {source_name}",
            f"- Source path: {rel_path}",
            f"- Delivered file: {target_rel_path}",
            "",
        ]
        if delivery_note:
            note_lines.append("## Note")
            note_lines.append(delivery_note)
            note_lines.append("")
        note_lines.append("## Action")
        note_lines.append(f"- Read the file via `read_file(path=\"{target_rel_path}\")`")
        note_path.write_text("\n".join(note_lines), encoding="utf-8")

        from app.models.audit import AuditLog
        async with async_session() as db:
            db.add(AuditLog(
                agent_id=from_agent_id,
                action="collaboration:file_send",
                details={
                    "to_agent": str(target_id),
                    "to_agent_name": target_name,
                    "source_file": rel_path,
                    "delivered_file": target_rel_path,
                },
            ))
            db.add(AuditLog(
                agent_id=target_id,
                action="collaboration:file_receive",
                details={
                    "from_agent": str(from_agent_id),
                    "from_agent_name": source_name,
                    "source_file": rel_path,
                    "delivered_file": target_rel_path,
                },
            ))
            await db.commit()

        await log_activity(
            from_agent_id,
            "agent_file_sent",
            f"Sent file to {target_name}",
            detail={"target_agent": target_name, "source_file": rel_path, "delivered_file": target_rel_path},
        )
        await log_activity(
            target_id,
            "agent_file_received",
            f"Received file from {source_name}",
            detail={"source_agent": source_name, "source_file": rel_path, "delivered_file": target_rel_path},
        )

        return (
            f"✅ File sent to {target_name}.\n"
            f"- Delivered to: {target_rel_path}\n"
            f"- Inbox note: workspace/inbox/{note_path.name}"
        )
    except Exception as e:
        return f"❌ Agent file send error: {str(e)[:200]}"


async def _send_message_to_agent(from_agent_id: uuid.UUID, args: dict) -> str:
    """Send a message to another digital employee. Uses a single request-response pattern:
    the source agent sends a message, the target agent replies once, and the result is returned.
    If the source agent needs to continue the conversation, it can call this tool again.
    """
    agent_name = args.get("agent_name", "").strip()
    message_text = args.get("message", "").strip()

    if not agent_name or not message_text:
        return "❌ Please provide target agent name and message content"

    try:
        from app.models.participant import Participant
        from datetime import datetime, timezone

        async with async_session() as db:
            # Look up source agent
            src_result = await db.execute(select(AgentModel).where(AgentModel.id == from_agent_id))

            source_agent = src_result.scalar_one_or_none()
            source_name = source_agent.name if source_agent else "Unknown agent"
            source_tenant_id = source_agent.tenant_id if source_agent else None

            # Build base filter: same tenant + not self
            base_filter = [AgentModel.id != from_agent_id]
            if source_tenant_id:
                base_filter.append(AgentModel.tenant_id == source_tenant_id)

            # Find target agent by name — exact match first, then fuzzy
            target = None
            exact_result = await db.execute(
                select(AgentModel).where(AgentModel.name == agent_name, *base_filter)
            )
            target = exact_result.scalars().first()
            if not target:
                safe_name = agent_name.replace("%", "").replace("_", r"\_")
                fuzzy_result = await db.execute(
                    select(AgentModel).where(AgentModel.name.ilike(f"%{safe_name}%"), *base_filter)
                )
                target = fuzzy_result.scalars().first()
            if not target:
                # Only show agents from relationships, not all agents
                rel_r = await db.execute(
                    select(AgentModel.name).join(
                        AgentAgentRelationship,
                        (AgentAgentRelationship.target_agent_id == AgentModel.id) & (AgentAgentRelationship.agent_id == from_agent_id)
                    )
                )
                rel_names = [n for (n,) in rel_r.all()]
                return f"❌ No agent found matching '{agent_name}'. Your connected colleagues: {', '.join(rel_names) if rel_names else 'none — ask your administrator to set up relationships'}"


            # Check if target agent has expired
            if target.is_expired or (target.expires_at and datetime.now(timezone.utc) >= target.expires_at):
                return f"⚠️ {target.name} is currently unavailable — their service period has ended. Please contact the platform administrator."

            # Enforce relationship: only allow communication with agents in relationships
            from app.models.org import AgentAgentRelationship
            rel_check = await db.execute(
                select(AgentAgentRelationship.id).where(
                    ((AgentAgentRelationship.agent_id == from_agent_id) & (AgentAgentRelationship.target_agent_id == target.id))
                    | ((AgentAgentRelationship.agent_id == target.id) & (AgentAgentRelationship.target_agent_id == from_agent_id))
                ).limit(1)
            )
            if not rel_check.scalar_one_or_none():
                return f"❌ You do not have a relationship with {target.name}. Only agents in your relationship list can be contacted. Ask your administrator to add a relationship if needed."

            src_part_r = await db.execute(select(Participant).where(Participant.type == "agent", Participant.ref_id == from_agent_id))
            src_participant = src_part_r.scalar_one_or_none()
            tgt_part_r = await db.execute(select(Participant).where(Participant.type == "agent", Participant.ref_id == target.id))
            tgt_participant = tgt_part_r.scalar_one_or_none()

            # Find or create ChatSession for this agent pair (ordered consistently)
            session_agent_id = min(from_agent_id, target.id, key=str)
            session_peer_id = max(from_agent_id, target.id, key=str)
            sess_r = await db.execute(
                select(ChatSession).where(
                    ChatSession.agent_id == session_agent_id,
                    ChatSession.peer_agent_id == session_peer_id,
                    ChatSession.source_channel == "agent",
                )
            )
            chat_session = sess_r.scalar_one_or_none()
            owner_id = source_agent.creator_id if source_agent else from_agent_id
            if not chat_session:
                src_part_id = src_participant.id if src_participant else None
                chat_session = ChatSession(
                    agent_id=session_agent_id,
                    user_id=owner_id,
                    title=f"{source_name} ↔ {target.name}",
                    source_channel="agent",
                    participant_id=src_part_id,
                    peer_agent_id=session_peer_id,
                )
                db.add(chat_session)
                await db.flush()

            session_id = str(chat_session.id)

            # ── OpenClaw target: queue message for gateway poll ──
            if getattr(target, "agent_type", "native") == "openclaw":
                # 1. Save the source message to the chat session
                db.add(ChatMessage(
                    agent_id=session_agent_id,
                    user_id=owner_id,
                    role="user",
                    content=message_text,
                    conversation_id=session_id,
                    participant_id=src_participant.id if src_participant else None,
                ))
                chat_session.last_message_at = datetime.now(timezone.utc)
                
                # 2. Queue for Gateway
                from app.models.gateway_message import GatewayMessage as GMsg
                gw_msg = GMsg(
                    agent_id=target.id,
                    sender_agent_id=from_agent_id,
                    sender_user_id=owner_id,
                    content=f"[From {source_name}] {message_text}",
                    status="pending",
                    conversation_id=session_id,
                )
                db.add(gw_msg)
                await db.commit()
                
                # 3. Log activity
                from app.services.activity_logger import log_activity
                await log_activity(
                    from_agent_id, "agent_msg_sent",
                    f"Sent message to {target.name} (queued)",
                    detail={"partner": target.name, "message": message_text[:200]},
                )

                online = target.openclaw_last_seen and (datetime.now(timezone.utc) - target.openclaw_last_seen).total_seconds() < 300
                status_hint = "online" if online else "offline (message will be delivered on next heartbeat)"
                return f"✅ Message sent to {target.name} (OpenClaw agent, currently {status_hint}). The message has been queued and will be delivered when the agent polls for updates."

            # Prepare target LLM
            from app.services.agent_context import build_agent_context
            from app.models.llm import LLMModel

            # Load primary model (with fallback support)
            target_model = None
            if target.primary_model_id:
                model_r = await db.execute(select(LLMModel).where(LLMModel.id == target.primary_model_id))
                target_model = model_r.scalar_one_or_none()

            # Config-level fallback: primary missing -> use fallback
            if not target_model and target.fallback_model_id:
                fb_r = await db.execute(select(LLMModel).where(LLMModel.id == target.fallback_model_id))
                target_model = fb_r.scalar_one_or_none()
                if target_model:
                    logger.warning(f"[A2A] Primary model unavailable for {target.name}, using fallback: {target_model.model}")

            if not target_model:
                return f"⚠️ {target.name} has no LLM model configured"

            # Build target system prompt
            target_static, target_dynamic = await build_agent_context(target.id, target.name, target.role_description or "")
            target_dynamic += (
                "\n\n--- Agent-to-Agent Message ---\n"
                "You are receiving a message from another digital employee. "
                "Reply concisely and helpfully. Focus on the request and provide a clear answer.\n"
                "\n** CRITICAL FILE DELIVERY RULE **\n"
                "After you write any file (report, document, analysis, etc.) that the requesting agent needs, "
                "you MUST call `send_file_to_agent(agent_name=\"<requester_name>\", file_path=\"<path>\")` "
                "to deliver it. The other agent CANNOT access your workspace. "
                "Never just tell them the path — always deliver explicitly.\n"
            )

            # Load recent history for context
            conversation_messages: list[dict] = []
            hist_result = await db.execute(
                select(ChatMessage)
                .where(
                    ChatMessage.conversation_id == session_id,
                    ChatMessage.agent_id == session_agent_id,
                )
                .order_by(ChatMessage.created_at.desc())
                .limit(20)
            )
            for m in reversed(hist_result.scalars().all()):
                if m.participant_id and src_participant and m.participant_id == src_participant.id:
                    role = "user"
                else:
                    role = "assistant"
                conversation_messages.append({"role": role, "content": m.content})

            # Add the new message from source
            conversation_messages.append({"role": "user", "content": f"[From {source_name}] {message_text}"})

            # Save source message
            owner_id = source_agent.creator_id if source_agent else from_agent_id
            db.add(ChatMessage(
                agent_id=session_agent_id,
                user_id=owner_id,
                role="user",
                content=message_text,
                conversation_id=session_id,
                participant_id=src_participant.id if src_participant else None,
            ))
            chat_session.last_message_at = datetime.now(timezone.utc)
            await db.commit()

            # Call target LLM with tool support (multi-round)
            import asyncio
            import random
            import httpx
            from app.services.llm_utils import (
                get_provider_base_url,
                create_llm_client,
                LLMMessage,
            )
            from app.services.llm_client import LLMError
            from app.services.agent_tools import get_agent_tools_for_llm, execute_tool
            base_url = get_provider_base_url(target_model.provider, target_model.base_url)
            if not base_url:
                return f"⚠️ {target.name}'s model has no API base URL configured"

            full_msgs: list[LLMMessage] = [LLMMessage(role="system", content=target_static, dynamic_content=target_dynamic)] + [
                LLMMessage(role=m["role"], content=m["content"]) for m in conversation_messages
            ]

            # Load tools for target agent
            tools_for_llm = await get_agent_tools_for_llm(target.id)

            max_tool_rounds = target.max_tool_rounds or 50
            target_reply = ""
            _a2a_accumulated_tokens = 0

            from app.services.token_tracker import record_token_usage, extract_usage_tokens, estimate_tokens_from_chars

            llm_client = create_llm_client(
                provider=target_model.provider,
                api_key=target_model.api_key_encrypted,
                model=target_model.model,
                base_url=base_url,
                timeout=120.0,
            )
            _A2A_RETRYABLE_MARKERS = (
                "http 408", "http 429", "http 500", "http 502", "http 503", "http 504",
                "timeout", "timed out", "connection failed", "temporarily unavailable", "rate limit",
            )
            _A2A_MAX_RETRIES = 3

            def _is_retryable_llm_error(exc: Exception) -> bool:
                """Determine whether an LLM exception is transient and worth retrying."""
                if isinstance(exc, (httpx.TimeoutException, httpx.TransportError)):
                    return True
                if isinstance(exc, LLMError):
                    lowered = (str(exc) or "").lower()
                    return any(m in lowered for m in _A2A_RETRYABLE_MARKERS)
                return False

            try:
                for _round in range(max_tool_rounds):
                    response = None
                    for attempt in range(1, _A2A_MAX_RETRIES + 1):
                        try:
                            response = await llm_client.complete(
                                messages=full_msgs,
                                tools=tools_for_llm if tools_for_llm else None,
                                temperature=target_model.temperature,
                                max_tokens=4096,
                            )
                            break
                        except Exception as llm_exc:
                            if not _is_retryable_llm_error(llm_exc) or attempt >= _A2A_MAX_RETRIES:
                                raise

                            err_text = str(llm_exc) or type(llm_exc).__name__
                            # Exponential backoff with jitter to prevent thundering herd
                            backoff = (2 ** (attempt - 1)) + random.uniform(0, 0.5)
                            logger.warning(
                                f"[A2A] LLM call failed for {target.name} (round={_round + 1}, "
                                f"attempt={attempt}/{_A2A_MAX_RETRIES}): {err_text[:200]}. "
                                f"Retrying in {backoff:.1f}s"
                            )
                            await asyncio.sleep(backoff)

                    if response is None:
                        raise RuntimeError("A2A LLM response is unexpectedly empty after retries")

                    # Track tokens from API response
                    real_tokens = extract_usage_tokens(response.usage)
                    if real_tokens:
                        _a2a_accumulated_tokens += real_tokens
                    else:
                        round_chars = sum(len(m.content or '') for m in full_msgs if isinstance(m.content, str))
                        _a2a_accumulated_tokens += estimate_tokens_from_chars(round_chars)

                    # Check for tool calls
                    if response.tool_calls:
                        # Add assistant message with tool calls to conversation
                        full_msgs.append(LLMMessage(
                            role="assistant",
                            content=response.content or None,
                            tool_calls=[{
                                "id": tc.get("id", ""),
                                "type": "function",
                                "function": tc.get("function", {}),
                            } for tc in response.tool_calls],
                            reasoning_content=response.reasoning_content,
                        ))

                        # Execute each tool call
                        for tc in response.tool_calls:
                            fn = tc.get("function", {})
                            tool_name = fn.get("name", "")
                            raw_args = fn.get("arguments", "{}")
                            if isinstance(raw_args, dict):
                                tool_args = raw_args
                            else:
                                try:
                                    tool_args = json.loads(raw_args) if raw_args else {}
                                except Exception:
                                    tool_args = {}

                            tool_result = await execute_tool(tool_name, tool_args, target.id, owner_id)

                            # Nudge: after write_file in A2A, remind to deliver via send_file_to_agent
                            if tool_name == "write_file" and isinstance(tool_result, str) and tool_result.startswith("\u2705"):
                                wrote_path = tool_args.get("path", "")
                                tool_result += (
                                    f"\n\n⚠️ REMINDER: The requesting agent ({source_name}) cannot access your workspace. "
                                    f"You MUST now call `send_file_to_agent(agent_name=\"{source_name}\", file_path=\"{wrote_path}\")` "
                                    f"to deliver this file to them."
                                )

                            # Save tool_call to DB so it appears in chat history
                            try:
                                async with async_session() as _tc_db:
                                    _tc_db.add(ChatMessage(
                                        agent_id=session_agent_id,
                                        user_id=owner_id,
                                        role="tool_call",
                                        content=json.dumps({
                                            "name": tool_name,
                                            "args": tool_args,
                                            "status": "done",
                                            "result": str(tool_result)[:500],
                                        }, ensure_ascii=False),
                                        conversation_id=session_id,
                                        participant_id=tgt_participant.id if tgt_participant else None,
                                    ))
                                    await _tc_db.commit()
                            except Exception as _tc_err:
                                logger.error(f"[A2A] Failed to save tool_call: {_tc_err}")

                            # Add tool result to conversation
                            full_msgs.append(LLMMessage(
                                role="tool",
                                tool_call_id=tc.get("id", ""),
                                content=str(tool_result)[:4000],
                            ))
                        continue  # Next LLM round

                    # No tool calls — this is the final text response
                    target_reply = response.content or ""
                    break
            finally:
                await llm_client.close()

            # Record accumulated A2A tokens for the target agent
            if _a2a_accumulated_tokens > 0:
                await record_token_usage(target.id, _a2a_accumulated_tokens)

            if not target_reply:
                return f"⚠️ {target.name} did not respond (LLM returned empty)"

            # Save target reply
            async with async_session() as db2:
                part_r = await db2.execute(select(Participant).where(Participant.type == "agent", Participant.ref_id == target.id))
                tgt_part = part_r.scalar_one_or_none()
                db2.add(ChatMessage(
                    agent_id=session_agent_id,
                    user_id=owner_id,
                    role="assistant",
                    content=target_reply,
                    conversation_id=session_id,
                    participant_id=tgt_part.id if tgt_part else None,
                ))
                await db2.commit()

            # Log activity
            from app.services.activity_logger import log_activity
            await log_activity(
                target.id, "agent_msg_sent",
                f"Replied to message from {source_name}",
                detail={"partner": source_name, "message": message_text[:200], "reply": target_reply[:200]},
            )
            await log_activity(
                from_agent_id, "agent_msg_sent",
                f"Sent message to {target.name} and received reply",
                detail={"partner": target.name, "message": message_text[:200], "reply": target_reply[:200]},
            )

            return f"💬 {target.name} replied:\n{target_reply}"

    except Exception as e:
        logger.exception(
            f"[A2A] send_message_to_agent failed: from={from_agent_id}, to={args.get('agent_name', '')}"
        )
        error_type = type(e).__name__
        error_detail = (str(e) or "").strip()
        if not error_detail:
            timeout_types = {"ReadTimeout", "ConnectTimeout", "TimeoutException"}
            if error_type in timeout_types:
                error_detail = "LLM request timed out while waiting for target agent response"
            else:
                error_detail = "No detailed error message returned from upstream"
        return f"❌ Message send error ({error_type}): {error_detail[:200]}"



# Plaza Tools — Agent Square social feed
# ═══════════════════════════════════════════════════════

async def _plaza_get_new_posts(agent_id: uuid.UUID, arguments: dict) -> str:
    """Get recent posts from the Agent Plaza, scoped to agent's tenant."""
    from app.models.plaza import PlazaPost, PlazaComment
    from app.models.agent import Agent as AgentModel
    from sqlalchemy import desc

    limit = min(arguments.get("limit", 10), 20)

    try:
        async with async_session() as db:
            # Resolve agent's tenant_id
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            tenant_id = agent.tenant_id if agent else None

            q = select(PlazaPost).order_by(desc(PlazaPost.created_at)).limit(limit)
            if tenant_id:
                q = q.where(PlazaPost.tenant_id == tenant_id)
            result = await db.execute(q)
            posts = result.scalars().all()

            if not posts:
                return "📭 No posts in the plaza yet. Be the first to share something!"

            output = []
            for p in posts:
                # Load comments
                cr = await db.execute(
                    select(PlazaComment).where(PlazaComment.post_id == p.id).order_by(PlazaComment.created_at).limit(5)
                )
                comments = cr.scalars().all()
                icon = "🤖" if p.author_type == "agent" else "👤"
                time_str = p.created_at.strftime("%m-%d %H:%M") if p.created_at else ""
                post_text = f"{icon} **{p.author_name}** ({time_str}) [post_id: {p.id}]\n{p.content}\n❤️ {p.likes_count}  💬 {p.comments_count}"
                if comments:
                    for c in comments:
                        c_icon = "🤖" if c.author_type == "agent" else "👤"
                        post_text += f"\n  └─ {c_icon} {c.author_name}: {c.content}"
                output.append(post_text)

            return "🏛️ Agent Plaza — Recent Posts:\n\n" + "\n\n---\n\n".join(output)

    except Exception as e:
        return f"❌ Failed to load plaza posts: {str(e)[:200]}"


async def _plaza_create_post(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new post in the Agent Plaza."""
    from app.models.plaza import PlazaPost
    from app.models.agent import Agent as AgentModel

    content = arguments.get("content", "").strip()
    if not content:
        return "Error: Post content cannot be empty."
    if len(content) > 500:
        content = content[:500]

    try:
        async with async_session() as db:
            # Get agent name
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            if not agent:
                return "Error: Agent not found."

            post = PlazaPost(
                author_id=agent_id,
                author_type="agent",
                author_name=agent.name,
                content=content,
                tenant_id=agent.tenant_id,
            )
            db.add(post)
            await db.flush()  # get post.id

            # Extract @mentions
            try:
                import re
                mentions = re.findall(r'@(\S+)', content)
                if mentions:
                    from app.services.notification_service import send_notification
                    a_q = select(AgentModel).where(AgentModel.id != agent_id)
                    if agent.tenant_id:
                        a_q = a_q.where(AgentModel.tenant_id == agent.tenant_id)
                    a_map = {a.name.lower(): a for a in (await db.execute(a_q)).scalars().all()}
                    notified = set()
                    for m in mentions:
                        ma = a_map.get(m.lower())
                        if ma and ma.id not in notified:
                            notified.add(ma.id)
                            await send_notification(
                                db, agent_id=ma.id,
                                type="mention",
                                title=f"{agent.name} mentioned you in a plaza post",
                                body=content[:150],
                                link=f"/plaza?post={post.id}",
                                ref_id=post.id,
                                sender_name=agent.name,
                            )
            except Exception:
                pass

            await db.commit()
            await db.refresh(post)
            return f"Post published! (ID: {post.id})"

    except Exception as e:
        return f"Failed to create post: {str(e)[:200]}"


async def _plaza_add_comment(agent_id: uuid.UUID, arguments: dict) -> str:
    """Add a comment to a plaza post."""
    from app.models.plaza import PlazaPost, PlazaComment
    from app.models.agent import Agent as AgentModel

    post_id = arguments.get("post_id", "")
    content = arguments.get("content", "").strip()
    if not content:
        return "Error: Comment content cannot be empty."
    if len(content) > 300:
        content = content[:300]

    try:
        pid = uuid.UUID(str(post_id))
    except Exception:
        return "Error: Invalid post_id format."

    try:
        async with async_session() as db:
            # Verify post exists
            pr = await db.execute(select(PlazaPost).where(PlazaPost.id == pid))
            post = pr.scalar_one_or_none()
            if not post:
                return "Error: Post not found."

            # Get agent name
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            if not agent:
                return "Error: Agent not found."

            comment = PlazaComment(
                post_id=pid,
                author_id=agent_id,
                author_type="agent",
                author_name=agent.name,
                content=content,
            )
            db.add(comment)
            post.comments_count = (post.comments_count or 0) + 1

            # Notify post author (if not self)
            if post.author_id != agent_id:
                try:
                    from app.services.notification_service import send_notification
                    if post.author_type == "agent":
                        await send_notification(
                            db, agent_id=post.author_id,
                            type="plaza_reply",
                            title=f"{agent.name} commented on your post",
                            body=content[:150],
                            link=f"/plaza?post={pid}",
                            ref_id=pid,
                            sender_name=agent.name,
                        )
                        # Also notify human creator
                        pa = (await db.execute(select(AgentModel).where(AgentModel.id == post.author_id))).scalar_one_or_none()
                        if pa and pa.creator_id:
                            await send_notification(
                                db, user_id=pa.creator_id,
                                type="plaza_comment",
                                title=f"{agent.name} commented on {pa.name}'s post",
                                body=content[:100],
                                link=f"/plaza?post={pid}",
                                ref_id=pid,
                                sender_name=agent.name,
                            )
                    elif post.author_type == "human":
                        await send_notification(
                            db, user_id=post.author_id,
                            type="plaza_reply",
                            title=f"{agent.name} commented on your post",
                            body=content[:150],
                            link=f"/plaza?post={pid}",
                            ref_id=pid,
                            sender_name=agent.name,
                        )
                except Exception:
                    pass

            # Notify other agents who commented on this post
            try:
                from app.services.notification_service import send_notification
                other_crs = await db.execute(
                    select(PlazaComment.author_id, PlazaComment.author_type)
                    .where(PlazaComment.post_id == pid)
                    .distinct()
                )
                notified = {post.author_id, agent_id}
                for row in other_crs.fetchall():
                    cid, ctype = row
                    if cid in notified:
                        continue
                    notified.add(cid)
                    if ctype == "agent":
                        await send_notification(
                            db, agent_id=cid,
                            type="plaza_reply",
                            title=f"{agent.name} also commented on a post you commented on",
                            body=content[:150],
                            link=f"/plaza?post={pid}",
                            ref_id=pid,
                            sender_name=agent.name,
                        )
            except Exception:
                pass

            # Extract @mentions
            try:
                import re
                mentions = re.findall(r'@(\S+)', content)
                if mentions:
                    from app.services.notification_service import send_notification
                    from app.models.user import User
                    # Load agents in tenant
                    a_q = select(AgentModel).where(AgentModel.id != agent_id)
                    if agent.tenant_id:
                        a_q = a_q.where(AgentModel.tenant_id == agent.tenant_id)
                    a_map = {a.name.lower(): a for a in (await db.execute(a_q)).scalars().all()}
                    notified_m = set()
                    for m in mentions:
                        ma = a_map.get(m.lower())
                        if ma and ma.id not in notified_m:
                            notified_m.add(ma.id)
                            await send_notification(
                                db, agent_id=ma.id,
                                type="mention",
                                title=f"{agent.name} mentioned you in a comment",
                                body=content[:150],
                                link=f"/plaza?post={pid}",
                                ref_id=pid,
                                sender_name=agent.name,
                            )
            except Exception:
                pass

            await db.commit()
            return f"Comment added to post by {post.author_name}."

    except Exception as e:
        return f"Failed to add comment: {str(e)[:200]}"


# ─── Code Execution ─────────────────────────────────────────────

# Dangerous patterns to block (for legacy fallback)
_DANGEROUS_BASH = [
    "rm -rf /", "rm -rf ~", "sudo ", "mkfs", "dd if=",
    ":(){ :", "chmod 777 /", "chown ", "shutdown", "reboot",
    "curl ", "wget ", "nc ", "ncat ", "ssh ", "scp ",
    "python3 -c", "python -c",
]

_DANGEROUS_PYTHON_IMPORTS = [
    "subprocess", "shutil.rmtree", "os.system", "os.popen",
    "os.exec", "os.spawn",
    "socket", "http.client", "urllib.request", "requests",
    "ftplib", "smtplib", "telnetlib", "ctypes",
    "__import__", "importlib",
]


def _check_code_safety(language: str, code: str) -> str | None:
    """Check code for dangerous patterns. Returns error message if unsafe, None if ok."""
    code_lower = code.lower()

    if language == "bash":
        for pattern in _DANGEROUS_BASH:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: dangerous command detected ({pattern.strip()})"
        # Block deep path traversal outside workspace
        if "../../" in code:
            return "❌ Blocked: directory traversal not allowed"

    elif language == "python":
        for pattern in _DANGEROUS_PYTHON_IMPORTS:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: unsafe operation detected ({pattern})"

    elif language == "node":
        dangerous_node = ["child_process", "fs.rmSync", "fs.rmdirSync", "process.exit",
                          "require('http')", "require('https')", "require('net')"]
        for pattern in dangerous_node:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: unsafe operation detected ({pattern})"

    return None


async def _execute_code(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Execute code using the configured sandbox backend."""
    language = arguments.get("language", "python")
    code = arguments.get("code", "")
    timeout = min(arguments.get("timeout", 30), 60)  # Max 60 seconds

    if not code.strip():
        return "❌ No code provided"

    if language not in ("python", "bash", "node"):
        return f"❌ Unsupported language: {language}. Use: python, bash, or node"

    # Working directory is the agent's root directory (must be absolute)
    # This allows code to access skills/, workspace/, memory/ etc. directly
    work_dir = ws.resolve()
    work_dir.mkdir(parents=True, exist_ok=True)

    try:
        # Import here to avoid circular imports
        from app.config import get_sandbox_config
        from app.services.sandbox.config import SandboxConfig
        from app.services.sandbox.registry import get_sandbox_backend

        # Get sandbox config: prefer tool config from DB, fallback to env vars
        fallback_config = get_sandbox_config()
        tool_config = await _get_tool_config(agent_id, "execute_code")

        if tool_config:
            sandbox_config = SandboxConfig.from_dict(tool_config, fallback_config)
        else:
            sandbox_config = fallback_config
            logger.info(f"[Sandbox] Using fallback config for agent {agent_id}")

        backend = get_sandbox_backend(sandbox_config)
        logger.info(f"[Sandbox] Executing code with backend: {backend.__class__.__name__}")
        result = await backend.execute(
            code=code,
            language=language,
            timeout=timeout,
            work_dir=str(work_dir),
        )

        # Format result for user display
        return backend._format_result(result)

    except ValueError as e:
        # Sandbox disabled or misconfigured - fall back to legacy subprocess
        logger.warning(f"[Sandbox] Config issue, falling back to legacy: {e}")
        return await _execute_code_legacy(ws, arguments)

    except Exception as e:
        logger.exception(f"[Sandbox] Execution failed for agent {agent_id}")
        # Try fallback to legacy subprocess
        try:
            return await _execute_code_legacy(ws, arguments)
        except Exception as fallback_error:
            logger.exception(f"[Sandbox] Fallback also failed for agent {agent_id}")
            return f"❌ Execution error: {str(e)[:200]}"


async def _execute_code_legacy(ws: Path, arguments: dict) -> str:
    """Legacy subprocess-based code execution (fallback)."""
    import asyncio

    language = arguments.get("language", "python")
    code = arguments.get("code", "")
    timeout = min(arguments.get("timeout", 30), 60)

    if not code.strip():
        return "❌ No code provided"

    if language not in ("python", "bash", "node"):
        return f"❌ Unsupported language: {language}. Use: python, bash, or node"

    # Security check
    safety_error = _check_code_safety(language, code)
    if safety_error:
        return safety_error

    # Working directory is the agent's root directory (must be absolute)
    # This allows code to access skills/, workspace/, memory/ etc. directly
    work_dir = ws.resolve()
    work_dir.mkdir(parents=True, exist_ok=True)

    # Determine command and file extension
    if language == "python":
        ext = ".py"
        cmd_prefix = ["python3"]
    elif language == "bash":
        ext = ".sh"
        cmd_prefix = ["bash"]
    elif language == "node":
        ext = ".js"
        cmd_prefix = ["node"]
    else:
        return f"❌ Unsupported language: {language}"

    # Write code to a temp file inside workspace
    script_path = work_dir / f"_exec_tmp{ext}"
    try:
        script_path.write_text(code, encoding="utf-8")

        # Inherit parent environment but override HOME to workspace
        safe_env = dict(os.environ)
        safe_env["HOME"] = str(work_dir)
        safe_env["PYTHONDONTWRITEBYTECODE"] = "1"

        proc = await asyncio.create_subprocess_exec(
            *cmd_prefix, str(script_path),
            cwd=str(work_dir),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            env=safe_env,
        )

        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            return f"❌ Code execution timed out after {timeout}s"

        stdout_str = stdout.decode("utf-8", errors="replace")[:10000]
        stderr_str = stderr.decode("utf-8", errors="replace")[:5000]

        result_parts = []
        if stdout_str.strip():
            result_parts.append(f"📤 Output:\n{stdout_str}")
        if stderr_str.strip():
            result_parts.append(f"⚠️ Stderr:\n{stderr_str}")
        if proc.returncode != 0:
            result_parts.append(f"Exit code: {proc.returncode}")

        if not result_parts:
            return "✅ Code executed successfully (no output)"

        return "\n\n".join(result_parts)

    except Exception as e:
        return f"❌ Execution error: {str(e)[:200]}"
    finally:
        # Clean up temp script
        try:
            script_path.unlink(missing_ok=True)
        except Exception:
            pass


# ─── Resource Discovery Executors ───────────────────────────────

async def _discover_resources(arguments: dict) -> str:
    """Search Smithery registry for MCP servers."""
    query = arguments.get("query", "")
    if not query:
        return "❌ Please provide a search query describing the capability you need."
    max_results = min(arguments.get("max_results", 5), 10)

    from app.services.resource_discovery import search_smithery
    return await search_smithery(query, max_results)


async def _import_mcp_server(agent_id: uuid.UUID, arguments: dict) -> str:
    """Import an MCP server — either from Smithery or by direct URL."""
    config = arguments.get("config") or {}
    reauthorize = arguments.get("reauthorize", False)
    mcp_url = config.pop("mcp_url", None) if isinstance(config, dict) else None

    if mcp_url:
        # Direct URL import — bypass Smithery
        from app.services.resource_discovery import import_mcp_direct
        server_name = arguments.get("server_id") or config.pop("server_name", None)
        api_key = config.pop("api_key", None)
        return await import_mcp_direct(mcp_url, agent_id, server_name, api_key)

    # Smithery import
    server_id = arguments.get("server_id", "")
    if not server_id:
        return "❌ Please provide a server_id (e.g. 'github'). Use discover_resources first to find available servers."

    from app.services.resource_discovery import import_mcp_from_smithery
    return await import_mcp_from_smithery(server_id, agent_id, config or None, reauthorize=reauthorize)


# ─── Trigger Management Handlers (Aware Engine) ────────────────────

MAX_TRIGGERS_PER_AGENT = 20
VALID_TRIGGER_TYPES = {"cron", "once", "interval", "poll", "on_message", "webhook"}


async def _handle_set_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new trigger for the agent."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    ttype = arguments.get("type", "").strip()
    config = arguments.get("config", {})
    reason = arguments.get("reason", "").strip()
    focus_ref = arguments.get("focus_ref", "") or arguments.get("agenda_ref", "")  # backward compat

    if not name:
        return "❌ Missing required argument 'name'"
    if ttype not in VALID_TRIGGER_TYPES:
        return f"❌ Invalid trigger type '{ttype}'. Valid types: {', '.join(VALID_TRIGGER_TYPES)}"
    if not reason:
        return "❌ Missing required argument 'reason'"

    # Validate type-specific config
    if ttype == "cron":
        expr = config.get("expr", "")
        if not expr:
            return "❌ cron trigger requires config.expr, e.g. {\"expr\": \"0 9 * * *\"}"
        try:
            from croniter import croniter
            croniter(expr)
        except Exception:
            return f"❌ Invalid cron expression: '{expr}'"
    elif ttype == "once":
        if not config.get("at"):
            return "❌ once trigger requires config.at, e.g. {\"at\": \"2026-03-10T09:00:00+08:00\"}"
    elif ttype == "interval":
        if not config.get("minutes"):
            return "❌ interval trigger requires config.minutes, e.g. {\"minutes\": 30}"
    elif ttype == "poll":
        if not config.get("url"):
            return "❌ poll trigger requires config.url"
    elif ttype == "on_message":
        if not config.get("from_agent_name") and not config.get("from_user_name"):
            return "❌ on_message trigger requires config.from_agent_name (for agents) or config.from_user_name (for human users on Feishu/Slack/Discord)"
        # Snapshot the latest message timestamp so we only detect NEW messages after this point
        # This prevents false positives from already-processed messages
        try:
            from app.models.audit import ChatMessage
            from app.models.chat_session import ChatSession
            from sqlalchemy import cast as sa_cast, String as SaString
            async with async_session() as _snap_db:
                _snap_q = select(ChatMessage.created_at).join(
                    ChatSession, ChatMessage.conversation_id == sa_cast(ChatSession.id, SaString)
                ).where(
                    ChatSession.agent_id == agent_id,
                    ChatMessage.created_at.isnot(None),
                ).order_by(ChatMessage.created_at.desc()).limit(1)
                _snap_r = await _snap_db.execute(_snap_q)
                _latest_ts = _snap_r.scalar_one_or_none()
                if _latest_ts:
                    config["_since_ts"] = _latest_ts.isoformat()
        except Exception:
            pass  # Fallback to trigger.created_at in the daemon
    elif ttype == "webhook":
        # Auto-generate a unique token for the webhook URL
        import secrets
        token = secrets.token_urlsafe(8)  # ~11 chars, URL-safe
        config["token"] = token

    try:
        async with async_session() as db:
            # Load agent to get per-agent trigger limit
            from app.models.agent import Agent as _AgentModel
            _a_result = await db.execute(select(_AgentModel).where(_AgentModel.id == agent_id))
            _agent_obj = _a_result.scalar_one_or_none()
            agent_max_triggers = (_agent_obj.max_triggers if _agent_obj else None) or MAX_TRIGGERS_PER_AGENT

            # Check max triggers
            from sqlalchemy import func as sa_func
            result = await db.execute(
                select(sa_func.count()).select_from(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.is_enabled == True,
                )
            )
            count = result.scalar() or 0
            if count >= agent_max_triggers:
                return f"❌ Maximum trigger limit reached ({agent_max_triggers}). Cancel some triggers first."

            # Check for duplicate name
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            existing = result.scalar_one_or_none()
            if existing:
                if existing.is_enabled:
                    return f"❌ Trigger '{name}' already exists and is active. Use update_trigger to modify it, or cancel_trigger first."
                else:
                    # Re-enable disabled trigger with new config (preserve fire history)
                    # For webhook triggers: reuse the old token so the URL stays stable
                    if ttype == "webhook":
                        old_token = (existing.config or {}).get("token")
                        if old_token:
                            config["token"] = old_token
                    existing.type = ttype
                    existing.config = config
                    existing.reason = reason
                    existing.focus_ref = focus_ref or None
                    existing.is_enabled = True
                    # Keep fire_count and last_fired_at — they are cumulative stats
                    await db.commit()
                    return f"✅ Trigger '{name}' re-enabled with new configuration ({ttype}, fired {existing.fire_count} times so far)"

            trigger = AgentTrigger(
                agent_id=agent_id,
                name=name,
                type=ttype,
                config=config,
                reason=reason,
                focus_ref=focus_ref or None,
            )
            db.add(trigger)
            await db.commit()

        # Activity log
        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_created", {
                "name": name, "type": ttype, "reason": reason[:100],
            }, agent_id=agent_id)
        except Exception:
            pass

        # Return webhook URL for webhook triggers
        if ttype == "webhook":
            from app.services.platform_service import platform_service
            base = await platform_service.get_public_base_url()
            webhook_url = f"{base.rstrip('/')}/api/webhooks/t/{config['token']}"

            return f"✅ Webhook trigger '{name}' created.\n\nWebhook URL: {webhook_url}\n\nTell the user to configure this URL in their external service (e.g. GitHub, Grafana). When the service sends a POST to this URL, you will be woken up with the payload as context."

        return f"✅ Trigger '{name}' created ({ttype}). It will fire according to your config and wake you up with the reason as context."

    except Exception as e:
        return f"❌ Failed to create trigger: {e}"


async def _handle_update_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    """Update an existing trigger's config or reason."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    if not name:
        return "❌ Missing required argument 'name'"

    new_config = arguments.get("config")
    new_reason = arguments.get("reason")

    if new_config is None and new_reason is None:
        return "❌ Provide at least one of 'config' or 'reason' to update"

    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            trigger = result.scalar_one_or_none()
            if not trigger:
                return f"❌ Trigger '{name}' not found"

            changes = []
            if new_config is not None:
                old_config = trigger.config
                trigger.config = new_config
                changes.append(f"config: {old_config} → {new_config}")
            if new_reason is not None:
                trigger.reason = new_reason
                changes.append(f"reason updated")

            await db.commit()

        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_updated", {
                "name": name, "changes": "; ".join(changes),
            }, agent_id=agent_id)
        except Exception:
            pass

        return f"✅ Trigger '{name}' updated: {'; '.join(changes)}"

    except Exception as e:
        return f"❌ Failed to update trigger: {e}"


async def _handle_cancel_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    """Cancel (disable) a trigger by name."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    if not name:
        return "❌ Missing required argument 'name'"

    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            trigger = result.scalar_one_or_none()
            if not trigger:
                return f"❌ Trigger '{name}' not found"
            if not trigger.is_enabled:
                return f"ℹ️ Trigger '{name}' is already disabled"

            trigger.is_enabled = False
            await db.commit()

        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_cancelled", {"name": name}, agent_id=agent_id)
        except Exception:
            pass

        return f"✅ Trigger '{name}' cancelled. It will no longer fire."

    except Exception as e:
        return f"❌ Failed to cancel trigger: {e}"


async def _handle_list_triggers(agent_id: uuid.UUID) -> str:
    """List all active triggers for the agent."""
    from app.models.trigger import AgentTrigger

    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                ).order_by(AgentTrigger.created_at.desc())
            )
            triggers = result.scalars().all()

        if not triggers:
            return "No triggers found. Use set_trigger to create one."

        lines = ["| Name | Type | Config | Reason | Status | Fires |", "|------|------|--------|--------|--------|-------|"]
        for t in triggers:
            status = "✅ active" if t.is_enabled else "⏸ disabled"
            config_str = str(t.config)[:50]
            reason_str = t.reason[:40] if t.reason else ""
            lines.append(f"| {t.name} | {t.type} | {config_str} | {reason_str} | {status} | {t.fire_count} |")

        return "\n".join(lines)

    except Exception as e:
        return f"❌ Failed to list triggers: {e}"


# ─── Image Upload (ImageKit CDN) ────────────────────────────────

async def _upload_image(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Upload an image to ImageKit CDN and return the public URL.

    Credential resolution order:
    1. Global tool config (admin-set, shared by all agents)
    2. Per-agent tool config override (agent-specific)
    """
    import httpx
    import base64

    file_path = arguments.get("file_path")
    url = arguments.get("url")
    file_name = arguments.get("file_name")
    folder = arguments.get("folder", "/clawith")

    if not file_path and not url:
        return "❌ Please provide either 'file_path' (workspace path) or 'url' (public image URL)"

    # ── Load ImageKit credentials (Agent > Company priority) ──
    private_key = ""
    url_endpoint = ""
    try:
        # Use standard _get_tool_config (Agent > Company, cached, schema-aware decryption)
        config = await _get_tool_config(agent_id, "upload_image") or {}
        private_key = config.get("private_key", "")
        url_endpoint = config.get("url_endpoint", "")
    except Exception as e:
        logger.error(f"[UploadImage] Config load error: {e}")

    if not private_key:
        return "❌ ImageKit Private Key not configured. Ask your admin to configure it in Enterprise Settings → Tools → Upload Image, or set it in your agent's tool config."

    # ── Prepare the file ──
    form_data = {}
    file_content = None

    if file_path:
        # Read from workspace
        full_path = (ws / file_path).resolve()
        if not str(full_path).startswith(str(ws)):
            return "❌ Access denied: path is outside the workspace"
        if not full_path.exists():
            return f"❌ File not found: {file_path}"
        if not full_path.is_file():
            return f"❌ Not a file: {file_path}"

        # Check file size (max 25MB for free plan)
        size_mb = full_path.stat().st_size / (1024 * 1024)
        if size_mb > 25:
            return f"❌ File too large ({size_mb:.1f}MB). Maximum is 25MB."

        file_content = full_path.read_bytes()
        if not file_name:
            file_name = full_path.name
    elif url:
        # Pass URL directly to ImageKit
        form_data["file"] = url
        if not file_name:
            from urllib.parse import urlparse
            file_name = urlparse(url).path.split("/")[-1] or "image.jpg"

    if not file_name:
        file_name = "image.png"

    form_data["fileName"] = file_name
    form_data["folder"] = folder
    form_data["useUniqueFileName"] = "true"

    # ── Upload to ImageKit V2 ──
    auth_string = base64.b64encode(f"{private_key}:".encode()).decode()

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            if file_content:
                # Binary upload via multipart
                files = {"file": (file_name, file_content)}
                resp = await client.post(
                    "https://upload.imagekit.io/api/v2/files/upload",
                    headers={"Authorization": f"Basic {auth_string}"},
                    data=form_data,
                    files=files,
                )
            else:
                # URL upload via form data
                resp = await client.post(
                    "https://upload.imagekit.io/api/v2/files/upload",
                    headers={"Authorization": f"Basic {auth_string}"},
                    data=form_data,
                )

        if resp.status_code in (200, 201):
            result = resp.json()
            cdn_url = result.get("url", "")
            file_id = result.get("fileId", "")
            size = result.get("size", 0)
            size_str = f"{size / 1024:.1f}KB" if size < 1024 * 1024 else f"{size / (1024 * 1024):.1f}MB"
            return (
                f"✅ Image uploaded successfully!\n\n"
                f"**CDN URL**: {cdn_url}\n"
                f"**File ID**: {file_id}\n"
                f"**Size**: {size_str}\n"
                f"**Name**: {result.get('name', file_name)}"
            )
        else:
            error_detail = resp.text[:300]
            return f"❌ Upload failed (HTTP {resp.status_code}): {error_detail}"

    except httpx.TimeoutException:
        return "❌ Upload timed out after 60s. The file may be too large or the network is slow."
    except Exception as e:
        return f"❌ Upload error: {type(e).__name__}: {str(e)[:300]}"



# ─── Image Generation (Multi-Provider) ────────────────────────────────────────

async def _generate_image(agent_id: uuid.UUID, ws: Path, arguments: dict, provider: str) -> str:
    """Generate an image using the configured provider and save to workspace.

    Supported providers:
    - siliconflow: OpenAI-compatible API (FLUX models, China-friendly)
    - openai: Native OpenAI API (GPT Image)
    - google: Google Gemini Native Image API (Nano Banana)

    The tool config is resolved via the standard _get_tool_config() hierarchy:
    global tool config (admin-set) -> per-agent tool config override.
    """
    import httpx
    from datetime import datetime

    prompt = arguments.get("prompt")
    if not prompt:
        return "❌ Missing required argument 'prompt' for generate_image"

    size = arguments.get("size", "1024x1024")
    save_path = arguments.get("save_path", "")

    # Load tool config (global -> per-agent override)
    tool_key = f"generate_image_{provider}"
    config = await _get_tool_config(agent_id, tool_key) or {}
    model = config.get("model", "")
    api_key = config.get("api_key", "")
    base_url = config.get("base_url", "")

    if not api_key:
        return (
            "❌ Image generation API key not configured. "
            "Ask your admin to configure it in Enterprise Settings → Tools → Generate Image."
        )

    # Generate the save path if not provided
    if not save_path:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        # Derive a short slug from the prompt for a more descriptive filename
        slug = "_".join(prompt.split()[:4]).lower()
        slug = "".join(c for c in slug if c.isalnum() or c == "_")[:40]
        save_path = f"workspace/images/{slug}_{ts}.png"

    # Ensure the target directory exists and path is within workspace
    full_save_path = (ws / save_path).resolve()
    if not str(full_save_path).startswith(str(ws.resolve())):
        return "❌ Access denied: save path is outside the workspace"
    full_save_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        if provider == "siliconflow":
            image_bytes = await _generate_image_siliconflow(
                api_key,
                model or "black-forest-labs/FLUX.1-schnell",
                base_url or "https://api.siliconflow.cn/v1",
                prompt, size,
            )
        elif provider == "openai":
            image_bytes = await _generate_image_openai(
                api_key,
                model or "gpt-image-1",
                base_url or "https://api.openai.com/v1",
                prompt, size,
            )
        elif provider == "google":
            image_bytes = await _generate_image_google(
                api_key,
                model or "gemini-2.5-flash-image",
                base_url or "https://generativelanguage.googleapis.com/v1beta",
                prompt, size,
            )
        else:
            return f"❌ Unknown image generation provider: {provider}. Supported: siliconflow, openai, google"

        if not image_bytes:
            return "❌ Image generation returned empty result. Please try a different prompt."

        # Save the generated image to workspace
        full_save_path.write_bytes(image_bytes)
        size_kb = len(image_bytes) / 1024

        # Build the API path for inline display in chat
        # The MarkdownRenderer will auto-inject the auth token for /api/agents/ paths
        api_image_path = f"/api/agents/{agent_id}/files/download?path={save_path}"

        return (
            f"✅ Image generated and saved to: {save_path}\n"
            f"Size: {size_kb:.1f} KB | Provider: {provider} | Model: {model or '(default)'}\n\n"
            f"Display this image to the user using this exact markdown:\n"
            f"![generated image]({api_image_path})"
        )
    except httpx.TimeoutException:
        logger.error(f"[GenerateImage] Timeout ({provider}): took longer than 120 seconds or network unreachable.")
        return (
            f"❌ Image generation failed ({provider}): API request timed out after 120 seconds. "
            f"This is usually caused by network issues or the model taking too long to generate."
        )
    except Exception as e:
        err_msg = str(e) or type(e).__name__
        logger.error(f"[GenerateImage] Error ({provider}): {err_msg}")
        return f"❌ Image generation failed ({provider}): {err_msg[:400]}"


async def _generate_image_siliconflow(
    api_key: str, model: str, base_url: str, prompt: str, size: str
) -> bytes:
    """Generate image via SiliconFlow (OpenAI-compatible images.generate API).

    SiliconFlow returns a temporary URL (expires in ~1 hour), so we download
    the image bytes immediately after generation.
    """
    import httpx
    import base64

    url = f"{base_url.rstrip('/')}/images/generations"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "prompt": prompt,
        "image_size": size,  # SiliconFlow uses 'image_size' instead of 'size'
        "n": 1,
    }

    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(url, json=payload, headers=headers)
        if resp.status_code != 200:
            # Extract API error message for better diagnostics
            try:
                err_body = resp.json()
                err_msg = err_body.get("message") or err_body.get("error", {}).get("message", resp.text[:300])
            except Exception:
                err_msg = resp.text[:300]
            raise ValueError(f"SiliconFlow API error ({resp.status_code}): {err_msg}")
        data = resp.json()

        # SiliconFlow may return url or b64_json
        image_data = data.get("data", [{}])[0]
        image_url = image_data.get("url")
        if image_url:
            # Download the temporary URL immediately
            img_resp = await client.get(image_url, timeout=60)
            img_resp.raise_for_status()
            return img_resp.content

        b64 = image_data.get("b64_json")
        if b64:
            return base64.b64decode(b64)

        raise ValueError(f"No image URL or b64_json in SiliconFlow response: {data}")


async def _generate_image_openai(
    api_key: str, model: str, base_url: str, prompt: str, size: str
) -> bytes:
    """Generate image via OpenAI GPT Image API.

    Requests b64_json format to avoid dealing with URL expiry.
    """
    import httpx
    import base64

    url = f"{base_url.rstrip('/')}/images/generations"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "prompt": prompt,
        "size": size,
        "n": 1,
        "response_format": "b64_json",
    }

    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(url, json=payload, headers=headers)
        if resp.status_code != 200:
            try:
                err_body = resp.json()
                err_msg = err_body.get("error", {}).get("message", resp.text[:300])
            except Exception:
                err_msg = resp.text[:300]
            raise ValueError(f"OpenAI API error ({resp.status_code}): {err_msg}")
        data = resp.json()

        image_data = data.get("data", [{}])[0]
        b64 = image_data.get("b64_json")
        if b64:
            return base64.b64decode(b64)

        # Fallback: try URL
        image_url = image_data.get("url")
        if image_url:
            img_resp = await client.get(image_url, timeout=60)
            img_resp.raise_for_status()
            return img_resp.content

        raise ValueError(f"No b64_json or URL in OpenAI response: {data}")


async def _generate_image_google(
    api_key: str, model: str, base_url: str, prompt: str, size: str
) -> bytes:
    """Generate image via Google Gemini Native Image API (Nano Banana) or Vertex AI.

    Uses the Gemini generateContent endpoint with responseModalities=["IMAGE"].
    Converts WxH size to aspect ratio format (e.g. 1024x1024 -> 1:1).
    Extracts the generated image from inlineData in the response parts.
    """
    import httpx
    import base64

    url = f"{base_url.rstrip('/')}/models/{model}:generateContent?key={api_key}"

    # Convert WxH size to aspect ratio for Gemini API
    # Supported: 1:1, 3:4, 4:3, 9:16, 16:9
    size_to_ratio = {
        "1024x1024": "1:1",
        "768x1024": "3:4",
        "1024x768": "4:3",
        "768x1366": "9:16",
        "1366x768": "16:9",
        "1024x1536": "3:4",
        "1536x1024": "4:3",
    }
    aspect_ratio = size_to_ratio.get(size, "1:1")

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "responseModalities": ["IMAGE"],
            "imageConfig": {
                "numberOfImages": 1,
                "aspectRatio": aspect_ratio,
            },
        },
    }

    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(
            url, json=payload, headers={"Content-Type": "application/json"}
        )
        if resp.status_code != 200:
            try:
                err_body = resp.json()
                err_msg = err_body.get("error", {}).get("message", resp.text[:300])
            except Exception:
                err_msg = resp.text[:300]
            raise ValueError(f"Google Gemini API error ({resp.status_code}): {err_msg}")
        data = resp.json()

        # Extract image from response candidates -> content -> parts
        candidates = data.get("candidates", [])
        if not candidates:
            raise ValueError(f"No candidates in Gemini response: {data}")

        parts = candidates[0].get("content", {}).get("parts", [])
        for part in parts:
            if "inlineData" in part:
                b64 = part["inlineData"]["data"]
                return base64.b64decode(b64)

        raise ValueError(
            f"No image (inlineData) found in Gemini response parts. "
            f"Parts: {[p.get('text', '(image)') if 'text' in p else '(inline)' for p in parts]}"
        )


# ─── Feishu Helper ────────────────────────────────────────────────────────────

async def _get_feishu_token(agent_id: uuid.UUID) -> tuple[str, str] | None:
    """Get (app_id, app_access_token) for the agent's configured Feishu channel."""
    import httpx
    from app.models.channel_config import ChannelConfig

    async with async_session() as db:
        result = await db.execute(
            select(ChannelConfig).where(
                ChannelConfig.agent_id == agent_id,
                ChannelConfig.channel_type == "feishu",
                ChannelConfig.is_configured == True,
            )
        )
        config = result.scalar_one_or_none()

    if not config or not config.app_id or not config.app_secret:
        return None

    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal",
            json={"app_id": config.app_id, "app_secret": config.app_secret},
        )
        token = resp.json().get("tenant_access_token", "")

    return (config.app_id, token) if token else None


async def _get_agent_calendar_id(token: str) -> tuple[str | None, str | None]:
    """Get (calendar_id, error_msg) for the agent app's primary calendar.

    Returns (calendar_id, None) on success, or (None, human_readable_error) on failure.
    """
    import httpx
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/calendar/v4/calendars/primary",
            headers={"Authorization": f"Bearer {token}"},
        )
    data = resp.json()
    code = data.get("code", -1)
    if code == 0:
        cals = data.get("data", {}).get("calendars", [])
        if cals:
            cal_id = cals[0].get("calendar", {}).get("calendar_id")
            return cal_id, None
        return None, "日历列表为空，请确认应用有 calendar:calendar 权限并已发布新版本"
    if code == 99991672:
        return None, (
            "❌ 飞书日历权限未开通（错误码 99991672）\n\n"
            "请在飞书开放平台为应用 cli_a9257c5136781ceb 开通以下权限并发布新版本：\n"
            "• calendar:calendar:readonly（应用身份权限）\n"
            "• calendar:calendar.event:create（应用身份权限）\n"
            "• calendar:calendar.event:read（用户身份权限）\n"
            "• calendar:calendar.event:update（用户身份权限）\n"
            "• calendar:calendar.event:delete（用户身份权限）\n\n"
            "开通步骤：飞书开放平台 → 权限管理 → 批量导入权限 → 添加以上权限 → 创建版本 → 确认发布"
        )
    return None, f"获取日历 ID 失败：{data.get('msg')} (code {code})"


async def _feishu_resolve_open_id(token: str, email: str) -> str | None:
    """Resolve a user's open_id from their email."""
    import httpx
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/contact/v3/users/batch_get_id",
            json={"emails": [email]},
            headers={"Authorization": f"Bearer {token}"},
            params={"user_id_type": "open_id"},
        )
    data = resp.json()
    if data.get("code") != 0:
        return None
    for u in data.get("data", {}).get("user_list", []):
        oid = u.get("user_id")
        if oid:
            return oid
    return None


def _iso_to_ts(iso_str: str) -> float:
    """Convert ISO 8601 string to Unix timestamp."""
    from datetime import datetime as _dt
    for fmt in ("%Y-%m-%dT%H:%M:%S%z", "%Y-%m-%dT%H:%M:%SZ", "%Y-%m-%dT%H:%M:%S"):
        try:
            if iso_str.endswith("Z"):
                d = _dt.fromisoformat(iso_str.replace("Z", "+00:00"))
            else:
                d = _dt.strptime(iso_str, fmt)
            return d.timestamp()
        except ValueError:
            continue
    raise ValueError(f"Cannot parse datetime: {iso_str!r}")


async def _get_feishu_credentials(agent_id: uuid.UUID) -> tuple[str, str]:
    """Retrieve Feishu app_id and app_secret for an agent.
    1. Try Agent-specific ChannelConfig
    2. Fallback to global settings (.env)
    """
    from app.models.channel_config import ChannelConfig
    from app.config import get_settings
    
    settings = get_settings()
    app_id = settings.FEISHU_APP_ID
    app_secret = settings.FEISHU_APP_SECRET
    
    try:
        async with async_session() as db:
            result = await db.execute(
                select(ChannelConfig).where(ChannelConfig.agent_id == agent_id, ChannelConfig.channel_type == "feishu")
            )
            config = result.scalar_one_or_none()
            if config and config.app_id and config.app_secret:
                app_id = config.app_id
                app_secret = config.app_secret
    except Exception:
        pass
        
    return app_id, app_secret


async def _get_feishu_tenant_doc_url(tenant_token: str, doc_token: str, doc_type: str = "docx") -> str:
    """Build a user-accessible document URL using the tenant's actual domain.

    The API gateway (open.feishu.cn) cannot serve user documents - we must use
    the tenant's own domain (e.g. xxx.feishu.cn or xxx.larksuite.com).
    Falls back to generating a search link if the tenant domain cannot be resolved.

    Args:
        tenant_token: A valid tenant_access_token.
        doc_token:    The document_id (docx) or wiki node token.
        doc_type:     'docx' or 'wiki' - controls the URL path prefix.
    Returns:
        A fully-formed URL string.
    """
    import httpx
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(
                "https://open.feishu.cn/open-apis/tenant/v2/tenant/query",
                headers={"Authorization": f"Bearer {tenant_token}"},
            )
        data = resp.json()
        if data.get("code") == 0:
            domain = data.get("data", {}).get("tenant", {}).get("domain", "")
            if domain:
                return f"https://{domain}/{doc_type}/{doc_token}"
    except Exception:
        pass
    # Fallback: construct a search URL so the user can locate the document
    return f"https://feishu.cn/{doc_type}/{doc_token}"




async def _get_feishu_bitable_url(tenant_token: str, app_token: str, table_id: str = "") -> str:
    """Build a user-accessible Bitable URL using the tenant's actual domain.

    Constructs https://{tenant_domain}/base/{app_token}?table={table_id}
    Falls back to https://feishu.cn/base/{app_token} if domain resolution fails.

    Args:
        tenant_token: A valid tenant_access_token.
        app_token:    The Bitable app token.
        table_id:     Optional table ID to deep-link to a specific sheet.
    Returns:
        A fully-formed URL string.
    """
    import httpx
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(
                "https://open.feishu.cn/open-apis/tenant/v2/tenant/query",
                headers={"Authorization": f"Bearer {tenant_token}"},
            )
        data = resp.json()
        if data.get("code") == 0:
            domain = data.get("data", {}).get("tenant", {}).get("domain", "")
            if domain:
                base_url = f"https://{domain}/base/{app_token}"
                if table_id:
                    base_url += f"?table={table_id}"
                return base_url
    except Exception:
        pass
    # Fallback
    base_url = f"https://feishu.cn/base/{app_token}"
    if table_id:
        base_url += f"?table={table_id}"
    return base_url


def _parse_feishu_url(url: str) -> dict:
    """Parse various Feishu URLs to extract tokens.
    Supports Bitable (table, view) and Docx.
    """
    import re
    result = {}
    
    # Bitable URL regex: e.g., https://example.feishu.cn/base/{app_token}?table={table_id}&view={view_id}
    base_match = re.search(r'/base/([a-zA-Z0-9_]+)', url)
    if base_match:
        result['app_token'] = base_match.group(1)
        
    table_match = re.search(r'table=([a-zA-Z0-9_]+)', url)
    if table_match:
        result['table_id'] = table_match.group(1)
    
    # support URL with /tblxxxxxx
    if not 'table_id' in result:
        tbl_match = re.search(r'/(tbl[a-zA-Z0-9_]+)', url)
        if tbl_match:
            result['table_id'] = tbl_match.group(1)
            
    view_match = re.search(r'view=([a-zA-Z0-9_]+)', url)
    if view_match:
        result['view_id'] = view_match.group(1)
        
    # Docx URL regex
    docx_match = re.search(r'/docx/([a-zA-Z0-9_]+)', url)
    if docx_match:
        result['document_token'] = docx_match.group(1)
        
    # Wiki URL regex
    wiki_match = re.search(r'/wiki/([a-zA-Z0-9_]+)', url)
    if wiki_match:
        result['wiki_token'] = wiki_match.group(1)
        
    return result


# ─── Feishu Bitable Tools ──────────────────────────────────────────

async def _resolve_bitable_app_token(agent_id: uuid.UUID, parsed_url: dict) -> str | None:
    app_token = parsed_url.get("app_token")
    if app_token:
        return app_token
    wiki_token = parsed_url.get("wiki_token")
    if wiki_token:
        app_id, app_secret = await _get_feishu_credentials(agent_id)
        if app_id and app_secret:
            from app.services.feishu_service import feishu_service
            token = await feishu_service.get_tenant_access_token(app_id, app_secret)
            node_info = await _feishu_wiki_get_node(wiki_token, token)
            if node_info and node_info.get("obj_token"):
                return node_info["obj_token"]
    return None

def _check_feishu_err(resp: dict) -> str | None:
    """Check Feishu API response for errors and return a user-friendly message.

    For permission-related errors, returns detailed step-by-step instructions
    guiding the user through the Feishu UI to grant the bot app access to
    the target document or Bitable.
    """
    code = resp.get("code")
    if code != 0:
        msg = str(resp.get("msg", ""))
        msg_lower = msg.lower()
        # Common Feishu permission/access error codes:
        #   99991663 - no permission to access resource
        #   99991661 - token invalid / access denied
        #   99991668 - permission denied on resource
        #   10006    - no permission
        #   91403    - Forbidden (bitable write without app access)
        #   91404    - bitable record not found (sometimes permission)
        #   1063001  - doc permission denied
        #   1063004  - doc operation forbidden
        _perm_codes = {99991663, 10006, 99991661, 99991668, 91403, 1063001, 1063004}
        _perm_keywords = ("permission", "forbidden", "no access", "access denied", "403")
        is_perm_error = code in _perm_codes or any(kw in msg_lower for kw in _perm_keywords)
        if is_perm_error:
            return (
                f"Failed: Permission denied (code: {code}, msg: {msg}). "
                "The bot app does not have access to this document/Bitable. "
                "CRITICAL INSTRUCTION for the Agent: You MUST explicitly output ALL of the following "
                "step-by-step instructions to the user in your reply. Do NOT omit or summarize any step:\n\n"
                "---\n"
                "**The bot needs to be added to this document/table as an authorized app. "
                "Please follow these steps:**\n\n"
                "1. Open the target document or Bitable in the Feishu client (web or desktop).\n"
                "2. Click the **「...」** menu button in the top-right corner of the page.\n"
                "3. In the dropdown menu, hover over **「更多」** (More) at the bottom.\n"
                "4. In the sub-menu that appears, click **「添加文档应用」** (Add Document App).\n"
                "5. In the search box, type the name of your Feishu bot app (the one bound to this Agent's channel), then click to add it.\n"
                "6. After adding, retry the same operation.\n\n"
                "If you cannot find 「添加文档应用」, it means the document owner may need to enable this option, "
                "or you can try: click **「分享」** (Share) button -> invite the bot app directly.\n"
                "---"
            )
        return f"Failed: API Error {code} - {msg}"
    return None

async def _bitable_list_tables(agent_id: uuid.UUID, arguments: dict) -> str:
    """List all tables in a Feishu Bitable app."""
    url = arguments.get("url", "")
    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    if not app_token:
        return "Failed: Could not extract Bitable app_token from the URL (also could not resolve wiki_token)."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."
        
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_list_tables(app_id, app_secret, app_token)
        err = _check_feishu_err(resp)
        if err: return err
        
        tables = resp.get("data", {}).get("items", [])
        if not tables:
            return "OK: No tables found in this Bitable."
        lines = [f"- {t.get('name')} (ID: {t.get('table_id')})" for t in tables]
        # Provide a user-accessible link so the user can open the Bitable directly
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token)
        return "OK: Tables in this Bitable:\n" + "\n".join(lines) + f"\n\n🔗 多维表格链接: {bitable_url}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _bitable_create_app(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new Feishu Bitable (多维表格) app.

    Calls the Bitable v1 apps API: POST /open-apis/bitable/v1/apps
    The API response includes a user-accessible URL with the tenant's own domain.
    """
    name = arguments.get("name", "").strip()
    if not name:
        return "Failed: Missing required argument 'name' — please provide a name for the new Bitable."

    folder_token = arguments.get("folder_token", "").strip()

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_create_app(app_id, app_secret, name, folder_token)
        err = _check_feishu_err(resp)
        if err:
            return err

        # API response structure: data.app.{app_token, name, url, default_table_id, folder_token}
        app_info = resp.get("data", {}).get("app", {})
        app_token = app_info.get("app_token", "")
        bitable_url = app_info.get("url", "")
        default_table_id = app_info.get("default_table_id", "")
        if not app_token:
            return f"Failed: Bitable created but could not extract app_token from response: {resp}"

        # Fallback URL resolution if the API didn't return one
        if not bitable_url:
            tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
            bitable_url = await _get_feishu_bitable_url(tenant_token, app_token)

        result = (
            f"OK: Bitable created successfully!\n"
            f"Name: {name}\n"
            f"App Token: {app_token}\n"
            f"URL: {bitable_url}"
        )
        if default_table_id:
            result += f"\nDefault Table ID: {default_table_id}"
        return result
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _bitable_list_fields(agent_id: uuid.UUID, arguments: dict) -> str:
    """List all fields (columns) in a specific Bitable table."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    
    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")
    
    if not app_token:
        return "Failed: Could not extract Bitable app_token from the URL."
    if not table_id:
        return "Failed: table_id is required. Provide it as a parameter or include it in the URL."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_list_fields(app_id, app_secret, app_token, table_id)
        err = _check_feishu_err(resp)
        if err: return err
        
        fields = resp.get("data", {}).get("items", [])
        if not fields:
            return "OK: No fields found in this table."
        lines = [f"- {f.get('field_name')} (type: {f.get('type')}, ID: {f.get('field_id')})" for f in fields]
        return "OK: Fields in this table:\n" + "\n".join(lines)
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_query_records(agent_id: uuid.UUID, arguments: dict) -> str:
    """Query records (rows) from a Bitable table, with optional FQL filter."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    filter_info = arguments.get("filter_info", "")
    max_results = arguments.get("max_results", 100)
    
    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")
    
    if not app_token or not table_id:
        return "Failed: Could not resolve app_token or table_id from the provided parameters/URL."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        import json
        filters_dict = {}
        if isinstance(filter_info, dict):
            filters_dict = filter_info
        elif isinstance(filter_info, str) and filter_info.strip():
            try:
                filters_dict = json.loads(filter_info)
            except:
                pass 
                
        resp = await feishu_service.bitable_query_records(app_id, app_secret, app_token, table_id, filters_dict)
        err = _check_feishu_err(resp)
        if err: return err
        
        records = resp.get("data", {}).get("items", [])
        if not records:
            return "OK: No matching records found."
        
        lines = []
        for r in records[:max_results]:
            lines.append(f"Record {r.get('record_id')}: {json.dumps(r.get('fields', {}), ensure_ascii=False)}")
        return "OK: Query results:\n" + "\n".join(lines)
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_create_record(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new record (row) in a Bitable table."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    fields_str = arguments.get("fields", "{}")
    
    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")
    
    if not app_token or not table_id:
        return "Failed: Could not resolve app_token or table_id from the provided parameters/URL."
        
    import json
    try:
        fields = json.loads(fields_str)
    except json.JSONDecodeError:
        return "Failed: The 'fields' parameter is not valid JSON."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_create_record(app_id, app_secret, app_token, table_id, fields)
        err = _check_feishu_err(resp)
        if err: return err
        
        record = resp.get("data", {}).get("record", {})
        # Provide a user-accessible link so they can verify the new row in the table
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token, table_id)
        return (
            f"OK: Record created. Record ID: {record.get('record_id')}\n"
            f"Fields: {json.dumps(record.get('fields', {}), ensure_ascii=False)}\n"
            f"🔗 多维表格链接: {bitable_url}"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_update_record(agent_id: uuid.UUID, arguments: dict) -> str:
    """Update an existing record in a Bitable table by record_id."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    record_id = arguments.get("record_id", "")
    fields_str = arguments.get("fields", "{}")
    
    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")
    
    if not app_token or not table_id or not record_id:
        return "Failed: Missing required parameters. Need app_token (from URL), table_id, and record_id."
        
    import json
    try:
        fields = json.loads(fields_str)
    except json.JSONDecodeError:
        return "Failed: The 'fields' parameter is not valid JSON."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_update_record(app_id, app_secret, app_token, table_id, record_id, fields)
        err = _check_feishu_err(resp)
        if err: return err
        
        record = resp.get("data", {}).get("record", {})
        # Provide a user-accessible link so they can verify the updated row
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token, table_id)
        return (
            f"OK: Record updated. Record ID: {record.get('record_id')}\n"
            f"Fields: {json.dumps(record.get('fields', {}), ensure_ascii=False)}\n"
            f"🔗 多维表格链接: {bitable_url}"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _bitable_delete_record(agent_id: uuid.UUID, arguments: dict) -> str:
    """Delete a record from a Bitable table by record_id."""
    url = arguments.get("url", "")
    table_id = arguments.get("table_id", "")
    record_id = arguments.get("record_id", "")
    
    parsed = _parse_feishu_url(url)
    app_token = await _resolve_bitable_app_token(agent_id, parsed)
    table_id = table_id or parsed.get("table_id")
    
    if not app_token or not table_id or not record_id:
        return "Failed: Missing required parameters. Need app_token (from URL), table_id, and record_id."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.bitable_delete_record(app_id, app_secret, app_token, table_id, record_id)
        err = _check_feishu_err(resp)
        if err: return err
        
        # Provide a user-accessible link so they can verify the deletion
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        bitable_url = await _get_feishu_bitable_url(tenant_token, app_token, table_id)
        return f"OK: Record {record_id} deleted successfully.\n🔗 多维表格链接: {bitable_url}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


# ─── Feishu Document Tools ──────────────────────────────────────────

async def _resolve_docx_document_token(agent_id: uuid.UUID, parsed_url: dict) -> str | None:
    doc_token = parsed_url.get("document_token")
    if doc_token:
        return doc_token
    wiki_token = parsed_url.get("wiki_token")
    if wiki_token:
        app_id, app_secret = await _get_feishu_credentials(agent_id)
        if app_id and app_secret:
            from app.services.feishu_service import feishu_service
            token = await feishu_service.get_tenant_access_token(app_id, app_secret)
            node_info = await _feishu_wiki_get_node(wiki_token, token)
            if node_info and node_info.get("obj_token"):
                return node_info["obj_token"]
    return None

async def _feishu_read_doc(agent_id: uuid.UUID, arguments: dict) -> str:
    """Read full text content of a Feishu Docx."""
    url = arguments.get("url", "")
    parsed = _parse_feishu_url(url)
    doc_token = await _resolve_docx_document_token(agent_id, parsed)
    if not doc_token:
        return "Failed: Could not extract Document token from the URL."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."
        
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.read_feishu_doc(app_id, app_secret, doc_token)
        err = _check_feishu_err(resp)
        if err: return err
        
        content = resp.get("data", {}).get("content", "")
        if not content:
            return "OK: Document is empty or content is unavailable."
        return f"OK: Document Content:\n{content}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _feishu_create_doc(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new blank Feishu Docx."""
    title = arguments.get("title", "Untitled Document")
    folder_token = arguments.get("folder_token", "")
    
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."
        
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.create_feishu_doc(app_id, app_secret, folder_token or None, title)
        err = _check_feishu_err(resp)
        if err: return err
        
        doc = resp.get("data", {}).get("document", {})
        doc_id = doc.get("document_id")
        # Get the tenant's actual domain (open.feishu.cn is the API gateway, not for users)
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        url = await _get_feishu_tenant_doc_url(tenant_token, doc_id)
        return f"OK: Document created perfectly. Document ID: {doc_id}\nURL: {url}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

async def _feishu_append_doc(agent_id: uuid.UUID, arguments: dict) -> str:
    """Append text to the bottom of a Feishu Docx."""
    url = arguments.get("url", "")
    content = arguments.get("content", "")
    if not content:
        return "Failed: Content to append cannot be empty."
        
    parsed = _parse_feishu_url(url)
    doc_token = await _resolve_docx_document_token(agent_id, parsed)
    if not doc_token:
        return "Failed: Could not extract Document token from the URL."
        
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."
        
    from app.services.feishu_service import feishu_service
    try:
        # Feishu uses the document_id as the root block_id to append entirely to the document
        resp = await feishu_service.append_feishu_doc(app_id, app_secret, doc_token, content)
        err = _check_feishu_err(resp)
        if err: return err
        
        return "OK: Content appended successfully to the end of the document."
    except Exception as e:
        return f"Failed: {str(e)[:300]}"

# ─── Feishu Wiki Tools ───────────────────────────────────────────────────────

async def _feishu_wiki_get_node(token_str: str, auth_token: str) -> dict | None:
    """Call wiki get_node API to resolve a wiki node token → {obj_token, space_id, has_child, title}.
    Returns None if the token is not a wiki node."""
    import httpx
    async with httpx.AsyncClient(timeout=5) as client:
        r = await client.get(
            "https://open.feishu.cn/open-apis/wiki/v2/spaces/get_node",
            headers={"Authorization": f"Bearer {auth_token}"},
            params={"token": token_str, "obj_type": "wiki"},
        )
    d = r.json()
    if d.get("code") != 0:
        return None
    node = d.get("data", {}).get("node", {})
    return {
        "obj_token": node.get("obj_token", ""),
        "space_id": node.get("origin_space_id", node.get("space_id", "")),
        "has_child": node.get("has_child", False),
        "title": node.get("title", ""),
        "node_token": node.get("node_token", token_str),
    }


async def _feishu_wiki_list(agent_id: uuid.UUID, arguments: dict) -> str:
    """List sub-pages of a Feishu Wiki node, optionally recursive."""
    import httpx

    node_token = (arguments.get("node_token") or "").strip()
    recursive = bool(arguments.get("recursive", False))

    if not node_token:
        return "❌ Missing required argument 'node_token'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)
    headers = {"Authorization": f"Bearer {token}"}

    # Resolve node → space_id
    node_info = await _feishu_wiki_get_node(node_token, token)
    if not node_info:
        return (
            f"❌ 无法解析 Wiki 节点 `{node_token}`。\n"
            "请确认 token 来自飞书知识库 URL（https://xxx.feishu.cn/wiki/NodeToken），"
            "而非普通文档 URL。"
        )

    space_id = node_info["space_id"]
    if not space_id:
        return f"❌ 无法获取知识库 space_id，请检查 token 是否正确。"

    async def _list_children(parent_token: str, depth: int) -> list[dict]:
        """Return flat list of {title, node_token, obj_token, has_child, depth}."""
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/nodes",
                headers=headers,
                params={"parent_node_token": parent_token, "page_size": 50},
            )
        data = resp.json()
        if data.get("code") != 0:
            return []
        items = data.get("data", {}).get("items", [])
        result = []
        for item in items:
            entry = {
                "title": item.get("title", "(无标题)"),
                "node_token": item.get("node_token", ""),
                "obj_token": item.get("obj_token", ""),
                "has_child": item.get("has_child", False),
                "depth": depth,
            }
            result.append(entry)
            if recursive and entry["has_child"] and depth < 2:
                children = await _list_children(entry["node_token"], depth + 1)
                result.extend(children)
        return result

    pages = await _list_children(node_token, 0)
    if not pages:
        return f"📂 Wiki 页面 `{node_token}` 下没有子页面。"

    lines = [f"📂 Wiki 页面 `{node_token}` 的子页面（共 {len(pages)} 个）：\n"]
    for p in pages:
        indent = "  " * p["depth"]
        child_hint = " _(有子页面)_" if p["has_child"] else ""
        lines.append(
            f"{indent}• **{p['title']}**{child_hint}\n"
            f"{indent}  node_token: `{p['node_token']}`\n"
            f"{indent}  obj_token: `{p['obj_token']}`"
        )
    lines.append(
        "\n💡 用 `feishu_doc_read(document_token=\"<node_token>\")` 读取每个子页面的内容。"
        "\n   对有子页面的条目，再次调用 `feishu_wiki_list(node_token=\"...\")` 继续展开。"
    )
    return "\n".join(lines)


async def _feishu_doc_read(agent_id: uuid.UUID, arguments: dict) -> str:
    document_token = arguments.get("document_token", "").strip()
    if not document_token:
        url = arguments.get("url", "")
        parsed = _parse_feishu_url(url)
        document_token = parsed.get("document_token", parsed.get("wiki_token", ""))
        
    if not document_token:
        return "Failed: Missing required argument 'document_token'"
    max_chars = min(int(arguments.get("max_chars", 6000)), 20000)

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
    
    read_token = document_token
    wiki_hint = ""
    node_info = await _feishu_wiki_get_node(document_token, tenant_token)
    if node_info and node_info.get("obj_token"):
        read_token = node_info["obj_token"]
        if node_info.get("has_child"):
            wiki_hint = (
                "\n\n> 💡 这是一个 Wiki 目录页，它有多个子页面。"
                "使用 `feishu_wiki_list` 工具（传入相同的 node_token）可以查看所有子页面列表。"
            )

    try:
        resp = await feishu_service.read_feishu_doc(app_id, app_secret, read_token)
        err = _check_feishu_err(resp)
        if err: return err
        
        content = resp.get("data", {}).get("content", "")
        if not content:
            return f"📄 Document '{document_token}' is empty.{wiki_hint}"

        truncated = ""
        if len(content) > max_chars:
            content = content[:max_chars]
            truncated = f"\n\n_(Truncated to {max_chars} chars)_"

        return f"📄 **Document content** (`{document_token}`):\n\n{content}{truncated}{wiki_hint}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _feishu_doc_create(agent_id: uuid.UUID, arguments: dict) -> str:
    title = arguments.get("title", "").strip()
    if not title:
        return "Failed: Missing required argument 'title'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."
        
    folder_token = arguments.get("folder_token")
    
    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.create_feishu_doc(app_id, app_secret, folder_token, title)
        err = _check_feishu_err(resp)
        if err: return err
        
        doc = resp.get("data", {}).get("document", {})
        doc_token = doc.get("document_id", "")
        # Get tenant-specific doc URL via tenant domain resolution
        tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
        doc_url = await _get_feishu_tenant_doc_url(tenant_token, doc_token)
        
        # Auto-share with the Feishu sender so they can access the document
        share_note = ""
        try:
            from app.api.websocket_chat import channel_feishu_sender_open_id
            sender_open_id = channel_feishu_sender_open_id.get(None)
            if sender_open_id and doc_token:
                import httpx
                tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)
                async with httpx.AsyncClient(timeout=10) as client:
                    share_resp = await client.post(
                        f"https://open.feishu.cn/open-apis/drive/v1/permissions/{doc_token}/members",
                        params={"type": "docx"},
                        json={
                            "member_type": "openid",
                            "member_id": sender_open_id,
                            "perm": "full_access",
                        },
                        headers={"Authorization": f"Bearer {tenant_token}"},
                    )
                sr = share_resp.json()
                if sr.get("code") == 0:
                    share_note = "\n✅ 已自动为你开通访问权限。"
                else:
                    share_note = f"\n⚠️ 自动授权失败（{sr.get('code')}），你可能需要手动在飞书前端搜索此文件。"
        except Exception as _e:
            share_note = f"\n⚠️ 自动授权异常: {_e}"

        return (
            f"✅ 文档创建成功！{share_note}\n"
            f"标题：{title}\n"
            f"Token：{doc_token}\n"
            f"🔗 访问链接：{doc_url}\n"
            f"下一步：调用 feishu_doc_append(document_token=\"{doc_token}\", content=\"...\") 写入正文内容。"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


def _parse_inline_markdown(text: str) -> list[dict]:
    """Parse inline markdown (bold, italic, strikethrough) into Feishu text_run elements.
    Note: inline `code` is deliberately NOT rendered as inline_code style because
    Feishu's API rejects inline_code inside heading blocks (field validation error).
    Instead, backtick-wrapped text is returned as plain text.
    Empty text_element_style dicts are intentionally omitted to avoid API validation errors.
    """
    import re as _re

    def _make_run(content: str, style: dict | None = None) -> dict:
        run: dict = {"content": content}
        if style:
            run["text_element_style"] = style
        return {"text_run": run}

    elements = []
    # Only handle **bold**, *italic*, ~~strikethrough~~; backticks become plain text
    pattern = r'(\*\*(.+?)\*\*|\*(.+?)\*|~~(.+?)~~|`(.+?)`)'
    pos = 0
    for m in _re.finditer(pattern, text):
        if m.start() > pos:
            elements.append(_make_run(text[pos:m.start()]))
        raw = m.group(0)
        if raw.startswith("**"):
            elements.append(_make_run(m.group(2), {"bold": True}))
        elif raw.startswith("~~"):
            elements.append(_make_run(m.group(4), {"strikethrough": True}))
        elif raw.startswith("`"):
            # Render as plain text to avoid inline_code validation issues in headings
            elements.append(_make_run(m.group(5)))
        else:
            elements.append(_make_run(m.group(3), {"italic": True}))
        pos = m.end()
    if pos < len(text):
        elements.append(_make_run(text[pos:]))
    if not elements:
        elements.append(_make_run(text or " "))
    return elements


def _markdown_to_feishu_blocks(markdown: str) -> list[dict]:
    """Convert Markdown text to Feishu docx v1 block list.

    Supported:
      # / ## / ### / ####  → heading1-4 (block_type 3-6)
      - / * / + text       → bullet      (block_type 12)
      1. text              → ordered     (block_type 13)
      > text               → quote       (block_type 15)
      --- / ***            → divider     (block_type 22)
      ``` ... ```          → code block  (block_type 14)
      plain text           → text        (block_type 2)
      inline **bold** *italic* `code` ~~strike~~  → text_element_style
    """
    import re as _re

    _HEADING_BLOCK = {1: (3, "heading1"), 2: (4, "heading2"),
                      3: (5, "heading3"), 4: (6, "heading4")}

    def _text_block(bt: int, key: str, line: str) -> dict:
        # Omit "style" entirely to avoid Feishu field validation errors on empty style dicts
        return {
            "block_type": bt,
            key: {"elements": _parse_inline_markdown(line)},
        }

    blocks: list[dict] = []
    lines = markdown.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]

        # ── Code fence ──────────────────────────────────────────────────────
        if line.strip().startswith("```"):
            lang = line.strip()[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            blocks.append({
                "block_type": 14,
                "code": {
                    "elements": [{"text_run": {"content": "\n".join(code_lines)}}],
                    "style": {"language": 1 if not lang else
                              {"python": 49, "javascript": 22, "js": 22,
                               "typescript": 56, "ts": 56, "bash": 4, "sh": 4,
                               "sql": 53, "java": 21, "go": 17, "rust": 51,
                               "json": 25, "yaml": 60, "html": 19, "css": 10,
                               }.get(lang.lower(), 1)},
                },
            })
            i += 1
            continue

        # ── Divider ──────────────────────────────────────────────────────────
        if _re.fullmatch(r'[-*_]{3,}', line.strip()):
            # block_type 22 = Divider; no extra fields allowed (empty dict causes validation error)
            blocks.append({"block_type": 22})
            i += 1
            continue

        # ── Headings ─────────────────────────────────────────────────────────
        hm = _re.match(r'^(#{1,4})\s+(.*)', line)
        if hm:
            level = min(len(hm.group(1)), 4)
            bt, key = _HEADING_BLOCK[level]
            blocks.append(_text_block(bt, key, hm.group(2)))
            i += 1
            continue

        # ── Bullet list ──────────────────────────────────────────────────────
        if _re.match(r'^[\-\*\+]\s+', line):
            text = _re.sub(r'^[\-\*\+]\s+', '', line)
            blocks.append(_text_block(12, "bullet", text))
            i += 1
            continue

        # ── Ordered list ─────────────────────────────────────────────────────
        if _re.match(r'^\d+\.\s+', line):
            text = _re.sub(r'^\d+\.\s+', '', line)
            blocks.append(_text_block(13, "ordered", text))
            i += 1
            continue

        # ── Blockquote ───────────────────────────────────────────────────────
        if line.startswith("> "):
            blocks.append(_text_block(15, "quote", line[2:]))
            i += 1
            continue

        # ── Empty line → empty text block ────────────────────────────────────
        if line.strip() == "":
            blocks.append({
                "block_type": 2,
                "text": {"elements": [{"text_run": {"content": " "}}]},
            })
            i += 1
            continue

        # ── Markdown table separator line (|---|---| ) → skip ───────────────
        if _re.match(r'^\|[\s\-:]+(\|[\s\-:]+)*\|?\s*$', line.strip()):
            i += 1
            continue

        # ── Markdown table row → plain text ──────────────────────────────────
        if line.strip().startswith("|") and line.strip().endswith("|"):
            # Strip pipe separators and render each cell as plain text
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            cell_text = "  |  ".join(c for c in cells if c)
            blocks.append(_text_block(2, "text", cell_text))
            i += 1
            continue

        # ── Plain text (with inline formatting) ──────────────────────────────
        blocks.append(_text_block(2, "text", line))
        i += 1

    return blocks


async def _feishu_doc_append(agent_id: uuid.UUID, arguments: dict) -> str:
    document_token = arguments.get("document_token", "").strip()
    if not document_token:
        url = arguments.get("url", "")
        parsed = _parse_feishu_url(url)
        document_token = parsed.get("document_token", parsed.get("wiki_token", ""))
        
    content = arguments.get("content", "").strip()
    if not document_token:
        return "Failed: Missing required argument 'document_token'"
    if not content:
        return "Failed: Missing required argument 'content'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "Failed: Feishu app credentials not configured for this agent."

    from app.services.feishu_service import feishu_service
    tenant_token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    # For wiki node tokens, use the obj_token for the docx API
    node_info = await _feishu_wiki_get_node(document_token, tenant_token)
    docx_token = node_info["obj_token"] if (node_info and node_info.get("obj_token")) else document_token

    try:
        import httpx
        async with httpx.AsyncClient(timeout=20) as client:
            meta_resp = (await client.get(
                f"https://open.feishu.cn/open-apis/docx/v1/documents/{docx_token}",
                headers={"Authorization": f"Bearer {tenant_token}"},
            )).json()
            err = _check_feishu_err(meta_resp)
            if err: return err

            body_block_id = (
                meta_resp.get("data", {}).get("document", {}).get("body", {}).get("block_id")
                or docx_token
            )

            children = _markdown_to_feishu_blocks(content)

            result = (await client.post(
                f"https://open.feishu.cn/open-apis/docx/v1/documents/{docx_token}/blocks/{body_block_id}/children",
                json={"children": children, "index": -1}, # -1 appends to end
                headers={"Authorization": f"Bearer {tenant_token}"},
            )).json()

            err = _check_feishu_err(result)
            if err: return err

        doc_url = await _get_feishu_tenant_doc_url(tenant_token, docx_token)
        return (
            f"✅ 已写入 {len(children)} 个段落到文档。\n"
            f"🔗 文档直链（原文发给用户，勿修改）：{doc_url}"
        )
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


# ─── Feishu Drive Share (All File Types) ────────────────────────────────────────

async def _feishu_drive_share(agent_id: uuid.UUID, arguments: dict) -> str:
    """Manage Feishu drive file collaborators.
    Automatically handles both regular docs/files (Drive permissions API)
    and Wiki node documents (Wiki space members API).
    """
    import httpx
    import re as _re

    document_token = (arguments.get("document_token") or "").strip()
    doc_type = (arguments.get("doc_type") or "docx").strip()
    action = (arguments.get("action") or "list").strip()
    permission = (arguments.get("permission") or "edit").strip()

    if not document_token:
        return "❌ Missing required argument 'document_token'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)
    headers = {"Authorization": f"Bearer {token}"}

    # ── Detect if this is a Wiki node token ─────────────────────────────────
    node_info = await _feishu_wiki_get_node(document_token, token)
    is_wiki = node_info is not None
    space_id = node_info.get("space_id", "") if node_info else ""
    obj_token = node_info.get("obj_token", "") if node_info else ""

    # Permission level mapping: Feishu API uses "view" / "edit" / "full_access"
    api_perm = {"view": "view", "edit": "edit", "full_access": "full_access"}.get(permission, "edit")
    # Wiki space role mapping: only "admin" / "member" are valid roles
    wiki_role = "admin" if api_perm in ("edit", "full_access") else "member"

    # ── LIST collaborators ────────────────────────────────────────────────────
    if action == "list":
        use_token = obj_token if (is_wiki and obj_token) else document_token
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://open.feishu.cn/open-apis/drive/v1/permissions/{use_token}/members",
                params={"type": doc_type},
                headers=headers,
            )
        data = resp.json()
        if data.get("code") != 0:
            _c = data.get("code")
            if _c == 1063003 and is_wiki:
                return (
                    f"ℹ️ 文档 `{document_token}` 是知识库页面，其权限由知识库空间统一管理。\n"
                    "知识库空间 ID：`" + space_id + "`\n"
                    "请直接在飞书知识库中管理成员权限。"
                )
            if _c in (99991672, 99991668):
                return (
                    f"❌ 权限不足（code {_c}）\n"
                    "需要在飞书开放平台开通：\n"
                    "• drive:drive（云文档权限管理）"
                )
            return f"❌ 获取协作者列表失败：{data.get('msg')} (code {_c})"

        members = data.get("data", {}).get("items", [])
        if not members:
            return f"📄 文档 `{document_token}` 当前没有其他协作者。"

        lines = [f"📄 文档 `{document_token}` 的协作者列表（共 {len(members)} 人）：\n"]
        for m in members:
            perm = m.get("perm", "")
            member_type = m.get("member_type", "")
            member_id = m.get("member_id", "")
            _type_label = {"openid": "用户", "openchat": "群组", "opendepartmentid": "部门"}.get(member_type, member_type)
            lines.append(f"• {_type_label} `{member_id}` | 权限: **{perm}**")
        return "\n".join(lines)

    # ── ADD / REMOVE collaborators ─────────────────────────────────────────────
    member_names: list[str] = list(arguments.get("member_names") or [])
    member_open_ids: list[str] = list(arguments.get("member_open_ids") or [])

    if not member_names and not member_open_ids:
        return "❌ 请提供 member_names（姓名列表）或 member_open_ids（open_id 列表）"

    # Resolve names → open_ids
    resolved: list[tuple[str, str]] = []  # (display_name, open_id)
    for name in member_names:
        sr = await _feishu_user_search(agent_id, {"name": name})
        m = _re.search(r'open_id: `(ou_[A-Za-z0-9]+)`', sr)
        if m:
            resolved.append((name, m.group(1)))
        else:
            resolved.append((name, ""))

    for oid in member_open_ids:
        if oid:
            resolved.append((oid, oid))

    results = []
    async with httpx.AsyncClient(timeout=15) as client:
        for display, oid in resolved:
            if not oid:
                results.append(f"❌ 无法找到「{display}」的 open_id，跳过")
                continue

            if action == "add":
                # ── Wiki node: use wiki space members API ──────────────────
                if is_wiki and space_id:
                    resp = await client.post(
                        f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/members",
                        json={"member_type": "openid", "member_id": oid, "member_role": wiki_role},
                        headers=headers,
                    )
                    d = resp.json()
                    _c = d.get("code")
                    if _c == 0:
                        results.append(f"✅ 已将「{display}」加入知识库空间（角色：{wiki_role}）")
                    elif _c == 131008:
                        results.append(f"ℹ️ 「{display}」已经是知识库成员，无需重复添加")
                    elif _c == 131101:
                        # Public wiki space — everyone already has access
                        results.append(
                            f"ℹ️ 这是一个**公开知识库**，所有人已可访问。\n"
                            f"「{display}」无需单独添加权限。"
                        )
                    else:
                        results.append(f"❌ 添加「{display}」到知识库失败：{d.get('msg')} (code {_c})")
                    continue

                # ── Regular docx: use Drive permissions API ────────────────
                body = {
                    "member_type": "openid",
                    "member_id": oid,
                    "perm": api_perm,
                }
                resp = await client.post(
                    f"https://open.feishu.cn/open-apis/drive/v1/permissions/{document_token}/members",
                    json=body,
                    headers=headers,
                    params={"type": doc_type},
                )
                d = resp.json()
                if d.get("code") == 0:
                    results.append(f"✅ 已将「{display}」添加为**{permission}**权限协作者")
                else:
                    _c = d.get("code")
                    if _c == 99992402:
                        # Feishu platform policy: you cannot add yourself as a collaborator via API.
                        # Permissions must be granted by others, or set manually in the UI.
                        results.append(
                            f"⚠️ 飞书平台安全限制：无法通过 API 为自己添加协作权限。\n"
                            f"请手动操作：打开文档 → 右上角「分享」→ 添加自己并设置权限。"
                        )
                    elif _c in (99991672, 99991668):
                        return (
                            f"❌ 权限不足（code {_c}）\n"
                            "需要在飞书开放平台开通：\n"
                            "• drive:drive（云文档权限管理）"
                        )
                    else:
                        results.append(f"❌ 添加「{display}」失败：{d.get('msg')} (code {_c})")

            elif action == "remove":
                if is_wiki and space_id:
                    resp = await client.delete(
                        f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/members/{oid}",
                        headers=headers,
                        params={"member_type": "openid"},
                    )
                    d = resp.json()
                    if d.get("code") == 0:
                        results.append(f"✅ 已将「{display}」从知识库移除")
                    else:
                        results.append(f"❌ 移除「{display}」失败：{d.get('msg')} (code {d.get('code')})")
                    continue

                resp = await client.delete(
                    f"https://open.feishu.cn/open-apis/drive/v1/permissions/{document_token}/members/{oid}",
                    headers=headers,
                    params={"type": doc_type, "member_type": "openid"},
                )
                d = resp.json()
                if d.get("code") == 0:
                    results.append(f"✅ 已移除「{display}」的协作权限")
                else:
                    results.append(f"❌ 移除「{display}」失败：{d.get('msg')} (code {d.get('code')})")

    return "\n".join(results) if results else "没有需要处理的成员"


# ─── Feishu Drive Delete ──────────────────────────────────────────────────────

async def _feishu_drive_delete(agent_id: uuid.UUID, arguments: dict) -> str:
    """Delete a file or folder from Feishu Drive (cloud space).
    The file is moved to the recycle bin, not permanently deleted.
    For folders, the deletion is asynchronous and returns a task_id.
    """
    import httpx

    file_token = (arguments.get("file_token") or "").strip()
    file_type = (arguments.get("file_type") or "").strip()

    if not file_token:
        return "❌ Missing required argument 'file_token'"
    if not file_type:
        return "❌ Missing required argument 'file_type'. Valid values: file, docx, bitable, folder, doc, sheet, mindnote, shortcut, slides"

    valid_types = {"file", "docx", "bitable", "folder", "doc", "sheet", "mindnote", "shortcut", "slides"}
    if file_type not in valid_types:
        return f"❌ Invalid file_type '{file_type}'. Valid values: {', '.join(sorted(valid_types))}"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    # Type label mapping for user-friendly output
    type_labels = {
        "file": "文件", "docx": "文档", "bitable": "多维表格",
        "folder": "文件夹", "doc": "旧版文档", "sheet": "电子表格",
        "mindnote": "思维笔记", "shortcut": "快捷方式", "slides": "幻灯片",
    }
    type_label = type_labels.get(file_type, file_type)

    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.delete(
                f"https://open.feishu.cn/open-apis/drive/v1/files/{file_token}",
                params={"type": file_type},
                headers={"Authorization": f"Bearer {token}"},
            )
        data = resp.json()
        code = data.get("code", -1)

        if code == 0:
            # Folder deletion returns a task_id for async tracking
            task_id = data.get("data", {}).get("task_id")
            if task_id:
                return (
                    f"✅ 已提交{type_label}删除任务（异步执行中）。\n"
                    f"📋 任务 ID: `{task_id}`\n"
                    f"文件夹删除为异步操作，文件会被移至回收站。"
                )
            return f"✅ {type_label} `{file_token}` 已删除（移至回收站）。"

        # Error handling with specific codes
        msg = data.get("msg", "Unknown error")
        if code == 1061003:
            return f"❌ 未找到文件 `{file_token}`。请确认文件 token 和类型是否正确。"
        elif code == 1061004:
            return (
                f"❌ 权限不足（code {code}）\n"
                "需要满足以下条件之一：\n"
                "• 文件所有者 + 父文件夹编辑权限\n"
                "• 父文件夹的所有者或 full_access 权限\n"
                "同时需要在飞书开放平台开通：drive:drive 或 space:document:delete"
            )
        elif code == 1061007:
            return f"❌ 文件 `{file_token}` 已被删除。"
        elif code == 1061045:
            return f"⚠️ 接口频率限制，请稍后重试。（每秒最多 5 次）"
        else:
            return f"❌ 删除{type_label}失败：{msg} (code {code})"

    except Exception as e:
        return f"❌ 删除文件异常: {str(e)[:300]}"


# ─── Feishu Calendar Tools ────────────────────────────────────────────────────

async def _feishu_calendar_list(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx
    import re as _re
    from datetime import timedelta as _td

    user_email = arguments.get("user_email", "").strip()

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    now = datetime.now(timezone.utc)

    def _to_iso(t: str | None, default: datetime) -> str:
        """Return an ISO-8601 string with timezone for freebusy API."""
        if not t:
            return default.strftime("%Y-%m-%dT%H:%M:%S+00:00")
        if _re.fullmatch(r'\d+', t.strip()):
            from datetime import datetime as _dt2
            return _dt2.fromtimestamp(int(t.strip()), tz=timezone.utc).strftime("%Y-%m-%dT%H:%M:%S+00:00")
        return t.strip()

    def _to_unix(t: str | None, default: datetime) -> str:
        """Convert ISO-8601 / Unix string / None to Unix timestamp string."""
        if not t:
            return str(int(default.timestamp()))
        if _re.fullmatch(r'\d+', t.strip()):
            return t.strip()
        try:
            from datetime import datetime as _dt2
            for fmt in ("%Y-%m-%dT%H:%M:%S%z", "%Y-%m-%dT%H:%M:%SZ", "%Y-%m-%dT%H:%M:%S"):
                try:
                    dt = _dt2.strptime(t.strip(), fmt)
                    if dt.tzinfo is None:
                        dt = dt.replace(tzinfo=timezone.utc)
                    return str(int(dt.timestamp()))
                except ValueError:
                    continue
            from dateutil import parser as _dp
            return str(int(_dp.parse(t).timestamp()))
        except Exception:
            return str(int(default.timestamp()))

    start_arg = arguments.get("start_time")
    end_arg = arguments.get("end_time")
    start_ts = _to_unix(start_arg, now)
    end_ts = _to_unix(end_arg, now + _td(days=7))
    start_iso = _to_iso(start_arg, now)
    end_iso = _to_iso(end_arg, now + _td(days=7))

    # ── 1. Query sender's real freebusy from Feishu Calendar ─────────────────
    sender_open_id = channel_feishu_sender_open_id.get(None)
    # Allow explicit override via argument
    if arguments.get("user_open_id"):
        sender_open_id = arguments["user_open_id"]
    elif user_email:
        resolved = await _feishu_resolve_open_id(token, user_email)
        if resolved:
            sender_open_id = resolved

    freebusy_section = ""
    if sender_open_id:
        try:
            async with httpx.AsyncClient(timeout=10) as fb_client:
                fb_resp = await fb_client.post(
                    "https://open.feishu.cn/open-apis/calendar/v4/freebusy/list",
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                    json={
                        "time_min": start_iso,
                        "time_max": end_iso,
                        "user_id": sender_open_id,
                    },
                )
            fb_data = fb_resp.json()
            if fb_data.get("code") == 0:
                busy_slots = fb_data.get("data", {}).get("freebusy_list", [])
                if busy_slots:
                    from datetime import datetime as _dt2
                    from zoneinfo import ZoneInfo
                    tz_cn = ZoneInfo("Asia/Shanghai")
                    busy_lines = []
                    for slot in sorted(busy_slots, key=lambda x: x.get("start_time", "")):
                        try:
                            s = _dt2.fromisoformat(slot["start_time"]).astimezone(tz_cn).strftime("%H:%M")
                            e = _dt2.fromisoformat(slot["end_time"]).astimezone(tz_cn).strftime("%H:%M")
                            busy_lines.append(f"  🔴 {s}–{e}")
                        except Exception:
                            busy_lines.append(f"  🔴 {slot.get('start_time')}–{slot.get('end_time')}")
                    freebusy_section = f"\n📌 **用户真实日历（忙碌时段）**：\n" + "\n".join(busy_lines)
                else:
                    freebusy_section = "\n📌 **用户真实日历**：该时段全部空闲。"
        except Exception as _fe:
            freebusy_section = f"\n⚠️ Freebusy 查询异常: {_fe}"

    # ── 2. Also list bot's own calendar events ───────────────────────────────
    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        # Return freebusy results even if bot calendar fails
        if freebusy_section:
            return freebusy_section.strip()
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    # Note: page_size is NOT a valid param for this API — omit it entirely
    params: dict = {}
    if start_ts:
        params["start_time"] = start_ts
    if end_ts:
        params["end_time"] = end_ts

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.get(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events",
            headers={"Authorization": f"Bearer {token}"},
            params=params,
        )

    data = resp.json()
    if data.get("code") != 0:
        if freebusy_section:
            return freebusy_section.strip()
        return f"❌ Calendar API error: {data.get('msg')} (code {data.get('code')})"

    items = data.get("data", {}).get("items", [])
    if not items and not freebusy_section:
        return "📅 该时间段内没有日程。"

    lines = []
    if items:
        lines.append(f"📅 Bot 日历共 {len(items)} 个日程：\n")
    for ev in items:
        summary = ev.get("summary", "(no title)")
        start = ev.get("start_time", {}).get("timestamp", "")
        end_t = ev.get("end_time", {}).get("timestamp", "")
        location = ev.get("location", {}).get("name", "")
        event_id = ev.get("event_id", "")
        try:
            from datetime import datetime as _dt
            s = _dt.fromtimestamp(int(start), tz=timezone.utc).strftime("%m-%d %H:%M") if start else "?"
            e = _dt.fromtimestamp(int(end_t), tz=timezone.utc).strftime("%H:%M") if end_t else "?"
        except Exception:
            s, e = start, end_t
        loc_str = f" | 📍{location}" if location else ""
        lines.append(f"- **{summary}** | 🕐{s}–{e}{loc_str}  (ID: `{event_id}`)")

    if freebusy_section:
        lines.append(freebusy_section)

    return "\n".join(lines) if lines else "📅 该时间段内没有日程。"


async def _feishu_calendar_create(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    user_email = arguments.get("user_email", "").strip()
    summary = arguments.get("summary", "").strip()
    start_time = arguments.get("start_time", "").strip()
    end_time = arguments.get("end_time", "").strip()

    for f, v in [("summary", summary), ("start_time", start_time), ("end_time", end_time)]:
        if not v:
            return f"❌ Missing required argument '{f}'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    # Resolve organizer open_id from email — soft failure
    organizer_open_id: str | None = None
    if user_email:
        organizer_open_id = await _feishu_resolve_open_id(token, user_email)
        if not organizer_open_id:
            logger.warning(f"[Feishu Calendar] Could not resolve open_id for '{user_email}', continuing without organizer invite")

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    tz = arguments.get("timezone", "Asia/Shanghai")
    body: dict = {
        "summary": summary,
        "start_time": {"timestamp": str(int(_iso_to_ts(start_time))), "timezone": tz},
        "end_time": {"timestamp": str(int(_iso_to_ts(end_time))), "timezone": tz},
    }
    if arguments.get("description"):
        body["description"] = arguments["description"]
    if arguments.get("location"):
        body["location"] = {"name": arguments["location"]}

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.post(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events",
            json=body,
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to create event: {data.get('msg')} (code {data.get('code')})"

    event_id = data.get("data", {}).get("event", {}).get("event_id", "")

    # Collect all attendee open_ids to invite
    attendee_open_ids: list[str] = []
    attendee_display: list[str] = []  # for summary message

    # 1. Direct open_ids provided by caller
    for oid in (arguments.get("attendee_open_ids") or []):
        if oid and oid not in attendee_open_ids:
            attendee_open_ids.append(oid)
            attendee_display.append(oid)

    # 2. Names → look up via feishu_user_search
    import re as _re_oid
    for aname in (arguments.get("attendee_names") or []):
        aname = aname.strip()
        if not aname:
            continue
        _sr = await _feishu_user_search(agent_id, {"name": aname})
        _m = _re_oid.search(r'open_id: `(ou_[A-Za-z0-9]+)`', _sr)
        if _m:
            _oid = _m.group(1)
            if _oid not in attendee_open_ids:
                attendee_open_ids.append(_oid)
                attendee_display.append(aname)
        else:
                logger.warning(f"[Calendar] Could not resolve attendee '{aname}': {_sr[:100]}")

    # 3. From explicit attendee_emails
    attendee_emails: list[str] = list(arguments.get("attendee_emails") or [])
    if user_email and user_email not in attendee_emails:
        attendee_emails.append(user_email)
    for email in attendee_emails[:20]:
        oid = await _feishu_resolve_open_id(token, email)
        if oid and oid not in attendee_open_ids:
            attendee_open_ids.append(oid)
            attendee_display.append(email)

    # 4. Auto-invite the Feishu message sender (from context var)
    sender_oid = channel_feishu_sender_open_id.get(None)
    if sender_oid and sender_oid not in attendee_open_ids:
        attendee_open_ids.append(sender_oid)

    if attendee_open_ids and event_id:
        async with httpx.AsyncClient(timeout=20) as client:
            for oid in attendee_open_ids:
                await client.post(
                    f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}/attendees",
                    json={"attendees": [{"type": "user", "user_id": oid}]},
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                )

    att_str = f"\n**参与人**: {', '.join(attendee_display)}" if attendee_display else ""
    invite_note = "\n（已向您发送日历邀请，请在飞书日历中确认）" if attendee_open_ids else ""
    return (
        f"✅ 日历事件已创建！\n"
        f"**标题**: {summary}\n"
        f"**时间**: {start_time} → {end_time}{att_str}\n"
        f"**Event ID**: `{event_id}`{invite_note}"
    )


async def _feishu_calendar_update(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    user_email = arguments.get("user_email", "").strip()
    event_id = arguments.get("event_id", "").strip()
    if not user_email or not event_id:
        return "❌ Both 'user_email' and 'event_id' are required."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    open_id = await _feishu_resolve_open_id(token, user_email)
    if not open_id:
        return f"❌ User '{user_email}' not found."

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    patch: dict = {}
    tz = arguments.get("timezone", "Asia/Shanghai")
    if arguments.get("summary"):
        patch["summary"] = arguments["summary"]
    if arguments.get("description"):
        patch["description"] = arguments["description"]
    if arguments.get("location"):
        patch["location"] = {"name": arguments["location"]}
    if arguments.get("start_time"):
        patch["start_time"] = {"timestamp": str(int(_iso_to_ts(arguments["start_time"]))), "timezone": tz}
    if arguments.get("end_time"):
        patch["end_time"] = {"timestamp": str(int(_iso_to_ts(arguments["end_time"]))), "timezone": tz}

    if not patch:
        return "ℹ️ No fields to update."

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.patch(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}",
            json=patch,
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to update: {data.get('msg')} (code {data.get('code')})"

    return f"✅ Event `{event_id}` updated. Changed: {', '.join(patch.keys())}."


async def _feishu_calendar_delete(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    user_email = arguments.get("user_email", "").strip()
    event_id = arguments.get("event_id", "").strip()
    if not user_email or not event_id:
        return "❌ Both 'user_email' and 'event_id' are required."

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    open_id = await _feishu_resolve_open_id(token, user_email)
    if not open_id:
        return f"❌ User '{user_email}' not found."

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.delete(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}",
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to delete: {data.get('msg')} (code {data.get('code')})"

    return f"✅ Event `{event_id}` deleted successfully."

# ─── Feishu Approval Tools ───────────────────────────────────────────────────

async def _feishu_approval_create(agent_id: uuid.UUID, arguments: dict) -> str:
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."

    approval_code = arguments.get("approval_code", "").strip()
    user_id = arguments.get("user_id", "").strip()
    form_data = arguments.get("form_data", "").strip()

    if not approval_code or not user_id or not form_data:
        return "❌ form_data, user_id and approval_code are required."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.create_approval_instance(app_id, app_secret, approval_code, user_id, form_data)
        err = _check_feishu_err(resp)
        if err: return err

        instance_code = resp.get("data", {}).get("instance_code", "")
        return f"✅ 审批发起成功！\n审批实例 ID: `{instance_code}`"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _feishu_approval_query(agent_id: uuid.UUID, arguments: dict) -> str:
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."

    approval_code = arguments.get("approval_code", "").strip()
    status = arguments.get("status")

    if not approval_code:
        return "❌ approval_code is required."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.query_approval_instances(app_id, app_secret, approval_code, status)
        err = _check_feishu_err(resp)
        if err: return err

        data = resp.get("data", {})
        instance_codes = data.get("instance_code_list", [])
        
        return f"✅ 查询完成。共发现 {len(instance_codes)} 个符合条件的审批实例。\n实例列表: {instance_codes}"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


async def _feishu_approval_get(agent_id: uuid.UUID, arguments: dict) -> str:
    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."

    instance_id = arguments.get("instance_id", "").strip()
    if not instance_id:
        return "❌ instance_id is required."

    from app.services.feishu_service import feishu_service
    try:
        resp = await feishu_service.get_approval_instance(app_id, app_secret, instance_id)
        err = _check_feishu_err(resp)
        if err: return err

        data = resp.get("data", {})
        import json
        return f"✅ 审批实例查询结果:\n```json\n{json.dumps(data, ensure_ascii=False, indent=2)}\n```"
    except Exception as e:
        return f"Failed: {str(e)[:300]}"


# ─── Feishu User Search ───────────────────────────────────────────────────────

async def _feishu_user_search(agent_id: uuid.UUID, arguments: dict) -> str:
    """Search for colleagues in the Feishu directory by name.

    Strategy:
    1. Search local contacts cache (populated when anyone messages the bot).
    2. Fall back to Contact v3 GET /users/{open_id} if we find a match by email.
    The cache is populated by feishu.py each time a message sender is resolved.
    """
    import httpx
    import json as _json
    import pathlib as _pl

    name = (arguments.get("name") or "").strip()
    if not name:
        return "❌ Missing required argument 'name'"

    app_id, app_secret = await _get_feishu_credentials(agent_id)
    if not app_id or not app_secret:
        return "❌ Agent has no Feishu channel configured."
    from app.services.feishu_service import feishu_service
    token = await feishu_service.get_tenant_access_token(app_id, app_secret)

    # ── Load local contacts cache ─────────────────────────────────────────────
    _cache_file = _pl.Path(f"/data/workspaces/{agent_id}/feishu_contacts_cache.json")
    _cached_users: list[dict] = []
    try:
        if _cache_file.exists():
            _raw = _json.loads(_cache_file.read_text())
            _cached_users = _raw.get("users", [])
    except Exception:
        pass

    name_lower = name.lower()

    def _matches(u: dict) -> bool:
        return (
            name_lower in (u.get("name") or "").lower()
            or name_lower in (u.get("en_name") or "").lower()
        )

    matched = [u for u in _cached_users if _matches(u)]

    if matched:
        lines = [f"🔍 找到 {len(matched)} 位匹配「{name}」的用户：\n"]
        for u in matched:
            open_id = u.get("open_id", "")
            user_id = u.get("user_id", "")
            display_name = u.get("name", "")
            en_name = u.get("en_name", "")
            email = u.get("email", "")
            lines.append(f"• **{display_name}**{'（' + en_name + '）' if en_name else ''}")
            if user_id:
                lines.append(f"  user_id: `{user_id}`")
            if open_id:
                lines.append(f"  open_id: `{open_id}`")
            if email:
                lines.append(f"  邮箱: {email}")
        return "\n".join(lines)

    # ── Cache miss: try OrgMember table first (has user_id from org sync) ──────
    try:
        from app.database import async_session as _async_session
        from sqlalchemy import select as _sa_select
        from app.models.org import OrgMember as _OrgMember
        async with _async_session() as _db:
            _r = await _db.execute(
                _sa_select(_OrgMember).where(_OrgMember.name.ilike(f"%{name}%"))
            )
            _org_members = _r.scalars().all()
        if _org_members:
            lines = [f"🔍 从通讯录找到 {len(_org_members)} 位匹配「{name}」的用户：\n"]
            for _om in _org_members:
                lines.append(f"• **{_om.name}**")
                if _om.external_id:
                    lines.append(f"  user_id: `{_om.external_id}`")
                if _om.open_id:
                    lines.append(f"  open_id: `{_om.open_id}`")
                if _om.email:
                    lines.append(f"  邮箱: {_om.email}")
                if _om.department_path:
                    lines.append(f"  部门: {_om.department_path}")
            return "\n".join(lines)
    except Exception:
        pass

    # ── Fallback: try User table ──────────────────────────────────────
    try:
        from app.database import async_session as _async_session
        from sqlalchemy import select as _sa_select
        from app.models.user import User as _User
        async with _async_session() as _db:
            _r = await _db.execute(
                _sa_select(_User).where(_User.display_name.ilike(f"%{name}%"))
            )
            _platform_users = _r.scalars().all()
        for _pu in _platform_users:
            _uid = getattr(_pu, "feishu_user_id", None)
            if _uid:
                result_lines = [f"🔍 找到匹配「{name}」的用户：\n", f"• **{_pu.display_name}**"]
                result_lines.append(f"  user_id: `{_uid}`")
                _email = getattr(_pu, "email", None)
                if _email:
                    result_lines.append(f"  邮箱: {_email}")
                return "\n".join(result_lines)
    except Exception:
        pass

    total = len(_cached_users)
    if total == 0:
        return (
            f"❌ 本地通讯录缓存为空，暂时无法搜索「{name}」。\n\n"
            "通讯录缓存会在同事向机器人发消息时自动建立。\n"
            "如果「覃睿」从未给机器人发过消息，可以请他先给机器人发一条消息，"
            "之后就能直接搜索到他了。\n\n"
            "或者，请直接告诉我「覃睿」的飞书 open_id 或邮箱，我可以立刻操作。"
        )
    return (
        f"❌ 未在本地通讯录（已缓存 {total} 人）中找到「{name}」。\n\n"
        "通讯录缓存来自给机器人发过消息的同事。\n"
        "如果「{name}」从未给机器人发消息，请他先发一条，之后即可自动识别。\n"
        "或者请直接提供其飞书 open_id / 工作邮箱。"
    )


async def _feishu_contacts_refresh(agent_id: uuid.UUID) -> None:
    """Force-clear the local contacts cache so next search re-fetches from API."""
    import pathlib as _pl
    _cache_file = _pl.Path("/data/workspaces") / str(agent_id) / "feishu_contacts_cache.json"
    try:
        if _cache_file.exists():
            _cache_file.unlink()
    except Exception:
        pass


# ─── Email Tool Helpers ─────────────────────────────────────

async def _get_email_config(agent_id: uuid.UUID) -> dict:
    """Retrieve per-agent email config from the send_email tool's AgentTool config."""
    from app.models.tool import Tool, AgentTool

    async with async_session() as db:
        # Find the send_email tool
        r = await db.execute(select(Tool).where(Tool.name == "send_email"))
        tool = r.scalar_one_or_none()
        if not tool:
            return {}

        # Get per-agent config
        at_r = await db.execute(
            select(AgentTool).where(
                AgentTool.agent_id == agent_id,
                AgentTool.tool_id == tool.id,
            )
        )
        at = at_r.scalar_one_or_none()
        agent_config = (at.config or {}) if at else {}
        # Merge global + agent override
        return {**(tool.config or {}), **agent_config}


# ── Pages: public HTML hosting ──────────────────────────

async def _publish_page(agent_id: uuid.UUID, user_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Publish an HTML file as a public page."""
    import secrets
    import re

    path = arguments.get("path", "")
    if not path:
        return "Missing required argument 'path'"

    # Validate file extension
    if not path.lower().endswith((".html", ".htm")):
        return "Only .html and .htm files can be published"

    # Resolve and check file exists
    full_path = (ws / path).resolve()
    if not str(full_path).startswith(str(ws.resolve())):
        return "Path traversal not allowed"
    if not full_path.exists() or not full_path.is_file():
        return f"File not found: {path}"

    # Extract title from HTML
    try:
        content = full_path.read_text(encoding="utf-8", errors="replace")
        title_match = re.search(r"<title[^>]*>(.*?)</title>", content, re.IGNORECASE | re.DOTALL)
        title = title_match.group(1).strip()[:200] if title_match else full_path.stem
    except Exception:
        title = full_path.stem

    # Generate short_id
    short_id = secrets.token_urlsafe(6)[:8]  # 8-char URL-safe string

    # Look up tenant_id
    tenant_id = None
    try:
        from app.models.agent import Agent as _AgModel
        async with async_session() as _db:
            _r = await _db.execute(select(_AgModel.tenant_id).where(_AgModel.id == agent_id))
            tenant_id = _r.scalar_one_or_none()
    except Exception:
        pass

    # Create record
    from app.models.published_page import PublishedPage
    try:
        async with async_session() as db:
            page = PublishedPage(
                short_id=short_id,
                agent_id=agent_id,
                user_id=user_id,
                tenant_id=tenant_id,
                source_path=path,
                title=title,
            )
            db.add(page)
            await db.commit()
    except Exception as e:
        return f"Failed to publish: {e}"

    # Build public URL using configured PUBLIC_BASE_URL
    from app.services.platform_service import platform_service
    async with async_session() as db2:
        public_base = await platform_service.get_public_base_url(db2)

    url = f"{public_base}/p/{short_id}" if public_base else f"/p/{short_id}"

    return f"Published successfully!\n\nPublic URL: {url}\nTitle: {title}\n\nAnyone can access this page without logging in."


async def _list_published_pages(agent_id: uuid.UUID) -> str:
    """List all published pages for this agent."""
    from app.models.published_page import PublishedPage

    try:
        async with async_session() as db:
            result = await db.execute(
                select(PublishedPage)
                .where(PublishedPage.agent_id == agent_id)
                .order_by(PublishedPage.created_at.desc())
            )
            pages = result.scalars().all()

        if not pages:
            return "No published pages yet."

        lines = [f"Published pages ({len(pages)} total):\n"]
        for p in pages:
            lines.append(f"- {p.title or 'Untitled'}")
            lines.append(f"  URL: /p/{p.short_id}")
            lines.append(f"  Source: {p.source_path}")
            lines.append(f"  Views: {p.view_count}")
            lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"Failed to list pages: {e}"


# ─── AgentBay Tool Handlers ─────────────────────────────────────

async def _agentbay_browser_navigate(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """AgentBay browser navigation.

    After navigating, always captures a screenshot.  Whether that screenshot is
    stored to disk or kept only in memory depends on save_to_workspace:
      - False (default): bytes are held in the process-level memory cache;
        the returned sentinel [ImageID: ...] is consumed by vision_inject.py
        in the same request cycle and then discarded — zero disk writes.
      - True: screenshot is written to workspace/ so the user can see it in
        their file manager, and a Markdown link is included in the return value.
    """
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    url = arguments.get("url", "")
    wait_for = arguments.get("wait_for", "")
    save_to_workspace = arguments.get("save_to_workspace", False)

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id)
        # Always request a screenshot for navigation so the model can observe the result
        result = await client.browser_navigate(url, wait_for=wait_for, screenshot=True)

        # Build text parts from the navigation result
        parts = [f"✅ 已访问: {url}"]
        if result.get("title"):
            parts.append(f"标题: {result['title']}")
        if result.get("content"):
            content = result["content"][:3000]
            parts.append(f"内容:\n{content}")
        logger.info(f"[AgentBay] Browser navigate result: {result.get('title')}")

        screenshot_data = result.get("screenshot")
        if screenshot_data:
            import base64 as _base64
            # Normalise to raw bytes regardless of whether it's a data URL or plain b64
            if isinstance(screenshot_data, str):
                if screenshot_data.startswith("data:image"):
                    screenshot_data = screenshot_data.split(",", 1)[1]
                raw_bytes = _base64.b64decode(screenshot_data)
            elif isinstance(screenshot_data, bytes):
                raw_bytes = screenshot_data
            else:
                raw_bytes = None

            if raw_bytes:
                if save_to_workspace:
                    # Persist to workspace/ so the user can see the file
                    import time as _time
                    rel_path = f"workspace/screenshot_{int(_time.time())}.png"
                    screenshot_path = ws / rel_path
                    screenshot_path.parent.mkdir(parents=True, exist_ok=True)
                    screenshot_path.write_bytes(raw_bytes)
                    parts.append(
                        f"截图已保存至 `{rel_path}`。\n"
                        f"![Browser Navigation Screenshot](/api/agents/{agent_id}/files/download?path={rel_path})\n"
                        f"CRITICAL: Do NOT call 'send_channel_file' or 'upload_image'. Just print the Markdown above exactly as shown."
                    )
                    logger.info(f"[AgentBay] Browser navigate screenshot saved to {rel_path}")
                else:
                    # Store in memory only — vision_inject.py will consume it
                    from app.services.vision_inject import store_temp_screenshot
                    img_id = store_temp_screenshot(raw_bytes)
                    parts.append(
                        f"Internal screenshot captured for analysis. [ImageID: {img_id}]\n"
                        f"NOTE: This screenshot is for YOUR eyes only (LLM vision). The user CANNOT see it. "
                        f"If the user asked to SEE a screenshot, call this tool again with save_to_workspace=true."
                    )
                    logger.info(f"[AgentBay] Browser navigate screenshot stored in memory (id={img_id})")

        return "\n\n".join(parts)

    except RuntimeError as e:
        return f"❌ {str(e)}。请先在 Agent 设置中配置 AgentBay 通道。"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser navigate failed for agent {agent_id}")
        return f"❌ AgentBay 浏览器访问失败: {str(e)[:200]}"


async def _agentbay_browser_screenshot(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Take a screenshot of the CURRENT browser page without navigating.

    Correct way to observe the result of a click, type, or form submit — never
    call browser_navigate again just to screenshot, that refreshes the page.

    By default (save_to_workspace=False) the image is held in the process-level
    memory cache and consumed once by the LLM vision pipeline — no disk write,
    nothing shown in the user's file manager or chat history.
    Set save_to_workspace=True to persist and display the image.
    """
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    save_to_workspace = arguments.get("save_to_workspace", False)

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id)
        result = await client.browser_screenshot()

        screenshot_data = result.get("screenshot")
        if not screenshot_data:
            return "❌ 截图失败：未返回图像数据"

        import base64 as _base64
        # Normalise to raw bytes
        if isinstance(screenshot_data, str):
            if screenshot_data.startswith("data:image"):
                screenshot_data = screenshot_data.split(",", 1)[1]
            raw_bytes = _base64.b64decode(screenshot_data)
        elif isinstance(screenshot_data, bytes):
            raw_bytes = screenshot_data
        else:
            return "❌ 截图失败：未知数据格式"

        if save_to_workspace:
            # Persist to workspace/ so the user can see the file
            import time as _time
            rel_path = f"workspace/screenshot_{int(_time.time())}.png"
            screenshot_path = ws / rel_path
            screenshot_path.parent.mkdir(parents=True, exist_ok=True)
            screenshot_path.write_bytes(raw_bytes)
            logger.info(f"[AgentBay] Browser screenshot saved to workspace: {rel_path}")
            return (
                f"✅ 截图已保存至 `{rel_path}`。\n"
                f"![Browser Screenshot](/api/agents/{agent_id}/files/download?path={rel_path})\n"
                f"CRITICAL: Do NOT call 'send_channel_file' or 'upload_image'. Just print the Markdown above exactly as shown."
            )
        else:
            # Store in memory only — vision_inject.py will consume it for LLM vision
            from app.services.vision_inject import store_temp_screenshot
            img_id = store_temp_screenshot(raw_bytes)
            logger.info(f"[AgentBay] Browser screenshot stored in memory (id={img_id})")
            return (
                f"Internal screenshot captured for analysis. [ImageID: {img_id}]\n"
                f"NOTE: This screenshot is for YOUR eyes only (LLM vision). The user CANNOT see it. "
                f"If the user asked to SEE a screenshot, call this tool again with save_to_workspace=true."
            )

    except RuntimeError as e:
        return f"❌ {str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser screenshot failed for agent {agent_id}")
        return f"❌ 截图失败: {str(e)[:200]}"


async def _agentbay_browser_click(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """AgentBay 浏览器点击。"""
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    selector = arguments.get("selector", "")

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id)
        await client.browser_click(selector)
        return f"✅ 已点击元素: {selector}"
    except RuntimeError as e:
        return f"❌ {str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser click failed")
        return f"❌ 点击失败: {str(e)[:200]}"


async def _agentbay_browser_type(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """AgentBay 浏览器输入。"""
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    selector = arguments.get("selector", "")
    text = arguments.get("text", "")

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id)
        await client.browser_type(selector, text)
        return f"✅ 已在 {selector} 输入文本"
    except RuntimeError as e:
        return f"❌ {str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Browser type failed")
        return f"❌ 输入失败: {str(e)[:200]}"


async def _agentbay_code_execute(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """在 AgentBay 代码空间执行代码。"""
    if not agent_id:
        return "❌ AgentBay 工具需要 agent 上下文"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    language = arguments.get("language", "python")
    code = arguments.get("code", "")
    timeout = arguments.get("timeout", 30)

    if not code.strip():
        return "❌ 请提供要执行的代码"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "code", session_id=_session_id)
        result = await client.code_execute(language, code, timeout)

        # 格式化返回结果
        parts = [f"✅ 代码执行完成 ({language})"]
        if result.get("stdout"):
            parts.append(f"📤 输出:\n{result['stdout']}")
        if result.get("stderr"):
            parts.append(f"⚠️ 错误输出:\n{result['stderr']}")
        if result.get("exit_code") != 0:
            parts.append(f"退出码: {result['exit_code']}")

        return "\n\n".join(parts)

    except RuntimeError as e:
        return f"❌ {str(e)}。请先在 Agent 设置中配置 AgentBay 通道。"
    except Exception as e:
        logger.exception(f"[AgentBay] Code execution failed for agent {agent_id}")
        return f"❌ 代码执行失败: {str(e)[:200]}"


async def _handle_email_tool(tool_name: str, agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Dispatch email tool calls to the email_service module."""
    from app.services.email_service import send_email, read_emails, reply_email

    config = await _get_email_config(agent_id)
    if not config.get("email_address") or not config.get("auth_code"):
        return (
            "❌ Email not configured for this agent.\n\n"
            "Please go to Agent → Tools → Send Email → Config to set up your email:\n"
            "1. Select your email provider\n"
            "2. Enter your email address\n"
            "3. Enter your authorization code (not your login password)"
        )

    try:
        if tool_name == "send_email":
            return await send_email(
                config=config,
                to=arguments.get("to", ""),
                subject=arguments.get("subject", ""),
                body=arguments.get("body", ""),
                cc=arguments.get("cc"),
                attachments=arguments.get("attachments"),
                workspace_path=ws,
            )
        elif tool_name == "read_emails":
            return await read_emails(
                config=config,
                limit=arguments.get("limit", 10),
                search=arguments.get("search"),
                folder=arguments.get("folder", "INBOX"),
            )
        elif tool_name == "reply_email":
            return await reply_email(
                config=config,
                message_id=arguments.get("message_id", ""),
                body=arguments.get("body", ""),
            )
        else:
            return f"❌ Unknown email tool: {tool_name}"
    except Exception as e:
        return f"❌ Email tool error: {str(e)[:200]}"


# ─── Skill Management Tools ────────────────────────────────────


async def _search_clawhub(agent_id: uuid.UUID, arguments: dict) -> str:
    """Search the ClawHub skill registry."""
    import httpx
    query = arguments.get("query", "").strip()
    if not query:
        return "Missing required argument 'query'"

    # Resolve tenant ClawHub API key
    from app.api.skills import _get_clawhub_key
    tenant_id = await _get_agent_tenant_id(agent_id)
    api_key = await _get_clawhub_key(tenant_id)
    headers = {"Authorization": f"Bearer {api_key}"} if api_key else {}

    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                "https://clawhub.ai/api/search",
                params={"q": query},
                headers=headers,
            )
            if resp.status_code != 200:
                return f"ClawHub search failed (HTTP {resp.status_code})"
            data = resp.json()
    except Exception as e:
        return f"❌ ClawHub search error: {str(e)[:200]}"

    results = data.get("results", [])
    if not results:
        return f"No skills found matching '{query}'."

    lines = [f"Found {len(results)} skill(s) matching '{query}':\n"]
    for r in results:
        name = r.get("displayName") or r.get("slug", "?")
        slug = r.get("slug", "")
        summary = (r.get("summary") or "")[:120]
        updated = ""
        if r.get("updatedAt"):
            from datetime import datetime
            try:
                dt = datetime.fromtimestamp(r["updatedAt"] / 1000)
                updated = f" | Updated: {dt.strftime('%Y-%m-%d')}"
            except Exception:
                pass
        lines.append(f"• **{name}** (`{slug}`){updated}")
        if summary:
            lines.append(f"  {summary}")
    lines.append("\nTo install a skill, use: install_skill(source=\"<slug>\")")
    return "\n".join(lines)


async def _install_skill(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Install a skill from ClawHub slug or GitHub URL into the agent's workspace."""
    import httpx
    source = arguments.get("source", "").strip()
    if not source:
        return "❌ Missing required argument 'source'. Provide a ClawHub slug (e.g. 'market-research') or a GitHub URL."

    is_url = source.startswith("http://") or source.startswith("https://")
    base = ws  # agent workspace dir (skills/ lives under workspace/)

    try:
        if is_url:
            # ── GitHub URL path ──
            from app.api.skills import _parse_github_url, _fetch_github_directory, _get_github_token

            parsed = _parse_github_url(source)
            if not parsed:
                return "❌ Invalid GitHub URL. Expected format: https://github.com/{owner}/{repo} or https://github.com/{owner}/{repo}/tree/{branch}/{path}"

            owner, repo, branch, path = parsed["owner"], parsed["repo"], parsed["branch"], parsed["path"]
            tenant_id = await _get_agent_tenant_id(agent_id)
            token = await _get_github_token(tenant_id)
            files = await _fetch_github_directory(owner, repo, path, branch, token)
            if not files:
                return "❌ No files found at the specified URL."

            folder_name = path.rstrip("/").split("/")[-1] if path else repo
        else:
            # ── ClawHub slug path ──
            slug = source
            from app.api.skills import _fetch_github_directory, _get_github_token, _get_clawhub_key

            # 1. Fetch metadata from ClawHub (with tenant API key)
            tenant_id = await _get_agent_tenant_id(agent_id)
            api_key = await _get_clawhub_key(tenant_id)
            ch_headers = {"Authorization": f"Bearer {api_key}"} if api_key else {}
            try:
                async with httpx.AsyncClient(timeout=15) as client:
                    resp = await client.get(f"https://clawhub.ai/api/v1/skills/{slug}", headers=ch_headers)
                    if resp.status_code == 404:
                        return f"Skill '{slug}' not found on ClawHub. Use search_clawhub to find available skills."
                    if resp.status_code != 200:
                        return f"ClawHub API error (HTTP {resp.status_code})"
                    meta = resp.json()
            except Exception as e:
                return f"Failed to connect to ClawHub: {str(e)[:200]}"

            owner_info = meta.get("owner", {})
            handle = owner_info.get("handle", "").lower()
            if not handle:
                return "❌ Could not determine skill owner from ClawHub metadata."

            # 2. Fetch files from GitHub
            github_path = f"skills/{handle}/{slug}"
            token = await _get_github_token(tenant_id)
            files = await _fetch_github_directory("openclaw", "skills", github_path, "main", token)
            if not files:
                return f"❌ No files found for skill '{slug}' in GitHub archive."

            folder_name = slug

        # 3. Write files to agent workspace
        skill_dir = base / "skills" / folder_name
        skill_dir.mkdir(parents=True, exist_ok=True)

        written = []
        for f in files:
            file_path = (skill_dir / f["path"]).resolve()
            if not str(file_path).startswith(str(base.resolve())):
                continue  # safety: skip path traversal
            file_path.parent.mkdir(parents=True, exist_ok=True)
            file_path.write_text(f["content"], encoding="utf-8")
            written.append(f["path"])

        return f"✅ Skill '{folder_name}' installed successfully ({len(written)} files written to skills/{folder_name}/).\n\nFiles: {', '.join(written)}"

    except Exception as e:
        return f"❌ Install failed: {str(e)[:300]}"


# ─── AgentBay: Browser Extract & Observe ────────────────────────────────

async def _agentbay_browser_extract(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Extract structured data from current browser page."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    instruction = arguments.get("instruction", "")
    selector = arguments.get("selector", "")

    if not instruction.strip():
        return "Missing required argument 'instruction'"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id)
        result = await client.browser_extract(instruction, selector=selector)

        if result.get("success"):
            import json
            data = result.get("data", {})
            data_str = json.dumps(data, ensure_ascii=False, indent=2) if isinstance(data, (dict, list)) else str(data)
            return f"Extraction successful:\n\n{data_str[:5000]}"
        else:
            return f"Extraction failed: {result}"

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Browser extract failed for agent {agent_id}")
        return f"Browser extract failed: {str(e)[:200]}"


async def _agentbay_browser_observe(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Observe the current browser page state."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    instruction = arguments.get("instruction", "")
    selector = arguments.get("selector", "")

    if not instruction.strip():
        return "Missing required argument 'instruction'"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id)
        result = await client.browser_observe(instruction, selector=selector)

        if result.get("success"):
            import json
            elements = result.get("elements", [])
            if not elements:
                return "No interactive elements found matching your instruction."
            elements_str = json.dumps(elements, ensure_ascii=False, indent=2)
            return f"Found {len(elements)} interactive element(s):\n\n{elements_str[:5000]}"
        else:
            return f"Observation failed: {result}"

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Browser observe failed for agent {agent_id}")
        return f"Browser observe failed: {str(e)[:200]}"


# ─── AgentBay: Command (Shell) ──────────────────────────────────────────

async def _agentbay_browser_login(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Perform an automated login using AgentBay's built-in login skill.

    Supports complex login flows including CAPTCHAs, OTP inputs,
    and multi-step authentication via AgentBay's AI-driven capability.
    """
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    url = arguments.get("url", "")
    login_config = arguments.get("login_config", "")

    if not url.strip():
        return "Missing required argument 'url'"
    if not login_config.strip():
        return "Missing required argument 'login_config' (JSON string with api_key + skill_id)"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "browser", session_id=_session_id)
        result = await client.browser_login(url, login_config)

        if result.get("success"):
            return f"Login completed successfully. {result.get('message', '')}"
        else:
            return f"Login failed: {result.get('message', 'Unknown error')}"

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Browser login failed for agent {agent_id}")
        return f"Login failed: {str(e)[:200]}"


async def _agentbay_command_exec(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Execute a shell command in the AgentBay environment."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    command = arguments.get("command", "")
    timeout_ms = arguments.get("timeout_ms", 50000)
    cwd = arguments.get("cwd", "")

    if not command.strip():
        return "Missing required argument 'command'"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "code", session_id=_session_id)
        result = await client.command_exec(command, timeout_ms=timeout_ms, cwd=cwd)

        parts = []
        if result.get("success"):
            parts.append(f"Command executed successfully (exit code: {result.get('exit_code', 0)})")
        else:
            parts.append(f"Command failed (exit code: {result.get('exit_code', -1)})")

        if result.get("stdout"):
            parts.append(f"stdout:\n{result['stdout'][:3000]}")
        if result.get("stderr"):
            parts.append(f"stderr:\n{result['stderr'][:1000]}")
        if result.get("error_message"):
            parts.append(f"Error: {result['error_message']}")

        return "\n\n".join(parts)

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Command exec failed for agent {agent_id}")
        return f"Command execution failed: {str(e)[:200]}"


# ─── AgentBay: Computer Use Handlers ────────────────────────────────────

def _save_screenshot_to_workspace(agent_id: uuid.UUID, ws: Path, data) -> str:
    """Save screenshot data to workspace and return markdown image link.

    Common helper for computer_screenshot and browser screenshot data.
    """
    import time
    import base64

    rel_path = f"workspace/desktop-screenshot-{int(time.time())}.png"
    screenshot_path = ws / rel_path
    screenshot_path.parent.mkdir(parents=True, exist_ok=True)

    # Handle various data formats from the SDK
    if isinstance(data, str):
        if data.startswith("data:image"):
            data = data.split(",", 1)[1]
        raw_bytes = base64.b64decode(data)
    elif isinstance(data, bytes):
        raw_bytes = data
    else:
        return ""

    screenshot_path.write_bytes(raw_bytes)
    return (
        f"Screenshot saved to `{rel_path}`.\n\n"
        f"![Desktop Screenshot](/api/agents/{agent_id}/files/download?path={rel_path})\n"
        f"CRITICAL: Do NOT call 'send_channel_file' or 'upload_image'. Just print the Markdown above exactly as shown."
    )


async def _agentbay_computer_screenshot(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Take a screenshot of the AgentBay cloud desktop.

    By default (save_to_workspace=False) the image is held in the process-level
    memory cache for LLM vision analysis only — no disk write, nothing shown in
    the user's file manager or chat history.
    Set save_to_workspace=True to persist and display the image.
    """
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    save_to_workspace = arguments.get("save_to_workspace", False)

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_screenshot()

        if not (result.get("success") and result.get("data")):
            return f"Screenshot failed: {result.get('error_message', 'Unknown error')}"

        raw_data = result["data"]

        # Normalise to raw bytes regardless of SDK return format
        import base64 as _base64
        if isinstance(raw_data, str):
            if raw_data.startswith("data:image"):
                raw_data = raw_data.split(",", 1)[1]
            raw_bytes = _base64.b64decode(raw_data)
        elif isinstance(raw_data, bytes):
            raw_bytes = raw_data
        else:
            return "Screenshot captured but data format is unrecognised."

        if save_to_workspace:
            # Persist to workspace/ for user visibility
            import time as _time
            rel_path = f"workspace/desktop-screenshot-{int(_time.time())}.png"
            screenshot_path = ws / rel_path
            screenshot_path.parent.mkdir(parents=True, exist_ok=True)
            screenshot_path.write_bytes(raw_bytes)
            logger.info(f"[AgentBay] Desktop screenshot saved to workspace: {rel_path}")
            return (
                f"Desktop screenshot saved to `{rel_path}`.\n"
                f"![Desktop Screenshot](/api/agents/{agent_id}/files/download?path={rel_path})\n"
                f"CRITICAL: Do NOT call 'send_channel_file' or 'upload_image'. Just print the Markdown above exactly as shown."
            )
        else:
            # Store in memory only — vision_inject.py will consume it for LLM vision
            from app.services.vision_inject import store_temp_screenshot
            img_id = store_temp_screenshot(raw_bytes)
            logger.info(f"[AgentBay] Desktop screenshot stored in memory (id={img_id})")
            return (
                f"Internal desktop screenshot captured for analysis. [ImageID: {img_id}]\n"
                f"NOTE: This screenshot is for YOUR eyes only (LLM vision). The user CANNOT see it. "
                f"If the user asked to SEE a screenshot, call this tool again with save_to_workspace=true."
            )

    except RuntimeError as e:
        return f"{str(e)}. Please configure AgentBay in Agent settings."
    except Exception as e:
        logger.exception(f"[AgentBay] Computer screenshot failed for agent {agent_id}")
        return f"Desktop screenshot failed: {str(e)[:200]}"


async def _agentbay_computer_click(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Click the mouse at specific coordinates on the desktop."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    x = arguments.get("x", 0)
    y = arguments.get("y", 0)
    button = arguments.get("button", "left")

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_click(x, y, button=button)
        if result.get("success"):
            return f"Clicked at ({x}, {y}) with {button} button"
        return f"Click failed at ({x}, {y})"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer click failed")
        return f"Click failed: {str(e)[:200]}"


async def _agentbay_computer_input_text(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Type text at the current cursor position."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    text = arguments.get("text", "")
    if not text:
        return "Missing required argument 'text'"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_input_text(text)
        if result.get("success"):
            return f"Typed text: {text[:100]}"
        return f"Text input failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer input_text failed")
        return f"Text input failed: {str(e)[:200]}"


async def _agentbay_computer_press_keys(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Press keyboard keys or shortcuts."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    keys = arguments.get("keys", [])
    hold = arguments.get("hold", False)

    if not keys:
        return "Missing required argument 'keys'"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_press_keys(keys, hold=hold)
        key_str = "+".join(keys)
        if result.get("success"):
            return f"Pressed keys: {key_str}" + (" (held)" if hold else "")
        return f"Key press failed: {key_str}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer press_keys failed")
        return f"Key press failed: {str(e)[:200]}"


async def _agentbay_computer_scroll(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Scroll the screen at a specific position."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    x = arguments.get("x", 0)
    y = arguments.get("y", 0)
    direction = arguments.get("direction", "down")
    amount = arguments.get("amount", 1)

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_scroll(x, y, direction=direction, amount=amount)
        if result.get("success"):
            return f"Scrolled {direction} by {amount} step(s) at ({x}, {y})"
        return f"Scroll failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer scroll failed")
        return f"Scroll failed: {str(e)[:200]}"


async def _agentbay_computer_move_mouse(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Move mouse to coordinates without clicking."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    x = arguments.get("x", 0)
    y = arguments.get("y", 0)

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_move_mouse(x, y)
        if result.get("success"):
            return f"Mouse moved to ({x}, {y})"
        return f"Mouse move failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer move_mouse failed")
        return f"Mouse move failed: {str(e)[:200]}"


async def _agentbay_computer_drag_mouse(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Drag mouse from one position to another."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    from_x = arguments.get("from_x", 0)
    from_y = arguments.get("from_y", 0)
    to_x = arguments.get("to_x", 0)
    to_y = arguments.get("to_y", 0)
    button = arguments.get("button", "left")

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_drag_mouse(from_x, from_y, to_x, to_y, button=button)
        if result.get("success"):
            return f"Dragged from ({from_x}, {from_y}) to ({to_x}, {to_y})"
        return f"Drag failed"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer drag_mouse failed")
        return f"Drag failed: {str(e)[:200]}"


async def _agentbay_computer_get_screen_size(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Get the screen resolution."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_get_screen_size()
        if result.get("success"):
            import json
            data = result.get("data")
            data_str = json.dumps(data, ensure_ascii=False) if isinstance(data, (dict, list)) else str(data)
            return f"Screen size: {data_str}"
        return f"Failed to get screen size: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer get_screen_size failed")
        return f"Get screen size failed: {str(e)[:200]}"


async def _agentbay_computer_start_app(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Start an application on the desktop."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    cmd = arguments.get("cmd", "")
    work_dir = arguments.get("work_dir", "")

    if not cmd.strip():
        return "Missing required argument 'cmd'"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_start_app(cmd, work_dir=work_dir)
        if result.get("success"):
            # result.data may contain non-serializable objects (e.g. Process),
            # so convert to string safely instead of json.dumps()
            data = result.get("data")
            if data is not None:
                try:
                    import json
                    data_str = json.dumps(data, ensure_ascii=False, indent=2) if isinstance(data, (dict, list, str, int, float, bool)) else str(data)
                except (TypeError, ValueError):
                    data_str = str(data)
            else:
                data_str = ""
            return f"Application started: {cmd}" + (f"\n\n{data_str[:1000]}" if data_str else "")
        return f"Failed to start application: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer start_app failed")
        return f"Start application failed: {str(e)[:200]}"


async def _agentbay_computer_get_cursor_position(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Get current cursor position."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_get_cursor_position()
        if result.get("success"):
            import json
            data = result.get("data")
            data_str = json.dumps(data, ensure_ascii=False) if isinstance(data, (dict, list)) else str(data)
            return f"Cursor position: {data_str}"
        return f"Failed to get cursor position: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer get_cursor_position failed")
        return f"Get cursor position failed: {str(e)[:200]}"


async def _agentbay_computer_get_active_window(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Get info about the currently active window."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_get_active_window()
        if result.get("success"):
            import json
            window = result.get("window")
            window_str = json.dumps(window, ensure_ascii=False, indent=2) if isinstance(window, dict) else str(window)
            return f"Active window:\n\n{window_str}"
        return f"Failed to get active window: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer get_active_window failed")
        return f"Get active window failed: {str(e)[:200]}"


async def _agentbay_computer_activate_window(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """Activate (bring to front) a window by its ID."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    window_id = arguments.get("window_id")
    if window_id is None:
        return "Missing required argument 'window_id'"

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_activate_window(int(window_id))
        if result.get("success"):
            return f"Window {window_id} activated (brought to front)"
        return f"Failed to activate window {window_id}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer activate_window failed")
        return f"Activate window failed: {str(e)[:200]}"


async def _agentbay_computer_list_visible_apps(agent_id: Optional[uuid.UUID], ws: Path, arguments: dict) -> str:
    """List currently visible/running applications."""
    if not agent_id:
        return "AgentBay tools require agent context"

    from app.services.agentbay_client import get_agentbay_client_for_agent

    try:
        _session_id = arguments.pop("_session_id", "")
        client = await get_agentbay_client_for_agent(agent_id, "computer", session_id=_session_id)
        result = await client.computer_list_visible_apps()
        if result.get("success"):
            import json
            apps = result.get("apps", [])
            if not apps:
                return "No visible applications running."
            apps_str = json.dumps(apps, ensure_ascii=False, indent=2)
            return f"Visible applications ({len(apps)}):\n\n{apps_str[:3000]}"
        return f"Failed to list applications: {result.get('error_message', 'Unknown error')}"
    except RuntimeError as e:
        return f"{str(e)}"
    except Exception as e:
        logger.exception(f"[AgentBay] Computer list_visible_apps failed")
        return f"List applications failed: {str(e)[:200]}"
